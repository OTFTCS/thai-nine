"""
Build M01-L003 — Pronouns and Who Questions
~11 slides based on blueprint CSV content.
"""

import sys
sys.path.insert(0, '/home/ubuntu/immersion-thai-slides')

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from design_system import *

def build_l003():
    prs = create_presentation()
    
    # ── Slide 1: Lesson Opener ──────────────────────────────────────────
    add_lesson_opener(
        prs,
        lesson_id="M01-L003",
        lesson_title="Pronouns and Who Questions",
        module_title="Module 1 — First Contact and Courtesy",
        level="A0"
    )
    
    # ── Slide 2: What You Will Learn ────────────────────────────────────
    add_objectives_slide(prs, [
        "Use basic pronouns: ฉัน, ผม, คุณ, เขา",
        "Ask who someone is with ใคร",
        "Identify people with นี่คือ...",
        "Confirm identity with คุณคือ...ใช่ไหม",
    ])
    
    # ── Slide 3: Pronoun Overview ───────────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "1", "Basic Thai pronouns")
    
    add_textbox(slide, CONTENT_LEFT, Inches(0.75), CONTENT_WIDTH, Inches(0.4),
               text="Four pronouns for first-contact conversations:",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_MEDIUM)
    
    pronouns = [
        ("ฉัน", "chǎn", "I", "common / all-purpose", ACCENT_TEAL),
        ("ผม", "phǒm", "I (male)", "male speaker", ACCENT_CLAY),
        ("คุณ", "khun", "you", "polite address", ACCENT_GOLD),
        ("เขา", "khǎo", "he / she", "third person", ACCENT_TEAL),
    ]
    
    y = Inches(1.3)
    for thai, translit, eng, usage, color in pronouns:
        card_h = Inches(1.1)
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), card_h,
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        
        # Color tag
        tag = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            CONTENT_LEFT + Inches(0.15), y + Inches(0.15),
            Inches(1.5), Inches(0.75)
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()
        tag.adjustments[0] = 0.15
        
        tf = tag.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = thai
        run.font.name = FONT_THAI
        run.font.size = SIZE_THAI_MED
        run.font.color.rgb = WHITE
        run.font.bold = True
        
        # Translit + English
        add_textbox(slide, CONTENT_LEFT + Inches(1.9), y + Inches(0.1),
                   Inches(5.0), Inches(0.35),
                   text=f"{translit}  —  {eng}",
                   font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
                   font_color=INK_DARK, bold=True)
        add_textbox(slide, CONTENT_LEFT + Inches(1.9), y + Inches(0.5),
                   Inches(5.0), Inches(0.3),
                   text=usage,
                   font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
                   font_color=INK_LIGHT)
        
        y += card_h + Inches(0.12)
    
    # ── Slide 4: Who Question — ใคร ─────────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "2", "Ask who someone is")
    
    # Key word
    add_phrase_card(
        slide, CONTENT_LEFT, Inches(1.0), Inches(6.5),
        thai="ใคร",
        translit="khrai",
        english="who",
        thai_size=SIZE_THAI_LARGE,
        accent_color=ACCENT_GOLD
    )
    
    add_textbox(slide, CONTENT_LEFT, Inches(3.1), CONTENT_WIDTH, Inches(0.5),
               text="Use ใคร to ask about identity — who is this person?",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_MEDIUM)
    
    # Question pattern
    q_card = add_rounded_rect(slide, CONTENT_LEFT, Inches(3.8), Inches(7.2), Inches(1.4),
                             fill_color=HIGHLIGHT_BG, border_color=ACCENT_SOFT_GOLD)
    add_accent_bar(slide, CONTENT_LEFT + Inches(0.12), Inches(4.0),
                  height=Inches(1.0), color=ACCENT_GOLD)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(3.9), Inches(6.0), Inches(0.55),
               text="เขาคือใคร", font_name=FONT_THAI, font_size=SIZE_THAI_MED,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(4.45), Inches(6.0), Inches(0.3),
               text="khǎo khʉʉ khrai  —  who is he/she?",
               font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM)
    
    # ── Slide 5: นี่คือ — This is... ────────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "2", "Identify people with นี่คือ")
    
    # Core words
    words = [
        ("นี่", "nîi", "this"),
        ("คือ", "khʉʉ", "is / this is"),
    ]
    
    y = Inches(1.0)
    for thai, translit, eng in words:
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), Inches(1.1),
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        add_accent_bar(slide, CONTENT_LEFT + Inches(0.1), y + Inches(0.15),
                      height=Inches(0.8), color=ACCENT_TEAL)
        add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.05),
                   Inches(2.5), Inches(0.5),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_MED,
                   font_color=INK_DARK, bold=True)
        add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.55),
                   Inches(2.5), Inches(0.3),
                   text=translit, font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
                   font_color=INK_MEDIUM, italic=True)
        add_textbox(slide, CONTENT_LEFT + Inches(3.2), y + Inches(0.2),
                   Inches(3.5), Inches(0.4),
                   text=eng, font_name=FONT_LATIN, font_size=SIZE_BODY,
                   font_color=INK_DARK, bold=True)
        y += Inches(1.25)
    
    # Combined pattern
    pattern_card = add_rounded_rect(slide, CONTENT_LEFT, y + Inches(0.2), Inches(7.2), Inches(1.4),
                                   fill_color=HIGHLIGHT_BG, border_color=ACCENT_SOFT_GOLD)
    add_accent_bar(slide, CONTENT_LEFT + Inches(0.12), y + Inches(0.4),
                  height=Inches(1.0), color=ACCENT_GOLD)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), y + Inches(0.3), Inches(6.0), Inches(0.55),
               text="นี่คือ...", font_name=FONT_THAI, font_size=SIZE_THAI_MED,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), y + Inches(0.85), Inches(6.0), Inches(0.3),
               text="nîi khʉʉ...  —  this is...",
               font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM)
    
    # ── Slide 6: Confirmation Pattern ───────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "3", "Confirm identity")
    
    add_textbox(slide, CONTENT_LEFT, Inches(0.75), CONTENT_WIDTH, Inches(0.4),
               text="Use ใช่ไหม to check if you are right:",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_MEDIUM)
    
    # Question
    q_card = add_rounded_rect(slide, CONTENT_LEFT, Inches(1.3), Inches(7.2), Inches(1.5),
                             fill_color=HIGHLIGHT_BG, border_color=ACCENT_SOFT_GOLD)
    add_accent_bar(slide, CONTENT_LEFT + Inches(0.12), Inches(1.5),
                  height=Inches(1.1), color=ACCENT_GOLD)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(1.35), Inches(1.5), Inches(0.3),
               text="Question", font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=ACCENT_GOLD, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(1.7), Inches(6.0), Inches(0.55),
               text="คุณคือ...ใช่ไหม", font_name=FONT_THAI, font_size=SIZE_THAI_MED,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(2.25), Inches(6.0), Inches(0.3),
               text="khun khʉʉ... châi mái  —  you are... right?",
               font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM)
    
    # Answers
    y = Inches(3.1)
    answers = [
        ("ใช่ครับ/ค่ะ", "châi khráp/khâ", "Yes, that's right.", ACCENT_TEAL),
        ("ไม่ใช่ครับ/ค่ะ", "mâi châi khráp/khâ", "No, that's not right.", ACCENT_CLAY),
    ]
    
    for thai, translit, eng, color in answers:
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), Inches(1.2),
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        add_accent_bar(slide, CONTENT_LEFT + Inches(0.12), y + Inches(0.15),
                      height=Inches(0.9), color=color)
        add_textbox(slide, CONTENT_LEFT + Inches(0.4), y + Inches(0.05),
                   Inches(5.5), Inches(0.5),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
                   font_color=INK_DARK, bold=True)
        add_textbox(slide, CONTENT_LEFT + Inches(0.4), y + Inches(0.5),
                   Inches(5.5), Inches(0.3),
                   text=f"{translit}  —  {eng}",
                   font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
                   font_color=INK_MEDIUM)
        y += Inches(1.35)
    
    # ── Slide 7: Drill — Point and Identify ─────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "3", "Point, identify, and ask")
    
    add_textbox(slide, CONTENT_LEFT, Inches(0.75), CONTENT_WIDTH, Inches(0.4),
               text="Practise using the patterns together:",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_MEDIUM)
    
    drills = [
        ("1", "Point and say:", "นี่คือครูครับ", "nîi khʉʉ khruu khráp", "This is the teacher."),
        ("2", "Ask who:", "เขาคือใครคะ", "khǎo khʉʉ khrai khá", "Who is he/she?"),
        ("3", "Confirm:", "คุณคือนักเรียนใช่ไหมครับ", "khun khʉʉ nák-rian châi mái khráp", "You are a student, right?"),
    ]
    
    y = Inches(1.3)
    colors = [ACCENT_TEAL, ACCENT_GOLD, ACCENT_CLAY]
    for num, instruction, thai, translit, eng in drills:
        row_h = Inches(1.5)
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), row_h,
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            CONTENT_LEFT + Inches(0.15), y + Inches(0.3), Inches(0.45), Inches(0.45)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors[int(num)-1]
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.name = FONT_LATIN
        run.font.size = Pt(16)
        run.font.color.rgb = WHITE
        run.font.bold = True
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.8), y + Inches(0.05),
                   Inches(6.0), Inches(0.3),
                   text=instruction, font_name=FONT_LATIN, font_size=SIZE_LABEL,
                   font_color=colors[int(num)-1], bold=True)
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.8), y + Inches(0.35),
                   Inches(6.0), Inches(0.5),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
                   font_color=INK_DARK, bold=True)
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.8), y + Inches(0.85),
                   Inches(6.0), Inches(0.3),
                   text=f"{translit}  —  {eng}",
                   font_name=FONT_TRANSLIT, font_size=Pt(14),
                   font_color=INK_MEDIUM)
        
        y += row_h + Inches(0.1)
    
    # ── Slide 8: Roleplay ───────────────────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, None, "Roleplay: Who is that?")
    
    add_textbox(slide, CONTENT_LEFT, Inches(0.7), CONTENT_WIDTH, Inches(0.4),
               text="At a class — asking about people around you",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_MEDIUM, italic=True)
    
    turns = [
        ("A", "สวัสดีค่ะ เขาคือใครคะ", "sà-wàt-dii khâ khǎo khʉʉ khrai khá", "Hello. Who is he/she?", False),
        ("B", "นี่คือครูครับ", "nîi khʉʉ khruu khráp", "This is the teacher.", True),
        ("A", "คุณคือนักเรียนใช่ไหมคะ", "khun khʉʉ nák-rian châi mái khá", "You are a student, right?", False),
        ("B", "ใช่ครับ ผมคือนักเรียนครับ", "châi khráp phǒm khʉʉ nák-rian khráp", "Yes. I am a student.", True),
    ]
    
    y = Inches(1.2)
    for speaker, thai, translit, english, is_learner in turns:
        h = add_dialogue_turn(slide, CONTENT_LEFT, y, Inches(7.2),
                             speaker, thai, translit, english,
                             is_learner=is_learner)
        y += h + Inches(0.1)
    
    # ── Slide 9: Recap ──────────────────────────────────────────────────
    add_recap_slide(prs, [
        "Use ฉัน, ผม, คุณ, and เขา for basic conversations",
        "Ask ใคร to find out who someone is",
        "Use นี่คือ... to identify people",
        "Confirm with ใช่ไหม and answer with ใช่ or ไม่ใช่",
    ])
    
    # ── Slide 10: Closing ───────────────────────────────────────────────
    add_closing_slide(
        prs,
        lesson_id="M01-L003",
        message="You can now ask who someone is\nand identify people around you.",
        takeaway_lines=[
            "Pronouns + คือ + ใคร = a working identity toolkit.",
            "Next lesson: numbers 0 to 10.",
        ]
    )
    
    output_path = "/home/ubuntu/immersion-thai-slides/M01-L003-deck.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")
    return prs


if __name__ == "__main__":
    build_l003()
