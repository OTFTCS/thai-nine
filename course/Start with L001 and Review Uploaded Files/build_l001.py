"""
Build M01-L001 — Hello and Basic Courtesy
~12 slides following the brief, spoken script, and visual script.
"""

import sys
sys.path.insert(0, '/home/ubuntu/immersion-thai-slides')

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from design_system import *

def build_l001():
    prs = create_presentation()
    
    # ── Slide 1: Lesson Opener ──────────────────────────────────────────
    add_lesson_opener(
        prs,
        lesson_id="M01-L001",
        lesson_title="Hello and Basic Courtesy",
        module_title="Module 1 — First Contact and Courtesy",
        level="A0"
    )
    
    # ── Slide 2: What You Will Learn ────────────────────────────────────
    add_objectives_slide(prs, [
        "Say hello politely with a correct polite ending",
        "Use ขอบคุณ and ขอโทษ in the right social moments",
        "Distinguish ครับ, ค่ะ, and คะ at a practical level",
        "Answer with ใช่ or ไม่ใช่ and soften with ไม่เป็นไร",
        "Deliver a short first-contact exchange with politeness",
    ])
    
    # ── Slide 3: Greeting Focus — สวัสดี ────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "1", "Say hello the Thai way")
    
    # Main phrase card
    add_phrase_card(
        slide, CONTENT_LEFT, Inches(1.5), Inches(6.5),
        thai="สวัสดี",
        translit="sà-wàt-dii",
        english="hello",
        thai_size=SIZE_THAI_LARGE,
        accent_color=ACCENT_GOLD
    )
    
    # Explanation text
    add_textbox(slide, CONTENT_LEFT, Inches(3.6), CONTENT_WIDTH, Inches(1.0),
               text="This is your safe, polite greeting for a first meeting.\nBut real Thai sounds warmer when a polite ending is attached.",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_MEDIUM)
    
    # ── Slide 4: Greeting Stack — with polite endings ───────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "1", "The greeting stack")
    
    # Subtitle
    add_textbox(slide, CONTENT_LEFT, Inches(0.9), CONTENT_WIDTH, Inches(0.4),
               text="Learn the greeting as a full social line, not just the base word.",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_MEDIUM)
    
    # Three cards stacked
    cards = [
        ("สวัสดี", "sà-wàt-dii", "hello — base form"),
        ("สวัสดีครับ", "sà-wàt-dii khráp", "hello — male speaker"),
        ("สวัสดีค่ะ", "sà-wàt-dii khâ", "hello — female speaker"),
    ]
    
    y = Inches(1.5)
    colors = [ACCENT_GOLD, ACCENT_TEAL, ACCENT_CLAY]
    for i, (thai, translit, eng) in enumerate(cards):
        # Card background
        card_h = Inches(1.35)
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.0), card_h,
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        
        # Accent bar
        add_accent_bar(slide, CONTENT_LEFT + Inches(0.12), y + Inches(0.2),
                      height=Inches(0.95), color=colors[i])
        
        # Thai
        add_textbox(slide, CONTENT_LEFT + Inches(0.4), y + Inches(0.1),
                   Inches(5.5), Inches(0.55),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_MED,
                   font_color=INK_DARK, bold=True)
        
        # Translit + English
        add_textbox(slide, CONTENT_LEFT + Inches(0.4), y + Inches(0.7),
                   Inches(5.5), Inches(0.5),
                   text=f"{translit}  —  {eng}",
                   font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
                   font_color=INK_MEDIUM)
        
        y += card_h + Inches(0.15)
    
    # ── Slide 5: Thank & Apologise — Contrast Board ────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "2", "Thank, apologise, and soften")
    
    # Two-column layout
    col_w = Inches(3.4)
    col1_x = CONTENT_LEFT
    col2_x = CONTENT_LEFT + col_w + Inches(0.3)
    
    # Column headers
    # Left column: Close warmly
    header_bg_1 = add_rounded_rect(slide, col1_x, Inches(1.3), col_w, Inches(0.45),
                                   fill_color=ACCENT_TEAL, border_color=ACCENT_TEAL)
    add_textbox(slide, col1_x, Inches(1.32), col_w, Inches(0.4),
               text="Close warmly", font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Right column: Interrupt or repair
    header_bg_2 = add_rounded_rect(slide, col2_x, Inches(1.3), col_w, Inches(0.45),
                                   fill_color=ACCENT_CLAY, border_color=ACCENT_CLAY)
    add_textbox(slide, col2_x, Inches(1.32), col_w, Inches(0.4),
               text="Interrupt or repair", font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Left card: ขอบคุณ
    card1 = add_rounded_rect(slide, col1_x, Inches(1.95), col_w, Inches(1.6),
                            fill_color=CARD_BG, border_color=CARD_BORDER)
    add_accent_bar(slide, col1_x + Inches(0.1), Inches(2.15),
                  height=Inches(1.2), color=ACCENT_TEAL)
    add_textbox(slide, col1_x + Inches(0.35), Inches(2.05), col_w - Inches(0.5), Inches(0.55),
               text="ขอบคุณ", font_name=FONT_THAI, font_size=SIZE_THAI_MED,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, col1_x + Inches(0.35), Inches(2.6), col_w - Inches(0.5), Inches(0.3),
               text="khàawp-khun", font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM, italic=True)
    add_textbox(slide, col1_x + Inches(0.35), Inches(2.95), col_w - Inches(0.5), Inches(0.3),
               text="thank you", font_name=FONT_LATIN, font_size=SIZE_ENGLISH,
               font_color=INK_LIGHT)
    
    # Right card: ขอโทษ
    card2 = add_rounded_rect(slide, col2_x, Inches(1.95), col_w, Inches(1.6),
                            fill_color=CARD_BG, border_color=CARD_BORDER)
    add_accent_bar(slide, col2_x + Inches(0.1), Inches(2.15),
                  height=Inches(1.2), color=ACCENT_CLAY)
    add_textbox(slide, col2_x + Inches(0.35), Inches(2.05), col_w - Inches(0.5), Inches(0.55),
               text="ขอโทษ", font_name=FONT_THAI, font_size=SIZE_THAI_MED,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, col2_x + Inches(0.35), Inches(2.6), col_w - Inches(0.5), Inches(0.3),
               text="khǎaw-thôot", font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM, italic=True)
    add_textbox(slide, col2_x + Inches(0.35), Inches(2.95), col_w - Inches(0.5), Inches(0.3),
               text="excuse me / sorry", font_name=FONT_LATIN, font_size=SIZE_ENGLISH,
               font_color=INK_LIGHT)
    
    # Arrow + ไม่เป็นไร response flow
    add_horizontal_line(slide, CONTENT_LEFT, Inches(3.85), Inches(7.1),
                       color=DIVIDER_COLOR, weight=Pt(1))
    
    # Flow label
    add_textbox(slide, CONTENT_LEFT, Inches(4.05), Inches(7.0), Inches(0.35),
               text="Natural reply to both →",
               font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=INK_LIGHT, italic=True)
    
    # ไม่เป็นไร card (centered, spanning both columns)
    mai_card = add_rounded_rect(slide, CONTENT_LEFT + Inches(1.0), Inches(4.5),
                               Inches(5.1), Inches(1.5),
                               fill_color=HIGHLIGHT_BG, border_color=ACCENT_SOFT_GOLD)
    add_accent_bar(slide, CONTENT_LEFT + Inches(1.15), Inches(4.7),
                  height=Inches(1.1), color=ACCENT_GOLD)
    add_textbox(slide, CONTENT_LEFT + Inches(1.5), Inches(4.6), Inches(4.4), Inches(0.55),
               text="ไม่เป็นไร", font_name=FONT_THAI, font_size=SIZE_THAI_MED,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(1.5), Inches(5.15), Inches(4.4), Inches(0.3),
               text="mâi bpen rai  —  it is okay / no problem",
               font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM)
    add_textbox(slide, CONTENT_LEFT + Inches(1.5), Inches(5.5), Inches(4.4), Inches(0.3),
               text="Smooths the moment after thanks or a small apology",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_LIGHT, italic=True)
    
    # ── Slide 6: Social Flow — when to use each ────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "2", "How these phrases work together")
    
    # Explanation text
    add_textbox(slide, CONTENT_LEFT, Inches(1.2), CONTENT_WIDTH, Inches(0.5),
               text="Each phrase has a social job in the interaction:",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_MEDIUM)
    
    # Three function rows
    functions = [
        ("ขอบคุณ", "khàawp-khun", "Closes a moment warmly", "After help, service, or a small kindness", ACCENT_TEAL),
        ("ขอโทษ", "khǎaw-thôot", "Opens a repair", "To interrupt gently or apologise", ACCENT_CLAY),
        ("ไม่เป็นไร", "mâi bpen rai", "Smooths everything out", "Answers a thank-you or small apology", ACCENT_GOLD),
    ]
    
    y = Inches(1.9)
    for thai, translit, function, when, color in functions:
        row_h = Inches(1.35)
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), row_h,
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        add_accent_bar(slide, CONTENT_LEFT + Inches(0.1), y + Inches(0.2),
                      height=Inches(0.95), color=color)
        
        # Thai + translit
        add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.1),
                   Inches(3.0), Inches(0.5),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
                   font_color=INK_DARK, bold=True)
        add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.55),
                   Inches(3.0), Inches(0.3),
                   text=translit, font_name=FONT_TRANSLIT, font_size=Pt(14),
                   font_color=INK_MEDIUM, italic=True)
        
        # Function + when
        add_textbox(slide, CONTENT_LEFT + Inches(3.2), y + Inches(0.15),
                   Inches(3.8), Inches(0.4),
                   text=function, font_name=FONT_LATIN, font_size=SIZE_BODY,
                   font_color=INK_DARK, bold=True)
        add_textbox(slide, CONTENT_LEFT + Inches(3.2), y + Inches(0.6),
                   Inches(3.8), Inches(0.4),
                   text=when, font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
                   font_color=INK_LIGHT)
        
        y += row_h + Inches(0.15)
    
    # ── Slide 7: Polite Endings — Choice System ────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "3", "Pick the right polite ending")
    
    add_textbox(slide, CONTENT_LEFT, Inches(1.1), CONTENT_WIDTH, Inches(0.4),
               text="Thai polite endings do real work — they shape the tone of the whole line.",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_MEDIUM)
    
    # Three choice chips
    endings = [
        ("ครับ", "khráp", "male polite ending", "Use for statements and questions", ACCENT_TEAL),
        ("ค่ะ", "khâ", "female statement ending", "Use for statements and answers", ACCENT_CLAY),
        ("คะ", "khá", "female question ending", "Use when asking questions", ACCENT_GOLD),
    ]
    
    y = Inches(1.7)
    for thai, translit, eng, usage, color in endings:
        chip_h = Inches(1.2)
        chip = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), chip_h,
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        
        # Color tag on left
        tag = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            CONTENT_LEFT + Inches(0.15), y + Inches(0.2),
            Inches(1.6), Inches(0.8)
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()
        tag.adjustments[0] = 0.15
        
        # Thai in tag
        tf = tag.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = thai
        run.font.name = FONT_THAI
        run.font.size = SIZE_THAI_MED
        run.font.color.rgb = WHITE
        run.font.bold = True
        
        # Translit + English + Usage
        add_textbox(slide, CONTENT_LEFT + Inches(2.0), y + Inches(0.1),
                   Inches(5.0), Inches(0.35),
                   text=f"{translit}  —  {eng}",
                   font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
                   font_color=INK_DARK, bold=True)
        add_textbox(slide, CONTENT_LEFT + Inches(2.0), y + Inches(0.5),
                   Inches(5.0), Inches(0.35),
                   text=usage,
                   font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
                   font_color=INK_MEDIUM)
        
        y += chip_h + Inches(0.15)
    
    # Practical rule note
    add_textbox(slide, CONTENT_LEFT, y + Inches(0.15), CONTENT_WIDTH, Inches(0.5),
               text="Beginner rule: learn ครับ and ค่ะ first for safe polite statements.",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=ACCENT_CLAY, bold=True)
    
    # ── Slide 8: Polite Endings — Substitution Drill ───────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "3", "Substitution drill")
    
    add_textbox(slide, CONTENT_LEFT, Inches(1.1), CONTENT_WIDTH, Inches(0.4),
               text="Say each line with the correct polite ending:",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_MEDIUM)
    
    # Drill lines
    drills = [
        ("สวัสดี___", "sà-wàt-dii ___", "Add ครับ or ค่ะ"),
        ("ขอบคุณ___", "khàawp-khun ___", "Add ครับ or ค่ะ"),
        ("ขอโทษ___", "khǎaw-thôot ___", "Add ครับ or ค่ะ"),
        ("ใช่___", "châi ___", "Add ครับ or ค่ะ"),
        ("ไม่เป็นไร___", "mâi bpen rai ___", "Add ครับ or ค่ะ"),
    ]
    
    y = Inches(1.7)
    for i, (thai, translit, instruction) in enumerate(drills):
        row_h = Inches(0.85)
        
        # Alternating background
        if i % 2 == 0:
            bg = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), row_h,
                                 fill_color=HIGHLIGHT_BG, border_color=RGBColor(0xF0, 0xE8, 0xD8))
        
        # Number
        add_textbox(slide, CONTENT_LEFT + Inches(0.15), y + Inches(0.15),
                   Inches(0.4), Inches(0.4),
                   text=str(i + 1), font_name=FONT_LATIN, font_size=SIZE_BODY,
                   font_color=ACCENT_GOLD, bold=True)
        
        # Thai
        add_textbox(slide, CONTENT_LEFT + Inches(0.5), y + Inches(0.08),
                   Inches(3.5), Inches(0.45),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
                   font_color=INK_DARK, bold=True)
        
        # Translit
        add_textbox(slide, CONTENT_LEFT + Inches(0.5), y + Inches(0.5),
                   Inches(3.5), Inches(0.3),
                   text=translit, font_name=FONT_TRANSLIT, font_size=Pt(14),
                   font_color=INK_MEDIUM, italic=True)
        
        # Instruction
        add_textbox(slide, CONTENT_LEFT + Inches(4.5), y + Inches(0.2),
                   Inches(2.5), Inches(0.4),
                   text=instruction, font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
                   font_color=ACCENT_CLAY, italic=True)
        
        y += row_h + Inches(0.08)
    
    # ── Slide 9: Yes / No / No Problem — Response Ladder ───────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "4", "Answer yes, no, and no problem")
    
    add_textbox(slide, CONTENT_LEFT, Inches(1.1), CONTENT_WIDTH, Inches(0.4),
               text="Three response tools — two answer the idea, one softens the moment.",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_MEDIUM)
    
    # Response ladder — two answer cards + one softener card
    # ใช่
    y = Inches(1.7)
    
    # "Answers" group label
    answers_label = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(1.8), Inches(0.35),
                                    fill_color=ACCENT_TEAL, border_color=ACCENT_TEAL)
    add_textbox(slide, CONTENT_LEFT, y + Inches(0.02), Inches(1.8), Inches(0.3),
               text="Direct answers", font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    y += Inches(0.5)
    
    # ใช่ card
    card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), Inches(1.2),
                           fill_color=CARD_BG, border_color=CARD_BORDER)
    add_accent_bar(slide, CONTENT_LEFT + Inches(0.1), y + Inches(0.15),
                  height=Inches(0.9), color=ACCENT_TEAL)
    add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.05),
               Inches(2.5), Inches(0.5),
               text="ใช่", font_name=FONT_THAI, font_size=SIZE_THAI_MED,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.55),
               Inches(2.5), Inches(0.3),
               text="châi", font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM, italic=True)
    add_textbox(slide, CONTENT_LEFT + Inches(3.0), y + Inches(0.15),
               Inches(4.0), Inches(0.4),
               text="yes / that is right",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(3.0), y + Inches(0.55),
               Inches(4.0), Inches(0.4),
               text="Confirms the idea",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_LIGHT)
    
    y += Inches(1.35)
    
    # ไม่ใช่ card
    card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), Inches(1.2),
                           fill_color=CARD_BG, border_color=CARD_BORDER)
    add_accent_bar(slide, CONTENT_LEFT + Inches(0.1), y + Inches(0.15),
                  height=Inches(0.9), color=ACCENT_TEAL)
    add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.05),
               Inches(2.5), Inches(0.5),
               text="ไม่ใช่", font_name=FONT_THAI, font_size=SIZE_THAI_MED,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.55),
               Inches(2.5), Inches(0.3),
               text="mâi châi", font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM, italic=True)
    add_textbox(slide, CONTENT_LEFT + Inches(3.0), y + Inches(0.15),
               Inches(4.0), Inches(0.4),
               text="no / not right",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(3.0), y + Inches(0.55),
               Inches(4.0), Inches(0.4),
               text="Corrects the idea — not aggressive",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_LIGHT)
    
    y += Inches(1.55)
    
    # "Softener" group label
    soft_label = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(2.0), Inches(0.35),
                                 fill_color=ACCENT_GOLD, border_color=ACCENT_GOLD)
    add_textbox(slide, CONTENT_LEFT, y + Inches(0.02), Inches(2.0), Inches(0.3),
               text="Social softener", font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    y += Inches(0.5)
    
    # ไม่เป็นไร card
    card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), Inches(1.2),
                           fill_color=HIGHLIGHT_BG, border_color=ACCENT_SOFT_GOLD)
    add_accent_bar(slide, CONTENT_LEFT + Inches(0.1), y + Inches(0.15),
                  height=Inches(0.9), color=ACCENT_GOLD)
    add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.05),
               Inches(2.5), Inches(0.5),
               text="ไม่เป็นไร", font_name=FONT_THAI, font_size=SIZE_THAI_MED,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.55),
               Inches(2.5), Inches(0.3),
               text="mâi bpen rai", font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM, italic=True)
    add_textbox(slide, CONTENT_LEFT + Inches(3.0), y + Inches(0.15),
               Inches(4.0), Inches(0.4),
               text="it is okay / no problem",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(3.0), y + Inches(0.55),
               Inches(4.0), Inches(0.4),
               text="Softens the situation — not a yes/no answer",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=ACCENT_CLAY, italic=True)
    
    # ── Slide 10: Roleplay — Reception Desk ────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, None, "Roleplay: First check-in")
    
    add_textbox(slide, CONTENT_LEFT, Inches(1.0), CONTENT_WIDTH, Inches(0.4),
               text="At a reception desk — practise the full exchange",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_MEDIUM, italic=True)
    
    # Dialogue turns
    turns = [
        ("Learner", "ขอโทษครับ สวัสดีครับ", "khǎaw-thôot khráp sà-wàt-dii khráp", "Excuse me, hello.", True),
        ("Staff", "สวัสดีค่ะ", "sà-wàt-dii khâ", "Hello.", False),
        ("Learner", "ที่นี่ใช่ไหมครับ", "thîi-nîi châi mái khráp", "Is this the right place?", True),
        ("Staff", "ใช่ค่ะ", "châi khâ", "Yes.", False),
        ("Learner", "ขอบคุณครับ", "khàawp-khun khráp", "Thank you.", True),
        ("Staff", "ไม่เป็นไรค่ะ", "mâi bpen rai khâ", "No problem.", False),
    ]
    
    y = Inches(1.6)
    for speaker, thai, translit, english, is_learner in turns:
        h = add_dialogue_turn(slide, CONTENT_LEFT, y, Inches(7.2),
                             speaker, thai, translit, english,
                             is_learner=is_learner)
        y += h + Inches(0.05)
    
    # ── Slide 11: Recap ────────────────────────────────────────────────
    add_recap_slide(prs, [
        "Start with สวัสดี and always add a polite ending",
        "Use ขอบคุณ to close warmly and ขอโทษ to repair politely",
        "Learn polite endings as social meaning, not decoration",
        "Use ใช่ and ไม่ใช่ for answers, ไม่เป็นไร to calm the moment",
        "Practise the full exchange — greeting to thank-you to no-problem",
    ])
    
    # ── Slide 12: Closing ──────────────────────────────────────────────
    add_closing_slide(
        prs,
        lesson_id="M01-L001",
        message="You can now greet, thank,\napologise, and respond politely.",
        takeaway_lines=[
            "You have the smallest courtesy toolkit that already works.",
            "Next lesson: names and countries — introducing yourself.",
        ]
    )
    
    # Save
    output_path = "/home/ubuntu/immersion-thai-slides/M01-L001-deck.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")
    return prs


if __name__ == "__main__":
    build_l001()
