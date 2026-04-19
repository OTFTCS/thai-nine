"""
Build M01-L004 — Numbers 0 to 10
~12 slides based on blueprint CSV content.
"""

import sys
sys.path.insert(0, '/home/ubuntu/immersion-thai-slides')

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from design_system import *

def build_l004():
    prs = create_presentation()
    
    # ── Slide 1: Lesson Opener ──────────────────────────────────────────
    add_lesson_opener(
        prs,
        lesson_id="M01-L004",
        lesson_title="Numbers 0 to 10",
        module_title="Module 1 — First Contact and Courtesy",
        level="A0"
    )
    
    # ── Slide 2: What You Will Learn ────────────────────────────────────
    add_objectives_slide(prs, [
        "Say Thai digits from 0 to 10",
        "Use numbers for age, room numbers, and basic counting",
        "Ask เบอร์อะไร (what number?)",
        "Say your age with ฉันอายุ...ปี",
    ])
    
    # ── Slide 3: Numbers 0–5 ────────────────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "1", "Numbers 0 to 5")
    
    numbers_1 = [
        ("ศูนย์", "sǔun", "0 — zero"),
        ("หนึ่ง", "nʉ̀ng", "1 — one"),
        ("สอง", "sǎawng", "2 — two"),
        ("สาม", "sǎam", "3 — three"),
        ("สี่", "sìi", "4 — four"),
        ("ห้า", "hâa", "5 — five"),
    ]
    
    y = Inches(0.9)
    for i, (thai, translit, eng) in enumerate(numbers_1):
        row_h = Inches(0.85)
        bg_color = HIGHLIGHT_BG if i % 2 == 0 else CARD_BG
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), row_h,
                               fill_color=bg_color, border_color=CARD_BORDER)
        
        # Number digit
        digit_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            CONTENT_LEFT + Inches(0.12), y + Inches(0.12), Inches(0.55), Inches(0.55)
        )
        digit_circle.fill.solid()
        digit_circle.fill.fore_color.rgb = ACCENT_GOLD
        digit_circle.line.fill.background()
        tf = digit_circle.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i)
        run.font.name = FONT_LATIN
        run.font.size = Pt(20)
        run.font.color.rgb = WHITE
        run.font.bold = True
        
        # Thai
        add_textbox(slide, CONTENT_LEFT + Inches(0.85), y + Inches(0.05),
                   Inches(2.5), Inches(0.5),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
                   font_color=INK_DARK, bold=True)
        
        # Translit
        add_textbox(slide, CONTENT_LEFT + Inches(3.3), y + Inches(0.15),
                   Inches(2.0), Inches(0.35),
                   text=translit, font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
                   font_color=INK_MEDIUM, italic=True)
        
        # English
        add_textbox(slide, CONTENT_LEFT + Inches(5.3), y + Inches(0.15),
                   Inches(1.8), Inches(0.35),
                   text=eng, font_name=FONT_LATIN, font_size=SIZE_ENGLISH,
                   font_color=INK_LIGHT)
        
        y += row_h + Inches(0.06)
    
    # ── Slide 4: Numbers 6–10 ───────────────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "1", "Numbers 6 to 10")
    
    numbers_2 = [
        ("หก", "hòk", "6 — six"),
        ("เจ็ด", "jèt", "7 — seven"),
        ("แปด", "bpàaet", "8 — eight"),
        ("เก้า", "gâo", "9 — nine"),
        ("สิบ", "sìp", "10 — ten"),
    ]
    
    y = Inches(0.9)
    for i, (thai, translit, eng) in enumerate(numbers_2):
        row_h = Inches(0.85)
        bg_color = HIGHLIGHT_BG if i % 2 == 0 else CARD_BG
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), row_h,
                               fill_color=bg_color, border_color=CARD_BORDER)
        
        digit_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            CONTENT_LEFT + Inches(0.12), y + Inches(0.12), Inches(0.55), Inches(0.55)
        )
        digit_circle.fill.solid()
        digit_circle.fill.fore_color.rgb = ACCENT_TEAL
        digit_circle.line.fill.background()
        tf = digit_circle.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 6)
        run.font.name = FONT_LATIN
        run.font.size = Pt(20)
        run.font.color.rgb = WHITE
        run.font.bold = True
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.85), y + Inches(0.05),
                   Inches(2.5), Inches(0.5),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
                   font_color=INK_DARK, bold=True)
        
        add_textbox(slide, CONTENT_LEFT + Inches(3.3), y + Inches(0.15),
                   Inches(2.0), Inches(0.35),
                   text=translit, font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
                   font_color=INK_MEDIUM, italic=True)
        
        add_textbox(slide, CONTENT_LEFT + Inches(5.3), y + Inches(0.15),
                   Inches(1.8), Inches(0.35),
                   text=eng, font_name=FONT_LATIN, font_size=SIZE_ENGLISH,
                   font_color=INK_LIGHT)
        
        y += row_h + Inches(0.06)
    
    # Practice note
    add_textbox(slide, CONTENT_LEFT, y + Inches(0.3), CONTENT_WIDTH, Inches(0.5),
               text="Practise: count from 0 to 10 slowly, then speed up.",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=ACCENT_CLAY, bold=True)
    
    # ── Slide 5: What Number? — เบอร์อะไร ──────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "2", "Ask what number")
    
    # Key chunk
    add_phrase_card(
        slide, CONTENT_LEFT, Inches(1.0), Inches(6.5),
        thai="เบอร์อะไร",
        translit="bəə à-rai",
        english="what number?",
        thai_size=SIZE_THAI_LARGE,
        accent_color=ACCENT_GOLD
    )
    
    add_textbox(slide, CONTENT_LEFT, Inches(3.1), CONTENT_WIDTH, Inches(0.5),
               text="Use this to ask for a room number, phone number, or any number.",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_MEDIUM)
    
    # Example
    ex_card = add_rounded_rect(slide, CONTENT_LEFT, Inches(3.8), Inches(7.2), Inches(1.3),
                              fill_color=CARD_BG, border_color=CARD_BORDER)
    add_accent_bar(slide, CONTENT_LEFT + Inches(0.12), Inches(4.0),
                  height=Inches(0.9), color=ACCENT_TEAL)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(3.85), Inches(1.5), Inches(0.3),
               text="Example", font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=ACCENT_TEAL, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(4.15), Inches(6.0), Inches(0.5),
               text="ห้องเบอร์อะไรครับ", font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(4.6), Inches(6.0), Inches(0.3),
               text="hâwng bəə à-rai khráp  —  What room number?",
               font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM)
    
    # ── Slide 6: Age Pattern — ฉันอายุ...ปี ─────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "2", "Say your age")
    
    # Pattern card
    pattern = add_rounded_rect(slide, CONTENT_LEFT, Inches(1.0), Inches(7.2), Inches(1.6),
                              fill_color=HIGHLIGHT_BG, border_color=ACCENT_SOFT_GOLD)
    add_accent_bar(slide, CONTENT_LEFT + Inches(0.12), Inches(1.2),
                  height=Inches(1.2), color=ACCENT_GOLD)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(1.05), Inches(1.5), Inches(0.3),
               text="Pattern", font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=ACCENT_GOLD, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(1.4), Inches(6.0), Inches(0.55),
               text="ฉัน/ผม  อายุ  ___  ปี", font_name=FONT_THAI, font_size=SIZE_THAI_MED,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(1.95), Inches(6.0), Inches(0.3),
               text="chǎn/phǒm  aa-yú  ___  bpii  —  I am ___ years old",
               font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM)
    
    # Key words
    words = [
        ("อายุ", "aa-yú", "age"),
        ("ปี", "bpii", "year"),
    ]
    
    y = Inches(2.9)
    for thai, translit, eng in words:
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(3.4), Inches(1.0),
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        add_accent_bar(slide, CONTENT_LEFT + Inches(0.1), y + Inches(0.15),
                      height=Inches(0.7), color=ACCENT_TEAL)
        add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.05),
                   Inches(2.5), Inches(0.45),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
                   font_color=INK_DARK, bold=True)
        add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.5),
                   Inches(2.5), Inches(0.3),
                   text=f"{translit}  —  {eng}",
                   font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
                   font_color=INK_MEDIUM)
        y += Inches(1.15)
    
    # Examples
    add_textbox(slide, CONTENT_LEFT + Inches(3.8), Inches(2.95), Inches(3.5), Inches(0.3),
               text="Examples:", font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=ACCENT_CLAY, bold=True)
    
    examples = [
        ("ฉันอายุห้าปี", "5 years old"),
        ("ผมอายุเก้าปี", "9 years old"),
        ("ฉันอายุสิบปี", "10 years old"),
    ]
    
    ey = Inches(3.35)
    for thai, eng in examples:
        add_textbox(slide, CONTENT_LEFT + Inches(3.8), ey, Inches(3.5), Inches(0.4),
                   text=thai, font_name=FONT_THAI, font_size=Pt(22),
                   font_color=INK_DARK, bold=True)
        add_textbox(slide, CONTENT_LEFT + Inches(6.0), ey + Inches(0.05), Inches(1.5), Inches(0.3),
                   text=eng, font_name=FONT_LATIN, font_size=Pt(13),
                   font_color=INK_LIGHT)
        ey += Inches(0.5)
    
    # ── Slide 7: Counting Drill ─────────────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "3", "Counting drill")
    
    add_textbox(slide, CONTENT_LEFT, Inches(0.75), CONTENT_WIDTH, Inches(0.4),
               text="Say each number aloud — check your pronunciation:",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_MEDIUM)
    
    # Grid layout: 3 columns, 4 rows
    all_nums = [
        ("ศูนย์", "0"), ("หนึ่ง", "1"), ("สอง", "2"),
        ("สาม", "3"), ("สี่", "4"), ("ห้า", "5"),
        ("หก", "6"), ("เจ็ด", "7"), ("แปด", "8"),
        ("เก้า", "9"), ("สิบ", "10"), ("", ""),
    ]
    
    col_w = Inches(2.2)
    row_h = Inches(1.1)
    start_y = Inches(1.3)
    
    colors_grid = [ACCENT_GOLD, ACCENT_TEAL, ACCENT_CLAY]
    
    for idx, (thai, digit) in enumerate(all_nums):
        if not thai:
            continue
        row = idx // 3
        col = idx % 3
        x = CONTENT_LEFT + col * (col_w + Inches(0.15))
        y = start_y + row * (row_h + Inches(0.1))
        
        card = add_rounded_rect(slide, x, y, col_w, row_h,
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        
        # Digit
        add_textbox(slide, x + Inches(0.1), y + Inches(0.05),
                   Inches(0.5), Inches(0.4),
                   text=digit, font_name=FONT_LATIN, font_size=Pt(22),
                   font_color=colors_grid[col], bold=True)
        
        # Thai
        add_textbox(slide, x + Inches(0.55), y + Inches(0.1),
                   Inches(1.5), Inches(0.5),
                   text=thai, font_name=FONT_THAI, font_size=Pt(26),
                   font_color=INK_DARK, bold=True)
    
    # ── Slide 8: Number in Context — Mini Scenarios ─────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "3", "Numbers in real situations")
    
    scenarios = [
        ("Room number", "ห้องเบอร์สามครับ", "hâwng bəə sǎam khráp", "Room number 3.", ACCENT_GOLD),
        ("Age", "ฉันอายุเจ็ดปีค่ะ", "chǎn aa-yú jèt bpii khâ", "I am 7 years old.", ACCENT_TEAL),
        ("Counting", "หนึ่ง สอง สาม สี่ ห้า", "nʉ̀ng sǎawng sǎam sìi hâa", "1, 2, 3, 4, 5", ACCENT_CLAY),
    ]
    
    y = Inches(0.9)
    for label, thai, translit, eng, color in scenarios:
        card_h = Inches(1.6)
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), card_h,
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        
        # Label pill
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            CONTENT_LEFT + Inches(0.15), y + Inches(0.15), Inches(1.8), Inches(0.32)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = color
        pill.line.fill.background()
        pill.adjustments[0] = 0.3
        tf = pill.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.name = FONT_LATIN
        run.font.size = SIZE_LABEL
        run.font.color.rgb = WHITE
        run.font.bold = True
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.15), y + Inches(0.55),
                   Inches(6.5), Inches(0.5),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
                   font_color=INK_DARK, bold=True)
        add_textbox(slide, CONTENT_LEFT + Inches(0.15), y + Inches(1.0),
                   Inches(6.5), Inches(0.3),
                   text=f"{translit}  —  {eng}",
                   font_name=FONT_TRANSLIT, font_size=Pt(14),
                   font_color=INK_MEDIUM)
        
        y += card_h + Inches(0.12)
    
    # ── Slide 9: Roleplay ───────────────────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, None, "Roleplay: What number?")
    
    add_textbox(slide, CONTENT_LEFT, Inches(0.7), CONTENT_WIDTH, Inches(0.4),
               text="At a hotel check-in — asking about room numbers and age",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_MEDIUM, italic=True)
    
    turns = [
        ("Staff", "สวัสดีค่ะ ห้องเบอร์อะไรคะ", "sà-wàt-dii khâ hâwng bəə à-rai khá", "Hello. What room number?", False),
        ("Guest", "ห้องเบอร์ห้าครับ", "hâwng bəə hâa khráp", "Room number 5.", True),
        ("Staff", "คุณอายุเท่าไหร่คะ", "khun aa-yú thâo-rài khá", "How old are you?", False),
        ("Guest", "ผมอายุสิบปีครับ", "phǒm aa-yú sìp bpii khráp", "I am 10 years old.", True),
        ("Staff", "ขอบคุณค่ะ", "khàawp-khun khâ", "Thank you.", False),
    ]
    
    y = Inches(1.2)
    for speaker, thai, translit, english, is_learner in turns:
        h = add_dialogue_turn(slide, CONTENT_LEFT, y, Inches(7.2),
                             speaker, thai, translit, english,
                             is_learner=is_learner)
        y += h + Inches(0.05)
    
    # ── Slide 10: Quick Check ───────────────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, None, "Quick check")
    
    add_textbox(slide, CONTENT_LEFT, Inches(0.75), CONTENT_WIDTH, Inches(0.4),
               text="Can you say these numbers in Thai?",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_MEDIUM)
    
    quiz_items = [
        ("1.", "3", "→  ___"),
        ("2.", "7", "→  ___"),
        ("3.", "10", "→  ___"),
        ("4.", "0", "→  ___"),
        ("5.", "Say your age:", "ฉัน/ผม อายุ ___ ปี"),
    ]
    
    y = Inches(1.4)
    for num, prompt, answer_hint in quiz_items:
        row_h = Inches(0.8)
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), row_h,
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.15), y + Inches(0.15),
                   Inches(0.4), Inches(0.4),
                   text=num, font_name=FONT_LATIN, font_size=SIZE_BODY,
                   font_color=ACCENT_GOLD, bold=True)
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.5), y + Inches(0.15),
                   Inches(2.5), Inches(0.4),
                   text=prompt, font_name=FONT_LATIN, font_size=Pt(22),
                   font_color=INK_DARK, bold=True)
        
        add_textbox(slide, CONTENT_LEFT + Inches(3.5), y + Inches(0.15),
                   Inches(3.5), Inches(0.4),
                   text=answer_hint, font_name=FONT_THAI, font_size=SIZE_BODY,
                   font_color=INK_LIGHT)
        
        y += row_h + Inches(0.08)
    
    # ── Slide 11: Recap ─────────────────────────────────────────────────
    add_recap_slide(prs, [
        "Count from ศูนย์ (0) to สิบ (10) confidently",
        "Ask เบอร์อะไร for any number you need",
        "Say your age with ฉัน/ผม อายุ...ปี",
        "Use numbers in real situations: rooms, age, counting",
    ])
    
    # ── Slide 12: Closing ───────────────────────────────────────────────
    add_closing_slide(
        prs,
        lesson_id="M01-L004",
        message="You can now count, ask for numbers,\nand say your age in Thai.",
        takeaway_lines=[
            "0 to 10 — the digits you need for age, rooms, and basic counting.",
            "Module 1 continues with where/what questions next.",
        ]
    )
    
    output_path = "/home/ubuntu/immersion-thai-slides/M01-L004-deck.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")
    return prs


if __name__ == "__main__":
    build_l004()
