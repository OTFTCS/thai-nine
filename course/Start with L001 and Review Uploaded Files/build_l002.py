"""
Build M01-L002 — Names and Countries
~11 slides following the script-master.json content.
"""

import sys
sys.path.insert(0, '/home/ubuntu/immersion-thai-slides')

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from design_system import *

def build_l002():
    prs = create_presentation()
    
    # ── Slide 1: Lesson Opener ──────────────────────────────────────────
    add_lesson_opener(
        prs,
        lesson_id="M01-L002",
        lesson_title="Names and Countries",
        module_title="Module 1 — First Contact and Courtesy",
        level="A0"
    )
    
    # ── Slide 2: What You Will Learn ────────────────────────────────────
    add_objectives_slide(prs, [
        "Ask someone's name with คุณชื่ออะไร",
        "Say your own name with ฉันชื่อ... or ผมชื่อ...",
        "State where you are from with มาจาก",
        "Build a short self-introduction from greeting to country",
    ])
    
    # ── Slide 3: Name Question Focus ────────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "1", "Ask and answer the name question")
    
    # Key word card
    add_phrase_card(
        slide, CONTENT_LEFT, Inches(1.2), Inches(6.5),
        thai="ชื่อ",
        translit="chûue",
        english="name",
        thai_size=SIZE_THAI_LARGE,
        accent_color=ACCENT_GOLD
    )
    
    add_textbox(slide, CONTENT_LEFT, Inches(3.3), CONTENT_WIDTH, Inches(0.5),
               text="This is the core word in every name question and answer.",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_MEDIUM)
    
    # ── Slide 4: Question-Answer Pair ───────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "1", "The name exchange")
    
    add_textbox(slide, CONTENT_LEFT, Inches(0.75), CONTENT_WIDTH, Inches(0.4),
               text="Learn the question and answer as a pair that belongs together.",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_MEDIUM)
    
    # Question card
    q_card = add_rounded_rect(slide, CONTENT_LEFT, Inches(1.3), Inches(7.2), Inches(1.5),
                             fill_color=HIGHLIGHT_BG, border_color=ACCENT_SOFT_GOLD)
    add_accent_bar(slide, CONTENT_LEFT + Inches(0.12), Inches(1.5),
                  height=Inches(1.1), color=ACCENT_GOLD)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(1.35), Inches(1.5), Inches(0.35),
               text="Question", font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=ACCENT_GOLD, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(1.7), Inches(6.0), Inches(0.55),
               text="คุณชื่ออะไร", font_name=FONT_THAI, font_size=SIZE_THAI_MED,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(2.25), Inches(6.0), Inches(0.3),
               text="khun chûue à-rai  —  what is your name?",
               font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM)
    
    # Answer cards
    y = Inches(3.1)
    answers = [
        ("ฉันชื่อ...", "chǎn chûue...", "my name is ...", ACCENT_TEAL),
        ("ผมชื่อ...", "phǒm chûue...", "my name is ... (male)", ACCENT_CLAY),
    ]
    
    for thai, translit, eng, color in answers:
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), Inches(1.3),
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        add_accent_bar(slide, CONTENT_LEFT + Inches(0.12), y + Inches(0.2),
                      height=Inches(0.9), color=color)
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.4), y + Inches(0.05), Inches(1.2), Inches(0.3),
                   text="Answer", font_name=FONT_LATIN, font_size=SIZE_LABEL,
                   font_color=color, bold=True)
        add_textbox(slide, CONTENT_LEFT + Inches(0.4), y + Inches(0.35), Inches(5.5), Inches(0.5),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_MED,
                   font_color=INK_DARK, bold=True)
        add_textbox(slide, CONTENT_LEFT + Inches(0.4), y + Inches(0.85), Inches(5.5), Inches(0.3),
                   text=f"{translit}  —  {eng}",
                   font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
                   font_color=INK_MEDIUM)
        y += Inches(1.5)
    
    # ── Slide 5: Pronouns — ฉัน vs ผม ──────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "2", "Use ฉัน and ผม naturally")
    
    add_textbox(slide, CONTENT_LEFT, Inches(0.75), CONTENT_WIDTH, Inches(0.4),
               text="Choose one pronoun confidently — the line pattern stays the same.",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_MEDIUM)
    
    # Two contrast cards side by side
    col_w = Inches(3.4)
    col1_x = CONTENT_LEFT
    col2_x = CONTENT_LEFT + col_w + Inches(0.3)
    
    # ฉัน card
    card1 = add_rounded_rect(slide, col1_x, Inches(1.3), col_w, Inches(2.5),
                            fill_color=CARD_BG, border_color=CARD_BORDER)
    add_accent_bar(slide, col1_x + Inches(0.1), Inches(1.5),
                  height=Inches(2.1), color=ACCENT_TEAL)
    
    # Label
    label1 = add_rounded_rect(slide, col1_x + Inches(0.3), Inches(1.45), Inches(2.0), Inches(0.32),
                              fill_color=ACCENT_TEAL, border_color=ACCENT_TEAL)
    add_textbox(slide, col1_x + Inches(0.3), Inches(1.47), Inches(2.0), Inches(0.3),
               text="Common / all-purpose", font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, col1_x + Inches(0.3), Inches(1.9), col_w - Inches(0.5), Inches(0.6),
               text="ฉัน", font_name=FONT_THAI, font_size=SIZE_THAI_LARGE,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, col1_x + Inches(0.3), Inches(2.5), col_w - Inches(0.5), Inches(0.3),
               text="chǎn  —  I", font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM)
    add_textbox(slide, col1_x + Inches(0.3), Inches(2.9), col_w - Inches(0.5), Inches(0.5),
               text="ฉันชื่อ...", font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
               font_color=INK_DARK)
    add_textbox(slide, col1_x + Inches(0.3), Inches(3.35), col_w - Inches(0.5), Inches(0.3),
               text="chǎn chûue...", font_name=FONT_TRANSLIT, font_size=Pt(14),
               font_color=INK_LIGHT, italic=True)
    
    # ผม card
    card2 = add_rounded_rect(slide, col2_x, Inches(1.3), col_w, Inches(2.5),
                            fill_color=CARD_BG, border_color=CARD_BORDER)
    add_accent_bar(slide, col2_x + Inches(0.1), Inches(1.5),
                  height=Inches(2.1), color=ACCENT_CLAY)
    
    label2 = add_rounded_rect(slide, col2_x + Inches(0.3), Inches(1.45), Inches(2.0), Inches(0.32),
                              fill_color=ACCENT_CLAY, border_color=ACCENT_CLAY)
    add_textbox(slide, col2_x + Inches(0.3), Inches(1.47), Inches(2.0), Inches(0.3),
               text="Male speaker", font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, col2_x + Inches(0.3), Inches(1.9), col_w - Inches(0.5), Inches(0.6),
               text="ผม", font_name=FONT_THAI, font_size=SIZE_THAI_LARGE,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, col2_x + Inches(0.3), Inches(2.5), col_w - Inches(0.5), Inches(0.3),
               text="phǒm  —  I (male)", font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM)
    add_textbox(slide, col2_x + Inches(0.3), Inches(2.9), col_w - Inches(0.5), Inches(0.5),
               text="ผมชื่อ...", font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
               font_color=INK_DARK)
    add_textbox(slide, col2_x + Inches(0.3), Inches(3.35), col_w - Inches(0.5), Inches(0.3),
               text="phǒm chûue...", font_name=FONT_TRANSLIT, font_size=Pt(14),
               font_color=INK_LIGHT, italic=True)
    
    # Listener reference
    add_textbox(slide, CONTENT_LEFT, Inches(4.2), CONTENT_WIDTH, Inches(0.5),
               text="Addressing someone:  คุณ  khun  —  you",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_MEDIUM)
    
    # ── Slide 6: Country Pattern — มาจาก ────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "3", "Say where you are from")
    
    # Pattern frame card
    frame_card = add_rounded_rect(slide, CONTENT_LEFT, Inches(1.0), Inches(7.2), Inches(1.5),
                                 fill_color=HIGHLIGHT_BG, border_color=ACCENT_SOFT_GOLD)
    add_accent_bar(slide, CONTENT_LEFT + Inches(0.12), Inches(1.2),
                  height=Inches(1.1), color=ACCENT_GOLD)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(1.05), Inches(1.5), Inches(0.3),
               text="Pattern", font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=ACCENT_GOLD, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(1.35), Inches(6.0), Inches(0.55),
               text="ฉัน/ผม  มาจาก  ___", font_name=FONT_THAI, font_size=SIZE_THAI_MED,
               font_color=INK_DARK, bold=True)
    add_textbox(slide, CONTENT_LEFT + Inches(0.4), Inches(1.9), Inches(6.0), Inches(0.3),
               text="chǎn/phǒm  maa-jàak  ___  —  I am from ___",
               font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM)
    
    # Country cards
    countries = [
        ("ไทย", "thai", "Thailand", ACCENT_TEAL),
        ("อังกฤษ", "ang-grìt", "England", ACCENT_CLAY),
        ("อเมริกา", "à-mee-rí-gaa", "America", ACCENT_GOLD),
    ]
    
    y = Inches(2.8)
    for thai, translit, eng, color in countries:
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), Inches(1.1),
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        add_accent_bar(slide, CONTENT_LEFT + Inches(0.1), y + Inches(0.15),
                      height=Inches(0.8), color=color)
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.05),
                   Inches(2.5), Inches(0.5),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_MED,
                   font_color=INK_DARK, bold=True)
        add_textbox(slide, CONTENT_LEFT + Inches(0.35), y + Inches(0.55),
                   Inches(2.5), Inches(0.3),
                   text=translit, font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
                   font_color=INK_MEDIUM, italic=True)
        add_textbox(slide, CONTENT_LEFT + Inches(3.2), y + Inches(0.2),
                   Inches(3.5), Inches(0.4),
                   text=eng, font_name=FONT_LATIN, font_size=SIZE_BODY,
                   font_color=INK_DARK, bold=True)
        
        y += Inches(1.25)
    
    # ── Slide 7: Country Substitution Drill ─────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "3", "Country substitution drill")
    
    add_textbox(slide, CONTENT_LEFT, Inches(0.75), CONTENT_WIDTH, Inches(0.4),
               text="Keep the frame, swap only the country:",
               font_name=FONT_LATIN, font_size=SIZE_BODY,
               font_color=INK_MEDIUM)
    
    drills = [
        ("ฉันมาจากไทย", "chǎn maa-jàak thai", "I am from Thailand"),
        ("ผมมาจากอังกฤษ", "phǒm maa-jàak ang-grìt", "I am from England"),
        ("ฉันมาจากอเมริกา", "chǎn maa-jàak à-mee-rí-gaa", "I am from America"),
    ]
    
    y = Inches(1.3)
    colors = [ACCENT_TEAL, ACCENT_CLAY, ACCENT_GOLD]
    for i, (thai, translit, eng) in enumerate(drills):
        row_h = Inches(1.3)
        bg_color = HIGHLIGHT_BG if i % 2 == 0 else CARD_BG
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), row_h,
                               fill_color=bg_color, border_color=CARD_BORDER)
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.15), y + Inches(0.12),
                   Inches(0.4), Inches(0.4),
                   text=str(i + 1), font_name=FONT_LATIN, font_size=SIZE_BODY,
                   font_color=colors[i], bold=True)
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.5), y + Inches(0.05),
                   Inches(6.0), Inches(0.55),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
                   font_color=INK_DARK, bold=True)
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.5), y + Inches(0.55),
                   Inches(6.0), Inches(0.3),
                   text=translit, font_name=FONT_TRANSLIT, font_size=Pt(14),
                   font_color=INK_MEDIUM, italic=True)
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.5), y + Inches(0.85),
                   Inches(6.0), Inches(0.3),
                   text=eng, font_name=FONT_LATIN, font_size=SIZE_ENGLISH,
                   font_color=INK_LIGHT)
        
        y += row_h + Inches(0.12)
    
    # ── Slide 8: Full Self-Introduction Sequence ────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, "4", "Build a full self-introduction")
    
    add_textbox(slide, CONTENT_LEFT, Inches(0.75), CONTENT_WIDTH, Inches(0.4),
               text="Connect greeting + name + country into one smooth sequence:",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_MEDIUM)
    
    steps = [
        ("1", "Greet", "สวัสดีครับ/ค่ะ", "sà-wàt-dii khráp/khâ", "hello", ACCENT_GOLD),
        ("2", "Name", "ฉันชื่อ...", "chǎn chûue...", "my name is ...", ACCENT_TEAL),
        ("3", "Country", "ฉันมาจาก...", "chǎn maa-jàak...", "I am from ...", ACCENT_CLAY),
        ("4", "Close", "ขอบคุณครับ/ค่ะ", "khàawp-khun khráp/khâ", "thank you", ACCENT_GOLD),
    ]
    
    y = Inches(1.3)
    for num, label, thai, translit, eng, color in steps:
        row_h = Inches(1.15)
        card = add_rounded_rect(slide, CONTENT_LEFT, y, Inches(7.2), row_h,
                               fill_color=CARD_BG, border_color=CARD_BORDER)
        
        # Step number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            CONTENT_LEFT + Inches(0.15), y + Inches(0.25), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.name = FONT_LATIN
        run.font.size = Pt(18)
        run.font.color.rgb = WHITE
        run.font.bold = True
        
        # Step label
        add_textbox(slide, CONTENT_LEFT + Inches(0.8), y + Inches(0.05),
                   Inches(1.2), Inches(0.3),
                   text=label, font_name=FONT_LATIN, font_size=SIZE_LABEL,
                   font_color=color, bold=True)
        
        # Thai
        add_textbox(slide, CONTENT_LEFT + Inches(0.8), y + Inches(0.3),
                   Inches(5.5), Inches(0.45),
                   text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
                   font_color=INK_DARK, bold=True)
        
        # Translit + English
        add_textbox(slide, CONTENT_LEFT + Inches(0.8), y + Inches(0.72),
                   Inches(5.5), Inches(0.3),
                   text=f"{translit}  —  {eng}",
                   font_name=FONT_TRANSLIT, font_size=Pt(14),
                   font_color=INK_MEDIUM)
        
        y += row_h + Inches(0.1)
    
    # ── Slide 9: Roleplay ───────────────────────────────────────────────
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    add_section_header(slide, None, "Roleplay: Meeting someone new")
    
    add_textbox(slide, CONTENT_LEFT, Inches(0.7), CONTENT_WIDTH, Inches(0.4),
               text="At a class or event — practise the full introduction",
               font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
               font_color=INK_MEDIUM, italic=True)
    
    turns = [
        ("A", "สวัสดีค่ะ", "sà-wàt-dii khâ", "Hello.", False),
        ("B", "สวัสดีครับ", "sà-wàt-dii khráp", "Hello.", True),
        ("A", "คุณชื่ออะไรคะ", "khun chûue à-rai khá", "What is your name?", False),
        ("B", "ผมชื่อแดนครับ", "phǒm chûue Daen khráp", "My name is Dan.", True),
        ("A", "คุณมาจากไหนคะ", "khun maa-jàak nǎi khá", "Where are you from?", False),
        ("B", "ผมมาจากอังกฤษครับ", "phǒm maa-jàak ang-grìt khráp", "I am from England.", True),
    ]
    
    y = Inches(1.2)
    for speaker, thai, translit, english, is_learner in turns:
        h = add_dialogue_turn(slide, CONTENT_LEFT, y, Inches(7.2),
                             speaker, thai, translit, english,
                             is_learner=is_learner)
        y += h + Inches(0.05)
    
    # ── Slide 10: Recap ─────────────────────────────────────────────────
    add_recap_slide(prs, [
        "Learn คุณชื่ออะไร and its answer as one fixed beginner pair",
        "Choose one stable self-pronoun: ฉัน or ผม",
        "Use มาจาก to say where you are from — keep the country slot flexible",
        "Carry over สวัสดี and ขอบคุณ so the introduction sounds social",
    ])
    
    # ── Slide 11: Closing ───────────────────────────────────────────────
    add_closing_slide(
        prs,
        lesson_id="M01-L002",
        message="You can now introduce yourself\nwith name and country.",
        takeaway_lines=[
            "Hello, my name is..., I am from... — a complete first contact.",
            "Next lesson: pronouns and who questions.",
        ]
    )
    
    output_path = "/home/ubuntu/immersion-thai-slides/M01-L002-deck.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")
    return prs


if __name__ == "__main__":
    build_l002()
