"""
Immersion Thai with Nine — Slide Design System
Shared constants, helpers, and reusable slide builders for all M01 lessons.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Canvas ──────────────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)   # 1920 px at 144 dpi → 13.333 in
SLIDE_HEIGHT = Inches(7.5)     # 1080 px at 144 dpi → 7.5 in
LEFT_ZONE_W = Inches(8.889)    # 66.67 %
RIGHT_ZONE_X = Inches(8.889)
RIGHT_ZONE_W = Inches(4.444)   # 33.33 %

# Comfortable content margins inside the left zone
CONTENT_LEFT = Inches(0.8)
CONTENT_TOP = Inches(1.2)
CONTENT_WIDTH = Inches(7.5)
CONTENT_BOTTOM_MARGIN = Inches(0.6)

# ── Palette ─────────────────────────────────────────────────────────────
BG_IVORY        = RGBColor(0xFA, 0xF6, 0xF0)
BG_SAND_LIGHT   = RGBColor(0xF0, 0xEB, 0xE3)
BG_WARM_WHITE   = RGBColor(0xFE, 0xFC, 0xF9)
INK_DARK        = RGBColor(0x2C, 0x24, 0x20)
INK_MEDIUM      = RGBColor(0x5A, 0x4E, 0x46)
INK_LIGHT       = RGBColor(0x8A, 0x7E, 0x76)
ACCENT_GOLD     = RGBColor(0xB8, 0x96, 0x3E)
ACCENT_CLAY     = RGBColor(0xA0, 0x52, 0x2D)
ACCENT_TEAL     = RGBColor(0x3B, 0x7A, 0x72)
ACCENT_SOFT_GOLD = RGBColor(0xD4, 0xBE, 0x8A)
WHITE           = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG         = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER     = RGBColor(0xE0, 0xD8, 0xCE)
HIGHLIGHT_BG    = RGBColor(0xFD, 0xF5, 0xE6)  # warm highlight behind Thai
DIVIDER_COLOR   = RGBColor(0xD8, 0xCE, 0xC0)

# ── Fonts ───────────────────────────────────────────────────────────────
FONT_THAI       = "Noto Sans Thai"
FONT_LATIN      = "Inter"
FONT_TRANSLIT   = "Inter"
FONT_MONO       = "Consolas"

# ── Font sizes ──────────────────────────────────────────────────────────
SIZE_TITLE       = Pt(36)
SIZE_SUBTITLE    = Pt(20)
SIZE_THAI_LARGE  = Pt(48)
SIZE_THAI_MED    = Pt(36)
SIZE_THAI_SMALL  = Pt(28)
SIZE_TRANSLIT    = Pt(18)
SIZE_ENGLISH     = Pt(16)
SIZE_BODY        = Pt(18)
SIZE_BODY_SMALL  = Pt(15)
SIZE_LABEL       = Pt(13)
SIZE_SECTION_NUM = Pt(14)

# ── Helpers ─────────────────────────────────────────────────────────────

def create_presentation():
    """Create a blank 16:9 presentation with correct dimensions."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_blank_slide(prs):
    """Add a blank slide layout."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    return slide


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_right_zone_tint(slide, color=BG_SAND_LIGHT, opacity=None):
    """Add a subtle tinted rectangle for the right camera zone."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        RIGHT_ZONE_X, Inches(0),
        RIGHT_ZONE_W, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Send to back
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_textbox(slide, left, top, width, height, text="",
                font_name=FONT_LATIN, font_size=SIZE_BODY,
                font_color=INK_DARK, bold=False, italic=False,
                alignment=PP_ALIGN.LEFT, word_wrap=True,
                anchor=MSO_ANCHOR.TOP):
    """Add a text box with a single run of styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except:
        pass
    
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """
    Add a text box with multiple styled lines.
    lines: list of dicts with keys: text, font_name, font_size, font_color, bold, italic, space_after
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    
    for i, line_spec in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.alignment = alignment
        if 'space_after' in line_spec:
            p.space_after = line_spec['space_after']
        if 'space_before' in line_spec:
            p.space_before = line_spec['space_before']
        
        run = p.add_run()
        run.text = line_spec.get('text', '')
        run.font.name = line_spec.get('font_name', FONT_LATIN)
        run.font.size = line_spec.get('font_size', SIZE_BODY)
        run.font.color.rgb = line_spec.get('font_color', INK_DARK)
        run.font.bold = line_spec.get('bold', False)
        run.font.italic = line_spec.get('italic', False)
    
    return txBox


def add_rounded_rect(slide, left, top, width, height,
                     fill_color=CARD_BG, border_color=CARD_BORDER,
                     border_width=Pt(1)):
    """Add a rounded rectangle card background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    # Adjust corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_horizontal_line(slide, left, top, width, color=DIVIDER_COLOR, weight=Pt(1)):
    """Add a thin horizontal divider line."""
    connector = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, weight
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = color
    connector.line.fill.background()
    return connector


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6),
                   color=ACCENT_GOLD):
    """Add a small vertical accent bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── Phrase Card Builder ─────────────────────────────────────────────────

def add_phrase_card(slide, left, top, width, thai, translit, english,
                    thai_size=SIZE_THAI_LARGE, card_bg=CARD_BG,
                    border_color=CARD_BORDER, accent_color=ACCENT_GOLD,
                    show_accent_bar=True):
    """
    Add a phrase card with Thai (large), transliteration (medium), English (small).
    Returns the card shape and text boxes for later editing.
    """
    card_height = Inches(1.8)
    
    # Card background
    card = add_rounded_rect(slide, left, top, width, card_height,
                           fill_color=card_bg, border_color=border_color)
    
    # Accent bar on left edge
    if show_accent_bar:
        add_accent_bar(slide, left + Inches(0.15), top + Inches(0.3),
                      height=Inches(1.2), color=accent_color)
    
    text_left = left + Inches(0.45) if show_accent_bar else left + Inches(0.25)
    text_width = width - Inches(0.7)
    
    # Thai text
    add_textbox(slide, text_left, top + Inches(0.2), text_width, Inches(0.7),
               text=thai, font_name=FONT_THAI, font_size=thai_size,
               font_color=INK_DARK, bold=True)
    
    # Transliteration
    add_textbox(slide, text_left, top + Inches(0.85), text_width, Inches(0.4),
               text=translit, font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM, italic=True)
    
    # English
    add_textbox(slide, text_left, top + Inches(1.25), text_width, Inches(0.4),
               text=english, font_name=FONT_LATIN, font_size=SIZE_ENGLISH,
               font_color=INK_LIGHT)
    
    return card


def add_compact_phrase(slide, left, top, width, thai, translit, english,
                       thai_size=SIZE_THAI_MED, accent_color=ACCENT_GOLD):
    """Add a compact phrase row (no card background) — Thai | translit | English."""
    # Accent dot
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left, top + Inches(0.18), Inches(0.12), Inches(0.12)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = accent_color
    dot.line.fill.background()
    
    text_left = left + Inches(0.25)
    text_width = width - Inches(0.3)
    
    add_textbox(slide, text_left, top, text_width, Inches(0.5),
               text=thai, font_name=FONT_THAI, font_size=thai_size,
               font_color=INK_DARK, bold=True)
    
    add_textbox(slide, text_left, top + Inches(0.5), text_width, Inches(0.3),
               text=f"{translit}  —  {english}",
               font_name=FONT_TRANSLIT, font_size=SIZE_TRANSLIT,
               font_color=INK_MEDIUM, italic=False)


def add_mini_phrase_row(slide, left, top, width, thai, translit, english,
                        thai_size=SIZE_THAI_SMALL):
    """Add a single-line compact phrase: Thai  translit  english."""
    # Thai
    add_textbox(slide, left, top, Inches(2.5), Inches(0.45),
               text=thai, font_name=FONT_THAI, font_size=thai_size,
               font_color=INK_DARK, bold=True)
    # translit + english
    add_textbox(slide, left + Inches(2.6), top + Inches(0.05), width - Inches(2.6), Inches(0.4),
               text=f"{translit}  ·  {english}",
               font_name=FONT_TRANSLIT, font_size=Pt(15),
               font_color=INK_MEDIUM)


# ── Section Header ──────────────────────────────────────────────────────

def add_section_header(slide, section_num, title, subtitle=""):
    """Add a section header bar at the top of a teaching slide."""
    # Top accent line
    add_horizontal_line(slide, Inches(0), Inches(0), LEFT_ZONE_W,
                       color=ACCENT_GOLD, weight=Pt(4))
    
    # Section number pill
    if section_num:
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            CONTENT_LEFT, Inches(0.25), Inches(1.0), Inches(0.32)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = ACCENT_GOLD
        pill.line.fill.background()
        pill.adjustments[0] = 0.3
        
        tf = pill.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"Section {section_num}"
        run.font.name = FONT_LATIN
        run.font.size = SIZE_LABEL
        run.font.color.rgb = WHITE
        run.font.bold = True
    
    # Title — use wider box and single line
    title_top = Inches(0.22) if not section_num else Inches(0.22)
    title_left = CONTENT_LEFT + Inches(1.15) if section_num else CONTENT_LEFT
    title_width = Inches(6.5) if section_num else Inches(7.5)
    add_textbox(slide, title_left, title_top, title_width, Inches(0.5),
               text=title, font_name=FONT_LATIN, font_size=Pt(30),
               font_color=INK_DARK, bold=True)
    
    if subtitle:
        add_textbox(slide, CONTENT_LEFT, Inches(0.72), CONTENT_WIDTH, Inches(0.4),
                   text=subtitle, font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL,
                   font_color=INK_MEDIUM)


# ── Dialogue Builder ────────────────────────────────────────────────────

def add_dialogue_turn(slide, left, top, width, speaker, thai, translit, english,
                      speaker_color=ACCENT_TEAL, is_learner=False):
    """Add one dialogue turn with speaker label, Thai, translit, English."""
    # Speaker label
    label_bg = ACCENT_TEAL if not is_learner else ACCENT_CLAY
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, Inches(1.2), Inches(0.32)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = label_bg
    pill.line.fill.background()
    pill.adjustments[0] = 0.3
    
    tf = pill.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = speaker
    run.font.name = FONT_LATIN
    run.font.size = SIZE_LABEL
    run.font.color.rgb = WHITE
    run.font.bold = True
    
    # Thai line
    add_textbox(slide, left + Inches(1.4), top - Inches(0.05), width - Inches(1.4), Inches(0.45),
               text=thai, font_name=FONT_THAI, font_size=SIZE_THAI_SMALL,
               font_color=INK_DARK, bold=True)
    
    # Translit + English
    add_textbox(slide, left + Inches(1.4), top + Inches(0.38), width - Inches(1.4), Inches(0.35),
               text=f"{translit}  —  {english}",
               font_name=FONT_TRANSLIT, font_size=Pt(14),
               font_color=INK_MEDIUM)
    
    return Inches(0.9)  # height consumed


# ── Lesson Opener ───────────────────────────────────────────────────────

def add_lesson_opener(prs, lesson_id, lesson_title, module_title, level="A0"):
    """Create the lesson title/opener slide."""
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    
    # Large gold accent line at top
    add_horizontal_line(slide, Inches(0), Inches(0), LEFT_ZONE_W,
                       color=ACCENT_GOLD, weight=Pt(6))
    
    # Module + level label
    add_textbox(slide, CONTENT_LEFT, Inches(1.8), CONTENT_WIDTH, Inches(0.4),
               text=f"{module_title}  ·  Level {level}",
               font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=ACCENT_GOLD, bold=True)
    
    # Lesson ID
    add_textbox(slide, CONTENT_LEFT, Inches(2.3), CONTENT_WIDTH, Inches(0.5),
               text=lesson_id,
               font_name=FONT_LATIN, font_size=SIZE_SUBTITLE,
               font_color=INK_MEDIUM, bold=False)
    
    # Lesson title
    add_textbox(slide, CONTENT_LEFT, Inches(2.9), CONTENT_WIDTH, Inches(1.0),
               text=lesson_title,
               font_name=FONT_LATIN, font_size=Pt(44),
               font_color=INK_DARK, bold=True)
    
    # Course name at bottom
    add_textbox(slide, CONTENT_LEFT, Inches(5.8), CONTENT_WIDTH, Inches(0.4),
               text="Immersion Thai with Nine",
               font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=INK_LIGHT, bold=False)
    
    # Bottom accent line
    add_horizontal_line(slide, CONTENT_LEFT, Inches(6.3), Inches(2.0),
                       color=ACCENT_GOLD, weight=Pt(3))
    
    return slide


# ── Objectives Slide ────────────────────────────────────────────────────

def add_objectives_slide(prs, objectives, title="What you will learn"):
    """Create the objectives/what-you-will-learn slide."""
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    
    add_section_header(slide, None, title)
    
    y = Inches(1.3)
    for i, obj in enumerate(objectives):
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            CONTENT_LEFT, y + Inches(0.05), Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_GOLD
        circle.line.fill.background()
        
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.name = FONT_LATIN
        run.font.size = Pt(16)
        run.font.color.rgb = WHITE
        run.font.bold = True
        
        # Objective text
        add_textbox(slide, CONTENT_LEFT + Inches(0.6), y, Inches(6.8), Inches(0.5),
                   text=obj, font_name=FONT_LATIN, font_size=SIZE_BODY,
                   font_color=INK_DARK)
        
        y += Inches(0.7)
    
    return slide


# ── Recap Slide ─────────────────────────────────────────────────────────

def add_recap_slide(prs, points, title="Recap"):
    """Create a recap slide with bullet points."""
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    
    add_section_header(slide, None, title)
    
    y = Inches(1.3)
    for point in points:
        # Checkmark or bullet
        add_textbox(slide, CONTENT_LEFT, y, Inches(0.3), Inches(0.4),
                   text="✓", font_name=FONT_LATIN, font_size=SIZE_BODY,
                   font_color=ACCENT_TEAL, bold=True)
        
        add_textbox(slide, CONTENT_LEFT + Inches(0.4), y, Inches(7.0), Inches(0.5),
                   text=point, font_name=FONT_LATIN, font_size=SIZE_BODY,
                   font_color=INK_DARK)
        
        y += Inches(0.65)
    
    return slide


# ── Closing Slide ───────────────────────────────────────────────────────

def add_closing_slide(prs, lesson_id, message, takeaway_lines=None):
    """Create the closing summary slide."""
    slide = add_blank_slide(prs)
    set_slide_bg(slide, BG_IVORY)
    add_right_zone_tint(slide, BG_SAND_LIGHT)
    
    add_horizontal_line(slide, Inches(0), Inches(0), LEFT_ZONE_W,
                       color=ACCENT_GOLD, weight=Pt(6))
    
    add_textbox(slide, CONTENT_LEFT, Inches(1.5), CONTENT_WIDTH, Inches(0.4),
               text=lesson_id,
               font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=ACCENT_GOLD, bold=True)
    
    add_textbox(slide, CONTENT_LEFT, Inches(2.1), CONTENT_WIDTH, Inches(1.0),
               text=message,
               font_name=FONT_LATIN, font_size=Pt(32),
               font_color=INK_DARK, bold=True)
    
    if takeaway_lines:
        y = Inches(3.5)
        for line in takeaway_lines:
            add_textbox(slide, CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.45),
                       text=line, font_name=FONT_LATIN, font_size=SIZE_BODY,
                       font_color=INK_MEDIUM)
            y += Inches(0.5)
    
    add_textbox(slide, CONTENT_LEFT, Inches(5.8), CONTENT_WIDTH, Inches(0.4),
               text="Immersion Thai with Nine",
               font_name=FONT_LATIN, font_size=SIZE_LABEL,
               font_color=INK_LIGHT)
    
    add_horizontal_line(slide, CONTENT_LEFT, Inches(6.3), Inches(2.0),
                       color=ACCENT_GOLD, weight=Pt(3))
    
    return slide
