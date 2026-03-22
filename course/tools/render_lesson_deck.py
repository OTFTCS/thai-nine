#!/usr/bin/env python3
"""Render a lesson deck PPTX and deck-source.json from script-master.json."""

from __future__ import annotations

import argparse
import csv
import datetime
import json
import math
import re
import sys
import tempfile
import urllib.error
import urllib.parse
import urllib.request
import xml.etree.ElementTree as ET
from io import BytesIO
from pathlib import Path
from typing import Any
from zipfile import ZIP_DEFLATED, ZipFile

from PIL import Image, ImageDraw, ImageOps, UnidentifiedImageError
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
# PiP camera placeholder (top-right, semi-transparent overlay)
PIP_W = Inches(4.2)
PIP_H = Inches(3.15)           # 4:3 aspect for camera
PIP_X = SLIDE_WIDTH - PIP_W    # Right-aligned
PIP_Y = Inches(0)              # Top-right
# Content width: constrained beside PiP (Y < PIP_H), full width below
LEFT_ZONE_W = SLIDE_WIDTH  # Full width (legacy alias, kept for opener divider)
CONTENT_LEFT = Inches(0.8)
CONTENT_TOP = Inches(1.15)
CONTENT_WIDTH_BESIDE_PIP = PIP_X - CONTENT_LEFT - Inches(0.6)  # ~7.7" (beside PiP, with margin)
CONTENT_WIDTH = Inches(11.0)                                     # Full width (below PiP)
PIP_BOTTOM = PIP_Y + PIP_H  # Y threshold: below this, full width is safe

BG_IVORY = RGBColor(0xF5, 0xF7, 0xF8)
BG_SAND_LIGHT = RGBColor(0xDB, 0xE7, 0xEC)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xD7, 0xDE, 0xE3)
INK_DARK = RGBColor(0x24, 0x33, 0x3D)
INK_MEDIUM = RGBColor(0x56, 0x68, 0x73)
INK_LIGHT = RGBColor(0x82, 0x90, 0x9A)
ACCENT_GOLD = RGBColor(0x35, 0x5C, 0x7D)
ACCENT_TEAL = RGBColor(0x7A, 0x8F, 0x54)
ACCENT_CLAY = RGBColor(0xC9, 0x6F, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER_COLOR = RGBColor(0xC9, 0xD7, 0xDE)

FONT_THAI = "Sarabun"
FONT_LATIN = "Sarabun"
FONT_TRANSLIT = "Sarabun"

SIZE_TITLE = Pt(36)
SIZE_THAI_LARGE = Pt(36)
SIZE_THAI_MED = Pt(30)
SIZE_THAI_SMALL = Pt(26)
SIZE_TRANSLIT = Pt(17)
SIZE_ENGLISH = Pt(17)
SIZE_BODY = Pt(22)
SIZE_BODY_SMALL = Pt(17)
SIZE_LABEL = Pt(12)

OPENVERSE_ENDPOINT = "https://api.openverse.org/v1/images/"
WIKIMEDIA_ENDPOINT = "https://commons.wikimedia.org/w/api.php"
EMU_PER_INCH = 914400
CANVA_BG_WIDTH = 1920
CANVA_BG_HEIGHT = 1080


def lesson_parts(lesson_id: str) -> tuple[str, str]:
    match = re.fullmatch(r"(M\d{2})-(L\d{3})", lesson_id)
    if not match:
        raise ValueError(f"Invalid lesson id: {lesson_id}")
    return match.group(1), match.group(2)


def lesson_dir(repo_root: Path, lesson_id: str) -> Path:
    module_id, lesson_key = lesson_parts(lesson_id)
    return repo_root / "course" / "modules" / module_id / lesson_key


def artifact_name(lesson_id: str, base_name: str) -> str:
    return f"{lesson_id}-{base_name}"


def read_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def read_blueprint_row(repo_root: Path, lesson_id: str) -> dict[str, str]:
    csv_path = repo_root / "course" / "exports" / "full-thai-course-blueprint.csv"
    with csv_path.open("r", encoding="utf-8", newline="") as handle:
        for row in csv.DictReader(handle):
            if row.get("lesson_id", "").strip() == lesson_id:
                return {key: (value or "").strip() for key, value in row.items()}
    raise ValueError(f"Lesson {lesson_id} not found in full-thai-course-blueprint.csv")


def to_hex(color: RGBColor) -> str:
    return f"#{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def slugify(value: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "-", value.strip().lower()).strip("-")
    return slug or "asset"


def _you(text: str) -> str:
    """Replace 'Learner' with 'You' in user-facing text (2nd person for 1:1 feel)."""
    return (text
            .replace("Learner can ", "You can ")
            .replace("learner can ", "you can ")
            .replace("The learner ", "You ")
            .replace("the learner ", "you ")
            .replace("A learner ", "You ")
            .replace("a learner ", "you ")
            .replace("Learner ", "You ")
            .replace("learner ", "you "))


def contains_thai(text: str) -> bool:
    return any("\u0E00" <= char <= "\u0E7F" for char in text)


def inline_phrase(thai: str, translit: str) -> str:
    thai = thai.strip()
    translit = translit.strip()
    if thai and translit:
        return f"{thai} ({translit})"
    return thai or translit


def inline_variants(items: list[tuple[str, str]], separator: str = " / ") -> str:
    parts = [inline_phrase(thai, translit) for thai, translit in items if thai.strip() or translit.strip()]
    return separator.join(part for part in parts if part)


def translit_font_size_value(size_pt: float) -> float:
    return max(float(SIZE_TRANSLIT.pt), round(size_pt * 0.62, 1))


def build_translit_entries(script: dict[str, Any]) -> list[tuple[str, str]]:
    entries: dict[str, str] = {}

    def add_entry(thai: str, translit: str) -> None:
        thai = str(thai or "").strip()
        translit = str(translit or "").strip()
        if thai and translit:
            entries[thai] = translit

    for section in script.get("sections", []):
        for lex in section.get("languageFocus", []):
            add_entry(lex.get("thai", ""), lex.get("translit", ""))
        for thai, translit, _english in parse_triplet_lines(section.get("onScreenBullets", [])):
            add_entry(thai, translit)

    for line in script.get("roleplay", {}).get("lines", []):
        add_entry(line.get("thai", ""), line.get("translit", ""))

    return sorted(entries.items(), key=lambda item: len(item[0]), reverse=True)


def inline_tokens(text: str, translit_entries: list[tuple[str, str]]) -> list[dict[str, str]]:
    if not text or not contains_thai(text) or not translit_entries:
        return [{"kind": "plain", "text": text}]

    tokens: list[dict[str, str]] = []
    cursor = 0
    length = len(text)

    while cursor < length:
        best_match: tuple[int, str, str] | None = None
        for thai, translit in translit_entries:
            search_at = text.find(thai, cursor)
            while search_at != -1 and re.match(r"\s*\(", text[search_at + len(thai) :]):
                search_at = text.find(thai, search_at + len(thai))
            if search_at == -1:
                continue
            if best_match is None or search_at < best_match[0] or (
                search_at == best_match[0] and len(thai) > len(best_match[1])
            ):
                best_match = (search_at, thai, translit)

        if best_match is None:
            tokens.append({"kind": "plain", "text": text[cursor:]})
            break

        match_at, thai, translit = best_match
        if match_at > cursor:
            tokens.append({"kind": "plain", "text": text[cursor:match_at]})
        tokens.append({"kind": "thai", "thai": thai, "translit": translit})
        cursor = match_at + len(thai)

    return [token for token in tokens if token.get("text") or token.get("thai")]


def tokens_to_text(tokens: list[dict[str, str]]) -> str:
    parts: list[str] = []
    for token in tokens:
        if token["kind"] == "plain":
            parts.append(token["text"])
        else:
            parts.append(inline_phrase(token["thai"], token["translit"]))
    return "".join(parts)


def render_runs_from_tokens(
    paragraph,
    tokens: list[dict[str, str]],
    *,
    default_font_name: str,
    default_font_size,
    default_color,
    default_bold: bool = False,
    default_italic: bool = False,
) -> None:
    for token in tokens:
        if token["kind"] == "plain":
            if not token["text"]:
                continue
            run = paragraph.add_run()
            run.text = token["text"]
            run.font.name = default_font_name
            run.font.size = default_font_size
            run.font.color.rgb = default_color
            run.font.bold = default_bold
            run.font.italic = default_italic
            apply_run_font_metadata(run, font_name=default_font_name)
            continue

        thai_run = paragraph.add_run()
        thai_run.text = token["thai"]
        thai_run.font.name = FONT_THAI
        thai_run.font.size = default_font_size
        thai_run.font.color.rgb = default_color
        thai_run.font.bold = default_bold
        thai_run.font.italic = False
        apply_run_font_metadata(thai_run, font_name=FONT_THAI, lang="th-TH")

        if token["translit"]:
            translit_run = paragraph.add_run()
            translit_run.text = f" ({token['translit']})"
            translit_run.font.name = FONT_TRANSLIT
            translit_run.font.size = Pt(translit_font_size_value(float(default_font_size.pt)))
            translit_run.font.color.rgb = INK_MEDIUM
            translit_run.font.bold = False
            translit_run.font.italic = True
            apply_run_font_metadata(translit_run, font_name=FONT_TRANSLIT, lang="en-US")


def estimate_seconds(lines: list[str], minimum: int = 12) -> int:
    word_count = len(" ".join(lines).split())
    return max(minimum, min(60, math.ceil(word_count / 2.8)))


def choose_image_usage(visual_plan: dict[str, Any] | None) -> str:
    if not visual_plan:
        return "icon"
    image_support = visual_plan.get("imageSupport") or {}
    if image_support.get("helpful") is True:
        return "real-image"
    if image_support.get("helpful") is False and image_support.get("priority") == "avoid":
        return "text-only"
    return "icon"


def request_json(url: str) -> dict[str, Any] | None:
    req = urllib.request.Request(url, headers={"User-Agent": "thai-nine-deck-renderer/1.0"})
    try:
        with urllib.request.urlopen(req, timeout=12) as response:
            return json.loads(response.read().decode("utf-8"))
    except Exception:
        return None


def fetch_openverse_candidate(query: str) -> dict[str, str] | None:
    params = urllib.parse.urlencode({"q": query, "page_size": 10})
    payload = request_json(f"{OPENVERSE_ENDPOINT}?{params}")
    if not payload:
        return None

    for item in payload.get("results", []):
        image_url = str(item.get("url") or "").strip()
        if not image_url:
            continue
        source_url = str(item.get("foreign_landing_url") or image_url).strip()
        license_name = str(item.get("license") or "").strip()
        license_version = str(item.get("license_version") or "").strip()
        license_label = " ".join(part for part in [license_name, license_version] if part).strip()
        if not license_label:
            license_label = "Openverse listing"
        return {
            "provider": "openverse",
            "download_url": image_url,
            "source_url": source_url,
            "license": license_label,
        }
    return None


def fetch_wikimedia_candidate(query: str) -> dict[str, str] | None:
    params = urllib.parse.urlencode(
        {
            "action": "query",
            "format": "json",
            "generator": "search",
            "gsrsearch": query,
            "gsrnamespace": 6,
            "gsrlimit": 8,
            "prop": "imageinfo",
            "iiprop": "url|extmetadata",
        }
    )
    payload = request_json(f"{WIKIMEDIA_ENDPOINT}?{params}")
    if not payload:
        return None

    pages = payload.get("query", {}).get("pages", {})
    for page in pages.values():
        info = (page.get("imageinfo") or [{}])[0]
        download_url = str(info.get("url") or "").strip()
        if not download_url:
            continue
        ext = info.get("extmetadata") or {}
        license_label = (
            (ext.get("LicenseShortName") or {}).get("value")
            or (ext.get("UsageTerms") or {}).get("value")
            or "Wikimedia Commons"
        )
        source_url = str(info.get("descriptionurl") or download_url).strip()
        return {
            "provider": "wikimedia",
            "download_url": download_url,
            "source_url": source_url,
            "license": re.sub(r"<[^>]+>", "", str(license_label)).strip() or "Wikimedia Commons",
        }
    return None


def download_image(download_url: str, destination: Path) -> None:
    req = urllib.request.Request(download_url, headers={"User-Agent": "thai-nine-deck-renderer/1.0"})
    with urllib.request.urlopen(req, timeout=20) as response:
        raw = response.read()

    image = Image.open(BytesIO(raw))
    image = ImageOps.exif_transpose(image)
    if image.mode not in ("RGB", "RGBA"):
        image = image.convert("RGBA")
    image.thumbnail((1600, 900))
    destination.parent.mkdir(parents=True, exist_ok=True)
    image.save(destination, format="PNG")


def resolve_asset(
    slide_id: str,
    lesson_root: Path,
    assets_dir: Path,
    query: str,
    usage_notes: str,
) -> tuple[dict[str, Any], dict[str, Any]]:
    asset_id = f"{slide_id}-asset-1"
    providers = [fetch_openverse_candidate, fetch_wikimedia_candidate]

    for provider_fn in providers:
        candidate = provider_fn(query)
        if not candidate:
            continue

        destination = assets_dir / f"{slide_id}-{slugify(query)[:48]}.png"
        try:
            download_image(candidate["download_url"], destination)
        except (urllib.error.URLError, TimeoutError, UnidentifiedImageError, OSError):
            continue

        local_path = destination.relative_to(lesson_root).as_posix()
        asset = {
            "assetId": asset_id,
            "kind": "image",
            "query": query,
            "sourcePolicy": "internet-first",
            "status": "resolved",
            "sourceProvider": candidate["provider"],
            "sourceUrl": candidate["source_url"],
            "license": candidate["license"],
            "localPath": local_path,
            "usageNotes": usage_notes,
        }
        provenance = {
            "assetId": asset_id,
            "slideId": slide_id,
            "kind": "image",
            "status": "resolved",
            "sourceProvider": candidate["provider"],
            "sourceUrl": candidate["source_url"],
            "license": candidate["license"],
            "usage": f"{slide_id}:left-panel image",
            "query": query,
            "sourcePolicy": "internet-first",
            "localPath": local_path,
            "rationale": usage_notes,
        }
        return asset, provenance

    fallback_reason = "No acceptable Openverse or Wikimedia asset could be downloaded automatically."
    asset = {
        "assetId": asset_id,
        "kind": "image",
        "query": query,
        "sourcePolicy": "internet-first",
        "status": "fallback-text-only",
        "sourceProvider": "none",
        "fallbackReason": fallback_reason,
        "usageNotes": usage_notes,
    }
    provenance = {
        "assetId": asset_id,
        "slideId": slide_id,
        "kind": "image",
        "status": "fallback-text-only",
        "sourceProvider": "none",
        "usage": f"{slide_id}:left-panel image fallback",
        "query": query,
        "sourcePolicy": "internet-first",
        "fallbackReason": fallback_reason,
        "rationale": usage_notes,
    }
    return asset, provenance


def build_slide_id(index: int, title: str) -> str:
    return f"slide-{index:02d}-{slugify(title)[:24]}"


def _add_pip_placeholder(slide):
    """Add semi-transparent PiP camera placeholder in top-right corner."""
    pip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PIP_X, PIP_Y, PIP_W, PIP_H)
    pip.fill.solid()
    pip.fill.fore_color.rgb = BG_SAND_LIGHT
    # Set 30% opacity via XML
    try:
        sp_pr = pip._element.spPr
        solid_fill_el = sp_pr.find(qn("a:solidFill"))
        if solid_fill_el is not None:
            srgb = solid_fill_el.find(qn("a:srgbClr"))
            if srgb is not None:
                alpha = OxmlElement("a:alpha")
                alpha.set("val", "30000")  # 30% opacity
                srgb.append(alpha)
    except Exception:
        pass  # Opacity is cosmetic — don't fail on it
    # Dashed navy border
    pip.line.color.rgb = RGBColor(0x35, 0x5C, 0x7D)
    pip.line.width = Pt(1.5)
    pip.line.dash_style = 2  # MSO_LINE_DASH_STYLE.DASH = 2
    # Add "Nine" label
    tf = pip.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Nine"
    run.font.name = FONT_LATIN
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0x35, 0x5C, 0x7D)
    run.font.bold = True
    from pptx.enum.text import MSO_ANCHOR as _A
    tf.vertical_anchor = _A.MIDDLE
    pip.name = "__pip_placeholder__"
    return pip


def add_blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_IVORY
    # PiP placeholder added last so it renders on top of content
    _add_pip_placeholder(slide)
    return slide


def _ensure_rpr_typeface(rpr, tag: str, typeface: str) -> None:
    child = rpr.find(qn(f"a:{tag}"))
    if child is None:
        child = OxmlElement(f"a:{tag}")
        rpr.append(child)
    child.set("typeface", typeface)


def apply_run_font_metadata(run, *, font_name: str | None = None, lang: str | None = None) -> None:
    text = run.text or ""
    resolved_font = font_name or run.font.name or (FONT_THAI if contains_thai(text) else FONT_LATIN)
    resolved_lang = lang or ("th-TH" if contains_thai(text) else "en-US")
    run.font.name = resolved_font

    rpr = run._r.get_or_add_rPr()
    rpr.set("lang", resolved_lang)
    rpr.set("altLang", "en-US")
    _ensure_rpr_typeface(rpr, "latin", resolved_font)
    _ensure_rpr_typeface(rpr, "ea", resolved_font)
    _ensure_rpr_typeface(rpr, "cs", resolved_font)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_name=FONT_LATIN,
    font_size=SIZE_BODY,
    color=INK_DARK,
    bold=False,
    italic=False,
    alignment=PP_ALIGN.LEFT,
    translit_entries: list[tuple[str, str]] | None = None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = alignment
    tokens = inline_tokens(text, translit_entries or [])
    render_runs_from_tokens(
        paragraph,
        tokens,
        default_font_name=font_name,
        default_font_size=font_size,
        default_color=color,
        default_bold=bold,
        default_italic=italic,
    )
    return box


def add_paragraphs(
    slide,
    left,
    top,
    width,
    height,
    lines: list[dict[str, Any]],
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.clear()

    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = line.get("alignment", PP_ALIGN.LEFT)
        run = paragraph.add_run()
        run.text = line.get("text", "")
        run.font.name = line.get("font_name", FONT_LATIN)
        run.font.size = line.get("font_size", SIZE_BODY)
        run.font.color.rgb = line.get("color", INK_DARK)
        run.font.bold = line.get("bold", False)
        run.font.italic = line.get("italic", False)
        apply_run_font_metadata(run)
    return box


def add_card(slide, left, top, width, height, fill=None):
    """Invisible layout spacer — no visible box drawn."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.background()  # Transparent
    card.line.fill.background()  # No border
    return card


def add_divider(slide, left, top, width, color=DIVIDER_COLOR, height=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_phrase_card(slide, left, top, width, thai, translit, english, compact=False):
    card_h = Inches(0.75) if compact else Inches(1.72)
    accent_h = Inches(0.45) if compact else Inches(1.25)
    add_card(slide, left, top, width, card_h)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.15), top + Inches(0.12 if compact else 0.22), Inches(0.08), accent_h)
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_GOLD
    accent.line.fill.background()
    text_left = left + Inches(0.38)
    if compact:
        # Compact: Thai + translit on one line, English on next, smaller fonts
        add_textbox(
            slide, text_left, top + Inches(0.08), width - Inches(0.5), Inches(0.35),
            thai, font_name=FONT_THAI, font_size=SIZE_THAI_SMALL, bold=True,
            translit_entries=[(thai, translit)],
        )
        add_textbox(
            slide, text_left, top + Inches(0.42), width - Inches(0.5), Inches(0.25),
            english, font_name=FONT_LATIN, font_size=Pt(13), color=INK_LIGHT,
        )
    else:
        add_textbox(
            slide, text_left, top + Inches(0.12), width - Inches(0.5), Inches(0.72),
            thai, font_name=FONT_THAI, font_size=SIZE_THAI_LARGE, bold=True,
            translit_entries=[(thai, translit)],
        )
        add_textbox(
            slide, text_left, top + Inches(0.96), width - Inches(0.5), Inches(0.3),
            english, font_name=FONT_LATIN, font_size=SIZE_ENGLISH, color=INK_LIGHT,
        )


def add_section_header(slide, title, eyebrow, translit_entries: list[tuple[str, str]] | None = None):
    add_divider(slide, 0, 0, LEFT_ZONE_W, ACCENT_GOLD, Pt(5))
    # Headers are beside the PiP — use constrained width
    header_width = CONTENT_WIDTH_BESIDE_PIP
    if eyebrow:
        add_textbox(
            slide,
            CONTENT_LEFT,
            Inches(0.28),
            header_width,
            Inches(0.28),
            eyebrow,
            font_name=FONT_LATIN,
            font_size=SIZE_LABEL,
            color=ACCENT_GOLD,
            bold=True,
            translit_entries=translit_entries,
        )
    # Use smaller font for long titles to prevent wrapping overlap
    title_font = Pt(28) if len(title) > 35 else SIZE_TITLE
    add_textbox(
        slide,
        CONTENT_LEFT,
        Inches(0.58),
        header_width,
        Inches(0.7),
        title,
        font_name=FONT_LATIN,
        font_size=title_font,
        bold=True,
        translit_entries=translit_entries,
    )


def parse_triplet_lines(lines: list[str]) -> list[tuple[str, str, str]]:
    triplets: list[tuple[str, str, str]] = []
    for line in lines:
        parts = [part.strip() for part in line.split("|")]
        if len(parts) >= 3:
            triplets.append((parts[0], parts[1], parts[2]))
    return triplets


def add_triplet_rows(slide, left, top, width, lines: list[str]) -> float:
    parsed = parse_triplet_lines(lines[:6])
    compact = len(parsed) > 4
    spacing = Inches(1.92) if len(parsed) <= 3 else Inches(1.4) if len(parsed) <= 4 else Inches(0.88)
    current_top = top
    for thai, translit, english in parsed:
        add_phrase_card(slide, left, current_top, width, thai, translit, english, compact=compact)
        current_top += spacing
    return current_top


def add_bullet_block(
    slide,
    left,
    top,
    width,
    heading,
    lines,
    accent_color=ACCENT_TEAL,
    translit_entries: list[tuple[str, str]] | None = None,
):
    # Cap lines to fit within slide bottom (estimate conservatively)
    available_height = SLIDE_HEIGHT - top - Inches(0.3)
    max_lines = max(1, int((available_height - Inches(0.48)) / Inches(0.45)))
    display_lines = lines[:min(5, max_lines)]

    add_textbox(
        slide,
        left + Inches(0.22),
        top + Inches(0.16),
        width - Inches(0.4),
        Inches(0.25),
        heading,
        font_name=FONT_LATIN,
        font_size=SIZE_BODY_SMALL,
        color=accent_color,
        bold=True,
        translit_entries=translit_entries,
    )
    current_top = top + Inches(0.48)
    text_width = width - Inches(0.55)
    # Estimate ~7 chars per inch at body-small size for mixed Thai/English
    chars_per_line = max(1, int(text_width / 914400 * 7))
    for line in display_lines:
        # Estimate wrapped line count and set textbox height accordingly
        est_lines = max(1, -(-len(line) // chars_per_line))  # ceiling division
        line_height = Inches(0.22) * est_lines + Inches(0.08)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.22), current_top + Inches(0.1), Inches(0.09), Inches(0.09))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent_color
        dot.line.fill.background()
        add_textbox(
            slide,
            left + Inches(0.38),
            current_top,
            text_width,
            line_height,
            line,
            font_name=FONT_LATIN,
            font_size=SIZE_BODY_SMALL,
            color=INK_DARK,
            translit_entries=translit_entries,
        )
        current_top += line_height + Inches(0.06)
    return current_top


def add_dialogue_turn(slide, left, top, width, speaker, thai, translit, english, learner=False):
    # Truncate long speaker labels for the pill — keep only the role name
    pill_label = speaker.split("(")[0].strip() if "(" in speaker else speaker
    pill_width = max(Inches(1.15), Inches(0.12 * len(pill_label) + 0.3))
    pill_width = min(pill_width, Inches(2.5))  # Cap pill width
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, pill_width, Inches(0.3))
    pill.fill.solid()
    pill.fill.fore_color.rgb = ACCENT_CLAY if learner else ACCENT_TEAL
    pill.line.fill.background()
    pill.adjustments[0] = 0.28
    tf = pill.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = pill_label
    run.font.name = FONT_LATIN
    run.font.size = SIZE_LABEL
    run.font.color.rgb = WHITE
    run.font.bold = True
    apply_run_font_metadata(run)
    add_textbox(
        slide,
        left + Inches(1.35),
        top - Inches(0.03),
        width - Inches(1.35),
        Inches(0.42),
        thai,
        font_name=FONT_THAI,
        font_size=SIZE_THAI_SMALL,
        bold=True,
        translit_entries=[(thai, translit)],
    )
    add_textbox(
        slide,
        left + Inches(1.35),
        top + Inches(0.28),
        width - Inches(1.35),
        Inches(0.25),
        english,
        font_name=FONT_LATIN,
        font_size=Pt(13),
        color=INK_MEDIUM,
    )


def fit_picture(slide, image_path: Path, left, top, max_width, max_height):
    with Image.open(image_path) as img:
        width, height = img.size
    image_ratio = width / height
    box_ratio = float(max_width) / float(max_height)
    if image_ratio >= box_ratio:
        final_width = int(max_width)
        final_height = int(float(max_width) / image_ratio)
    else:
        final_height = int(max_height)
        final_width = int(float(max_height) * image_ratio)
    slide.shapes.add_picture(str(image_path), left, top, width=final_width, height=final_height)


def render_opener(slide, slide_data: dict[str, Any], row: dict[str, str], translit_entries: list[tuple[str, str]] | None = None):
    add_divider(slide, 0, 0, LEFT_ZONE_W, ACCENT_GOLD, Pt(6))
    eyebrow = f"{row.get('module_title', '').strip()}  ·  {row.get('cefr_band', '').strip()}".strip(" ·")
    add_textbox(slide, CONTENT_LEFT, Inches(1.75), CONTENT_WIDTH_BESIDE_PIP, Inches(0.25), eyebrow or "Immersion Thai with Nine", font_name=FONT_LATIN, font_size=SIZE_LABEL, color=ACCENT_GOLD, bold=True, translit_entries=translit_entries)
    add_textbox(slide, CONTENT_LEFT, Inches(2.18), CONTENT_WIDTH_BESIDE_PIP, Inches(0.36), row.get("lesson_id", slide_data["id"]), font_name=FONT_LATIN, font_size=Pt(20), color=INK_MEDIUM, translit_entries=translit_entries)
    add_textbox(slide, CONTENT_LEFT, Inches(2.78), CONTENT_WIDTH_BESIDE_PIP, Inches(1.2), slide_data["title"], font_name=FONT_LATIN, font_size=Pt(32), bold=True, translit_entries=translit_entries)
    note_lines = slide_data["textBlocks"][0]["lines"][:2]
    add_bullet_block(slide, CONTENT_LEFT, Inches(4.5), CONTENT_WIDTH_BESIDE_PIP, "Lesson focus", note_lines, ACCENT_GOLD, translit_entries)
    add_textbox(slide, CONTENT_LEFT, Inches(6.2), CONTENT_WIDTH, Inches(0.25), "Immersion Thai with Nine", font_name=FONT_LATIN, font_size=SIZE_LABEL, color=INK_LIGHT)


def render_objectives(slide, slide_data: dict[str, Any], translit_entries: list[tuple[str, str]] | None = None):
    add_section_header(slide, slide_data["title"], "Lesson objectives", translit_entries)
    lines = slide_data["textBlocks"][0]["lines"]
    card_w = Inches(3.5)
    positions = [
        (CONTENT_LEFT, Inches(1.55)),
        (CONTENT_LEFT + Inches(3.85), Inches(1.55)),
        (CONTENT_LEFT, Inches(3.35)),
        (CONTENT_LEFT + Inches(3.85), Inches(3.35)),
    ]
    for idx, line in enumerate(lines[:4]):
        left, top = positions[idx]
        add_card(slide, left, top, card_w, Inches(1.35))
        add_textbox(slide, left + Inches(0.24), top + Inches(0.18), Inches(0.4), Inches(0.25), str(idx + 1), font_name=FONT_LATIN, font_size=SIZE_LABEL, color=ACCENT_GOLD, bold=True)
        add_textbox(slide, left + Inches(0.24), top + Inches(0.46), card_w - Inches(0.48), Inches(0.7), line, font_name=FONT_LATIN, font_size=SIZE_BODY_SMALL, color=INK_DARK, bold=False, translit_entries=translit_entries)


def render_teaching(
    slide,
    slide_data: dict[str, Any],
    lesson_root: Path,
    translit_entries: list[tuple[str, str]] | None = None,
):
    add_section_header(slide, slide_data["title"], "Teaching slide", translit_entries)
    content_left = CONTENT_LEFT
    top = Inches(1.45)
    # Teaching content is beside PiP — constrain to not overlap
    usable_width = CONTENT_WIDTH_BESIDE_PIP
    left_column_width = usable_width * 0.55
    right_column_left = CONTENT_LEFT + usable_width * 0.58
    right_column_width = usable_width * 0.40

    # Text-only slides — no image fetching or rendering. Add images in Canva if needed.
    triplet_blocks = next((block for block in slide_data["textBlocks"] if block["kind"] == "triplet-list"), None)
    drill_block = next((block for block in slide_data["textBlocks"] if block["kind"] in {"bullet-list", "note"}), None)

    content_bottom = top
    if triplet_blocks:
        # Use full constrained width for triplet cards
        content_bottom = add_triplet_rows(slide, content_left, top, usable_width, triplet_blocks["lines"])
    elif slide_data["thaiFocus"]:
        first = slide_data["thaiFocus"][0]
        add_phrase_card(slide, content_left, top, usable_width, first["thai"], first["translit"], first["english"])
        content_bottom = top + Inches(1.92)

    if drill_block:
        # Place drill block below content — but only if there's room on the slide
        drill_top = content_bottom + Inches(0.15)
        min_drill_height = Inches(0.9)  # heading + at least 1 line
        if drill_top + min_drill_height <= SLIDE_HEIGHT:
            add_bullet_block(
                slide,
                content_left,
                drill_top,
                CONTENT_WIDTH,
                drill_block.get("heading") or "Practice",
                drill_block["lines"][:4],
                ACCENT_CLAY,
                translit_entries=translit_entries,
            )
        else:
            # Not enough room — put drill text in speaker notes only
            notes = slide.notes_slide.notes_text_frame
            notes.text += "\n\n[Practice — moved to notes, no room on slide]\n"
            for line in drill_block["lines"][:4]:
                notes.text += f"• {line}\n"


def render_roleplay(slide, slide_data: dict[str, Any], translit_entries: list[tuple[str, str]] | None = None):
    # Use short "Roleplay" heading; full scenario goes to speaker notes only
    add_section_header(slide, "Roleplay", "Roleplay", translit_entries)
    current_top = Inches(1.2)
    for index, line in enumerate(slide_data["textBlocks"][0]["lines"][:6]):
        speaker, thai, translit, english = [part.strip() for part in line.split("|", 3)]
        # Use constrained width when beside PiP, full width when below
        turn_width = CONTENT_WIDTH_BESIDE_PIP if current_top < PIP_BOTTOM else CONTENT_WIDTH
        add_dialogue_turn(slide, CONTENT_LEFT, current_top, turn_width, speaker, thai, translit, english, learner=index % 2 == 1)
        current_top += Inches(0.85)


def render_recap(slide, slide_data: dict[str, Any], translit_entries: list[tuple[str, str]] | None = None):
    add_section_header(slide, slide_data["title"], "Recap", translit_entries)
    # Top block is beside PiP — constrained width
    block_bottom = add_bullet_block(slide, CONTENT_LEFT, Inches(1.55), CONTENT_WIDTH_BESIDE_PIP, "What you can now do", slide_data["textBlocks"][0]["lines"][:5], ACCENT_GOLD, translit_entries)
    takeaway = slide_data["speakerNotes"][-1] if slide_data["speakerNotes"] else "Ready to record."
    # Bottom block positioned below top block with gap — full width OK (below PiP)
    remember_top = max(block_bottom + Inches(0.3), Inches(4.55))
    add_bullet_block(slide, CONTENT_LEFT, remember_top, CONTENT_WIDTH, "Remember", [takeaway], ACCENT_TEAL, translit_entries)


def render_closing(slide, slide_data: dict[str, Any], translit_entries: list[tuple[str, str]] | None = None):
    add_section_header(slide, slide_data["title"], "Closing", translit_entries)
    add_bullet_block(slide, CONTENT_LEFT, Inches(1.75), CONTENT_WIDTH_BESIDE_PIP, "Next steps", slide_data["textBlocks"][0]["lines"][:4], ACCENT_CLAY, translit_entries)


def build_teaching_slide(
    slide_id: str,
    title: str,
    section: dict[str, Any],
    lesson_title: str,
    lesson_root: Path,
    assets_dir: Path,
) -> tuple[dict[str, Any], list[dict[str, Any]]]:
    visual_plan = section.get("visualPlan") or {}
    image_usage = choose_image_usage(visual_plan)
    text_blocks = [
        {
            "id": f"{slide_id}-triplets",
            "kind": "triplet-list",
            "heading": "Key language",
            "lines": section.get("onScreenBullets", [])[:6],
        },
        {
            "id": f"{slide_id}-practice",
            "kind": "bullet-list",
            "heading": "Practice",
            "lines": section.get("drills", [])[:4] or [section.get("purpose", "")],
        },
    ]
    assets: list[dict[str, Any]] = []
    provenance: list[dict[str, Any]] = []
    if image_usage == "real-image":
        search_queries = ((visual_plan.get("imageSupport") or {}).get("searchQueries") or [])
        query = search_queries[0] if search_queries else f"{lesson_title} {title} Thailand"
        usage_notes = (visual_plan.get("imageSupport") or {}).get(
            "rationale",
            "Use as a left-panel supporting image only.",
        )
        asset, provenance_entry = resolve_asset(slide_id, lesson_root, assets_dir, query, usage_notes)
        assets.append(asset)
        provenance.append(provenance_entry)

    slide = {
        "id": slide_id,
        "role": "teaching",
        "title": title,
        "estimatedSeconds": estimate_seconds(section.get("spokenNarration", []), 18),
        "layout": visual_plan.get("leftPanelLayout", "focus-card"),
        "speakerNotes": section.get("spokenNarration", []) or [section.get("purpose", title)],
        "textBlocks": text_blocks,
        "thaiFocus": [
            {
                "thai": lex.get("thai", ""),
                "translit": lex.get("translit", ""),
                "english": lex.get("english", ""),
            }
            for lex in section.get("languageFocus", [])
        ],
        "visualStrategy": {
            "onScreenGoal": visual_plan.get("onScreenGoal", section.get("purpose", title)),
            "teachingVisuals": visual_plan.get("teachingVisuals") or [title],
            "teacherCues": visual_plan.get("teacherCues") or ["Keep teaching content in the left zone."],
            "imageUsage": image_usage,
            "rationale": ((visual_plan.get("imageSupport") or {}).get("rationale")) or section.get("purpose", title),
        },
        "assets": assets,
    }
    return slide, provenance


def build_deck_source(
    repo_root: Path,
    lesson_id: str,
    script: dict[str, Any],
    row: dict[str, str],
    lesson_root: Path,
) -> tuple[dict[str, Any], dict[str, Any]]:
    assets_dir = lesson_root / "slide-assets"
    slides: list[dict[str, Any]] = []
    provenance_assets: list[dict[str, Any]] = []

    slides.append(
        {
            "id": "slide-01-opener",
            "role": "opener",
            "title": script.get("title", lesson_id),
            "estimatedSeconds": estimate_seconds([script.get("objective", "")], 12),
            "layout": "lesson-opener",
            "speakerNotes": [
                script.get("teachingFrame", {}).get("openingHook") or script.get("objective", ""),
                script.get("objective", ""),
            ],
            "textBlocks": [
                {
                    "id": "opener-focus",
                    "kind": "bullet-list",
                    "heading": "Lesson focus",
                    "lines": [
                        _you(row.get("lesson_primary_outcome", "") or script.get("objective", "")),
                        _you(row.get("speaking_target", "") or script.get("teachingFrame", {}).get("learnerTakeaway", "")),
                    ],
                }
            ],
            "thaiFocus": [],
            "visualStrategy": {
                "onScreenGoal": "Introduce the lesson clearly and set the recording tone.",
                "teachingVisuals": [script.get("title", lesson_id), row.get("module_title", "")],
                "teacherCues": ["Use the deck title as the recording opener.", "Keep the right third clear."],
                "imageUsage": "text-only",
                "rationale": "The opener should establish identity and objective without visual clutter.",
            },
            "assets": [],
        }
    )

    objective_lines = [
        _you(row.get("lesson_primary_outcome", "")),
        _you(row.get("lesson_secondary_outcome", "")),
        _you(row.get("speaking_target", "")),
        _you(row.get("listening_target", "")),
    ]
    objective_lines = [line for line in objective_lines if line][:4] or [_you(script.get("objective", ""))]
    slides.append(
        {
            "id": "slide-02-objectives",
            "role": "objectives",
            "title": "What you will learn",
            "estimatedSeconds": estimate_seconds(objective_lines, 14),
            "layout": "objectives-list",
            "speakerNotes": objective_lines,
            "textBlocks": [
                {
                    "id": "objectives-list",
                    "kind": "bullet-list",
                    "heading": "Objectives",
                    "lines": objective_lines,
                }
            ],
            "thaiFocus": [],
            "visualStrategy": {
                "onScreenGoal": "Make the lesson goals visible before the teaching slides begin.",
                "teachingVisuals": objective_lines,
                "teacherCues": ["Pause briefly after each objective.", "Keep the slide readable from a distance."],
                "imageUsage": "text-only",
                "rationale": "Objectives are clearer as cards than as decorative visuals.",
            },
            "assets": [],
        }
    )

    for index, section in enumerate(script.get("sections", []), start=3):
        slide_id = build_slide_id(index, section.get("heading", f"section-{index - 2}"))
        slide, provenance = build_teaching_slide(
            slide_id,
            section.get("heading", f"Section {index - 2}"),
            section,
            script.get("title", lesson_id),
            lesson_root,
            assets_dir,
        )
        slides.append(slide)
        provenance_assets.extend(provenance)

    roleplay_lines = [
        " | ".join(
            [
                # Normalize speaker: "Learner" → "You", strip parenthetical cues
                (lambda s: s.split("(")[0].strip())(line.get("speaker", "").replace("Learner", "You")),
                line.get("thai", ""),
                line.get("translit", ""),
                line.get("english", ""),
            ]
        )
        for line in script.get("roleplay", {}).get("lines", [])
    ]
    scenario = script.get("roleplay", {}).get("scenario", "Roleplay")
    slides.append(
        {
            "id": f"slide-{len(slides) + 1:02d}-roleplay",
            "role": "roleplay",
            "title": "Roleplay",
            "estimatedSeconds": estimate_seconds(roleplay_lines, 22),
            "layout": "roleplay-dialogue",
            "speakerNotes": [scenario] + (roleplay_lines or []),
            "textBlocks": [
                {
                    "id": "roleplay-dialogue",
                    "kind": "dialogue",
                    "heading": "Dialogue",
                    "lines": roleplay_lines[:6] or ["Teacher | สวัสดี | sà-wàt-dii | hello"],
                }
            ],
            "thaiFocus": [
                {
                    "thai": line.get("thai", ""),
                    "translit": line.get("translit", ""),
                    "english": line.get("english", ""),
                }
                for line in script.get("roleplay", {}).get("lines", [])
            ],
            "visualStrategy": {
                "onScreenGoal": "Support the roleplay recording with clean turn-taking.",
                "teachingVisuals": [script.get("roleplay", {}).get("scenario", "Roleplay ladder")],
                "teacherCues": ["Record against the dialogue ladder.", "Bottom-right PiP reserved for camera."],
                "imageUsage": "text-only",
                "rationale": "The dialogue itself is the teaching focus.",
            },
            "assets": [],
        }
    )

    recap_lines = [_you(line) for line in (script.get("recap", [])[:5] or [script.get("objective", "")])]
    slides.append(
        {
            "id": f"slide-{len(slides) + 1:02d}-recap",
            "role": "recap",
            "title": "Recap",
            "estimatedSeconds": estimate_seconds(recap_lines, 12),
            "layout": "recap-checklist",
            "speakerNotes": recap_lines,
            "textBlocks": [
                {
                    "id": "recap-list",
                    "kind": "recap-list",
                    "heading": "What you can now do",
                    "lines": recap_lines,
                }
            ],
            "thaiFocus": [],
            "visualStrategy": {
                "onScreenGoal": "Finish with a simple, memorable checklist.",
                "teachingVisuals": recap_lines,
                "teacherCues": ["Use this slide as the summary.", "Keep the ending calm and uncluttered."],
                "imageUsage": "text-only",
                "rationale": "A recap slide should be compact and memorable.",
            },
            "assets": [],
        }
    )

    closing_lines = [
        _you(script.get("teachingFrame", {}).get("learnerTakeaway", "")),
        f"Next: {row.get('module_title', '').strip()} continues.",
    ]
    closing_lines = [line for line in closing_lines if line]
    slides.append(
        {
            "id": f"slide-{len(slides) + 1:02d}-closing",
            "role": "closing",
            "title": "Closing",
            "estimatedSeconds": estimate_seconds(closing_lines, 10),
            "layout": "lesson-closing",
            "speakerNotes": closing_lines or [script.get("objective", "")],
            "textBlocks": [
                {
                    "id": "closing-list",
                    "kind": "closing-list",
                    "heading": "Takeaway",
                    "lines": closing_lines or [script.get("objective", "")],
                }
            ],
            "thaiFocus": [],
            "visualStrategy": {
                "onScreenGoal": "Give Nine a clean final slide to close the recording.",
                "teachingVisuals": closing_lines or [script.get("objective", "")],
                "teacherCues": ["Use the closing as the final spoken wrap-up."],
                "imageUsage": "text-only",
                "rationale": "The closing slide is a presenter anchor, not a new teaching beat.",
            },
            "assets": [],
        }
    )

    deck_source = {
        "schemaVersion": 1,
        "lessonId": lesson_id,
        "sourceScript": artifact_name(lesson_id, "script-master.json"),
        "canvas": {
            "width": 1920,
            "height": 1080,
            "leftTeachingFraction": 1.0,
            "rightCameraFraction": 0.0,
            "pipCameraWidth": 4.2,
            "pipCameraHeight": 3.15,
            "pipPosition": "bottom-right",
            "safeZoneLabel": "Bottom-right PiP reserved for recording / camera",
        },
        "theme": {
            "id": "thai-nine-canva-bridge-v1",
            "name": "Canva Bridge Deck",
            "thaiFont": FONT_THAI,
            "latinFont": FONT_LATIN,
            "backgroundColor": to_hex(BG_IVORY),
            "rightZoneTint": to_hex(BG_SAND_LIGHT),
            "accentColors": [to_hex(ACCENT_GOLD), to_hex(ACCENT_TEAL), to_hex(ACCENT_CLAY)],
        },
        "slides": slides,
    }

    provenance = {
        "schemaVersion": 1,
        "lessonId": lesson_id,
        "generatedAt": datetime.datetime.now(datetime.UTC).isoformat().replace("+00:00", "Z"),
        "assets": provenance_assets,
    }

    return deck_source, provenance


def render_deck(
    deck_source: dict[str, Any],
    row: dict[str, str],
    lesson_root: Path,
    output_pptx: Path,
    translit_entries: list[tuple[str, str]] | None = None,
) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for slide_data in deck_source["slides"]:
        slide = add_blank_slide(prs)
        role = slide_data["role"]
        if role == "opener":
            render_opener(slide, slide_data, row, translit_entries)
        elif role == "objectives":
            render_objectives(slide, slide_data, translit_entries)
        elif role == "teaching":
            render_teaching(slide, slide_data, lesson_root, translit_entries)
        elif role == "roleplay":
            render_roleplay(slide, slide_data, translit_entries)
        elif role == "recap":
            render_recap(slide, slide_data, translit_entries)
        else:
            render_closing(slide, slide_data, translit_entries)

    # Validate layout before saving
    layout_issues = validate_slide_layout(prs)
    if layout_issues:
        print("LAYOUT VALIDATION FAILED — deck not saved:", file=sys.stderr)
        for issue in layout_issues:
            print(f"  ✗ {issue}", file=sys.stderr)
        print(f"\n{len(layout_issues)} layout violation(s). Fix the renderer before producing decks.", file=sys.stderr)
        sys.exit(1)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    patch_theme_fonts(output_pptx)


def _rects_overlap(r1: tuple, r2: tuple) -> bool:
    """Check if two (left, top, right, bottom) rectangles overlap."""
    return r1[0] < r2[2] and r1[2] > r2[0] and r1[1] < r2[3] and r1[3] > r2[1]


def validate_slide_layout(prs: Presentation) -> list[str]:
    """Check all slides for boundary violations, PiP zone intrusions, and text overlaps.

    Returns a list of issue strings. Any non-empty list = hard fail.
    """
    issues: list[str] = []
    pip_x_emu = PIP_X
    pip_bottom_emu = PIP_BOTTOM
    slide_w_emu = SLIDE_WIDTH
    slide_h_emu = SLIDE_HEIGHT

    for slide_idx, slide in enumerate(prs.slides):
        all_shapes = []
        text_shapes = []

        for shape in slide.shapes:
            # Skip PiP placeholder itself
            if getattr(shape, "name", "") == "__pip_placeholder__":
                continue

            left = shape.left
            top = shape.top
            right = left + shape.width
            bottom = top + shape.height

            # Skip thin decorative bars (top accent stripe, bottom dividers)
            # These are full-width design elements < 0.15" tall
            if shape.height < 137160 and not (shape.has_text_frame and shape.text_frame.text.strip()):
                continue
            label = ""
            if shape.has_text_frame:
                label = shape.text_frame.text.strip()[:50]
            if not label:
                label = getattr(shape, "name", "shape")

            # Check 1: PiP zone intrusion — shape crosses into camera zone
            # Camera zone: top-right rectangle from (PIP_X, 0) to (SLIDE_WIDTH, PIP_BOTTOM)
            if top < pip_bottom_emu and right > pip_x_emu:
                issues.append(
                    f"Slide {slide_idx + 1}: PiP zone intrusion — "
                    f"'{label}' extends into camera zone "
                    f"(right={right / 914400:.1f}\" > PiP left={pip_x_emu / 914400:.1f}\")"
                )

            # Check 2: Off-screen bottom
            if bottom > slide_h_emu:
                issues.append(
                    f"Slide {slide_idx + 1}: off-screen bottom — "
                    f"'{label}' extends below slide "
                    f"(bottom={bottom / 914400:.1f}\" > slide height={slide_h_emu / 914400:.1f}\")"
                )

            # Check 3: Off-screen right
            if right > slide_w_emu:
                issues.append(
                    f"Slide {slide_idx + 1}: off-screen right — "
                    f"'{label}' extends past slide edge "
                    f"(right={right / 914400:.1f}\" > slide width={slide_w_emu / 914400:.1f}\")"
                )

            # Collect text shapes for overlap checks
            if shape.has_text_frame and shape.text_frame.text.strip():
                rect = (left, top, right, bottom)
                text_shapes.append((shape, rect, label))

        # Check 4: Text-text overlaps
        for i, (s1, r1, t1) in enumerate(text_shapes):
            for j, (s2, r2, t2) in enumerate(text_shapes[i + 1:], start=i + 1):
                if _rects_overlap(r1, r2):
                    # Skip overlaps within the same dialogue turn:
                    # pill, Thai text, and English text share similar top positions
                    tops_close = abs(r1[1] - r2[1]) < 320000  # < ~0.35" apart vertically
                    if tops_close:
                        continue  # Same dialogue turn — intentional stacking
                    # For non-turn elements, flag significant overlaps
                    y_overlap = min(r1[3], r2[3]) - max(r1[1], r2[1])
                    if y_overlap > 137160:  # > 0.15 inch
                        issues.append(
                            f"Slide {slide_idx + 1}: text overlap — "
                            f"'{t1}' and '{t2}'"
                        )
    return issues


def hex_to_rgb(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def color_tuple_to_rgb(value: tuple[int, int, int]) -> RGBColor:
    return RGBColor(value[0], value[1], value[2])


def emu_to_px(value: int, total_emu: int, total_px: int) -> int:
    return int(round((value / total_emu) * total_px))


def load_canva_contract(repo_root: Path) -> dict[str, Any]:
    return read_json(repo_root / "course" / "layout-contracts" / "canva" / "l001-reference-v1.json")


def contract_family_map(contract: dict[str, Any]) -> dict[str, dict[str, Any]]:
    return {family["name"]: family for family in contract.get("slideFamilies", [])}


def role_style(contract: dict[str, Any], role_name: str) -> dict[str, Any]:
    return dict((contract.get("fontRoles") or {}).get(role_name) or {})


def slide_background_path(index: int) -> str:
    return f"canva-backgrounds/slide-{index:02d}.png"


def safe_list_get(items: list[Any], index: int, default: Any = None) -> Any:
    if 0 <= index < len(items):
        return items[index]
    return default


def first_text_block_lines(slide_data: dict[str, Any]) -> list[str]:
    text_blocks = slide_data.get("textBlocks") or []
    if not text_blocks:
        return []
    return [str(line).strip() for line in (text_blocks[0].get("lines") or []) if str(line).strip()]


def make_canva_slide(
    contract: dict[str, Any],
    family_name: str,
    slide_id: str,
    title: str,
    values: dict[str, Any],
    background_path: str,
) -> dict[str, Any]:
    families = contract_family_map(contract)
    family = families[family_name]
    colors = (contract.get("theme") or {}).get("colors") or {}
    elements: list[dict[str, Any]] = []

    for placeholder in family.get("placeholders", []):
        placeholder_id = placeholder["id"]
        raw_value = values.get(placeholder_id)
        if raw_value is None:
            continue

        if placeholder.get("kind") == "image":
            local_path = str(raw_value).strip()
            if not local_path:
                continue
            elements.append(
                {
                    "id": placeholder_id,
                    "kind": "image",
                    "x": placeholder["x"],
                    "y": placeholder["y"],
                    "w": placeholder["w"],
                    "h": placeholder["h"],
                    "localPath": local_path,
                }
            )
            continue

        text_value = str(raw_value).strip()
        if not text_value:
            continue

        font_role = role_style(contract, placeholder.get("fontRole", "body"))
        elements.append(
            {
                "id": placeholder_id,
                "kind": "text",
                "x": placeholder["x"],
                "y": placeholder["y"],
                "w": placeholder["w"],
                "h": placeholder["h"],
                "fontName": font_role.get("fontName", FONT_LATIN),
                "fontSizePt": font_role.get("fontSizePt", 14),
                "color": colors.get(font_role.get("color", "ink"), "#24333D"),
                "bold": bool(font_role.get("bold", False)),
                "italic": bool(font_role.get("italic", False)),
                "align": placeholder.get("align") or font_role.get("align", "left"),
                "value": text_value,
            }
        )

    return {
        "id": slide_id,
        "title": title,
        "layoutFamily": family_name,
        "backgroundPath": background_path,
        "elements": elements,
    }


def translit_partner_id(element_id: str) -> str | None:
    if "_thai" in element_id:
        return element_id.replace("_thai", "_translit")
    if element_id == "thai":
        return "translit"
    if "_translit" in element_id:
        return element_id.replace("_translit", "_thai")
    if element_id == "translit":
        return "thai"
    return None


def is_translit_placeholder_id(element_id: str) -> bool:
    return element_id == "translit" or "_translit" in element_id


def canva_segments_from_tokens(
    tokens: list[dict[str, str]],
    element: dict[str, Any],
) -> list[dict[str, Any]]:
    segments: list[dict[str, Any]] = []
    default_font_name = str(element.get("fontName", FONT_LATIN))
    default_font_size = float(element.get("fontSizePt", 14))
    default_color = str(element.get("color", "#24333D"))
    default_bold = bool(element.get("bold", False))
    default_italic = bool(element.get("italic", False))

    for token in tokens:
        if token["kind"] == "plain":
            text = token["text"]
            if not text:
                continue
            segments.append(
                {
                    "text": text,
                    "fontName": default_font_name,
                    "fontSizePt": default_font_size,
                    "color": default_color,
                    "bold": default_bold,
                    "italic": default_italic,
                }
            )
            continue

        segments.append(
            {
                "text": token["thai"],
                "fontName": FONT_THAI,
                "fontSizePt": default_font_size,
                "color": default_color,
                "bold": default_bold,
                "italic": False,
            }
        )
        if token["translit"]:
            segments.append(
                {
                    "text": f" ({token['translit']})",
                    "fontName": FONT_TRANSLIT,
                    "fontSizePt": translit_font_size_value(default_font_size),
                    "color": to_hex(INK_MEDIUM),
                    "bold": False,
                    "italic": True,
                }
            )

    return segments


def inline_canva_segments(
    thai_text: str,
    translit_text: str,
    element: dict[str, Any],
) -> list[dict[str, Any]]:
    default_font_size = float(element.get("fontSizePt", 14))
    default_color = str(element.get("color", "#24333D"))
    default_bold = bool(element.get("bold", False))

    segments = [
        {
            "text": thai_text,
            "fontName": FONT_THAI,
            "fontSizePt": default_font_size,
            "color": default_color,
            "bold": default_bold,
            "italic": False,
        }
    ]
    if translit_text:
        segments.append(
            {
                "text": f" ({translit_text})",
                "fontName": FONT_TRANSLIT,
                "fontSizePt": translit_font_size_value(default_font_size),
                "color": to_hex(INK_MEDIUM),
                "bold": False,
                "italic": True,
            }
        )
    return segments


def normalize_canva_elements(
    elements: list[dict[str, Any]],
    translit_entries: list[tuple[str, str]],
) -> list[dict[str, Any]]:
    element_map = {str(element.get("id", "")): element for element in elements}
    normalized: list[dict[str, Any]] = []
    skipped_ids: set[str] = set()

    for element in elements:
        element_id = str(element.get("id", ""))
        if element_id in skipped_ids:
            continue
        if element.get("kind") != "text":
            normalized.append(element)
            continue

        partner_id = translit_partner_id(element_id)
        if is_translit_placeholder_id(element_id) and partner_id and partner_id in element_map:
            continue

        if partner_id and partner_id in element_map and ("_thai" in element_id or element_id == "thai"):
            partner = element_map[partner_id]
            thai_value = str(element.get("value", "")).strip()
            translit_value = str(partner.get("value", "")).strip()
            if thai_value:
                merged = dict(element)
                merged["value"] = inline_phrase(thai_value, translit_value)
                merged["segments"] = inline_canva_segments(thai_value, translit_value, element)
                x0 = min(float(element["x"]), float(partner["x"]))
                y0 = min(float(element["y"]), float(partner["y"]))
                x1 = max(float(element["x"]) + float(element["w"]), float(partner["x"]) + float(partner["w"]))
                y1 = max(float(element["y"]) + float(element["h"]), float(partner["y"]) + float(partner["h"]))
                merged["x"] = x0
                merged["y"] = y0
                merged["w"] = x1 - x0
                merged["h"] = y1 - y0
                merged["fontName"] = FONT_THAI
                normalized.append(merged)
                skipped_ids.add(partner_id)
                continue

        text_value = str(element.get("value", ""))
        if "\n" in text_value:
            normalized.append(element)
            continue

        tokens = inline_tokens(text_value, translit_entries)
        if any(token["kind"] == "thai" for token in tokens):
            updated = dict(element)
            updated["value"] = tokens_to_text(tokens)
            updated["segments"] = canva_segments_from_tokens(tokens, element)
            normalized.append(updated)
            continue

        normalized.append(element)

    return normalized


def normalize_canva_slides(
    slides: list[dict[str, Any]],
    translit_entries: list[tuple[str, str]],
) -> list[dict[str, Any]]:
    return [
        {
            **slide,
            "elements": normalize_canva_elements(list(slide.get("elements", [])), translit_entries),
        }
        for slide in slides
    ]


def build_l001_reference_canva_slides(
    lesson_id: str,
    script: dict[str, Any],
    row: dict[str, str],
    contract: dict[str, Any],
) -> list[dict[str, Any]]:
    sections = script.get("sections") or []
    s1 = safe_list_get(sections, 0, {}) or {}
    s2 = safe_list_get(sections, 1, {}) or {}
    s3 = safe_list_get(sections, 2, {}) or {}
    s4 = safe_list_get(sections, 3, {}) or {}

    s1_focus = s1.get("languageFocus") or []
    s2_focus = s2.get("languageFocus") or []
    s3_focus = s3.get("languageFocus") or []
    s4_focus = s4.get("languageFocus") or []
    recap_lines = [str(line).strip() for line in script.get("recap", []) if str(line).strip()]
    roleplay_lines = script.get("roleplay", {}).get("lines", []) or []

    objectives = [
        "1. Say hello politely with a correct polite ending",
        "2. Use ขอบคุณ and ขอโทษ in the right social moments",
        "3. Distinguish ครับ, ค่ะ, and คะ at a practical level",
        "4. Answer with ใช่ or ไม่ใช่ and soften with ไม่เป็นไร",
        "5. Deliver a short first-contact exchange with politeness",
    ]

    slides: list[dict[str, Any]] = []
    lesson_title = row.get("lesson_title") or script.get("title", lesson_id)
    next_lesson_line = "Next lesson: names and countries - introducing yourself."

    slides.append(
        make_canva_slide(
            contract,
            "opener",
            "canva-slide-01-opener",
            lesson_title,
            {
                "eyebrow": f"{row.get('module_title', 'Module 1')}  ·  {row.get('cefr_band', 'Level A0')}",
                "lesson_id": lesson_id,
                "title": lesson_title,
                "brand": "Immersion Thai with Nine",
            },
            slide_background_path(1),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "objectives-list",
            "canva-slide-02-objectives",
            "What you will learn",
            {
                "title": "What you will learn",
                "objective_1": objectives[0],
                "objective_2": objectives[1],
                "objective_3": objectives[2],
                "objective_4": objectives[3],
                "objective_5": objectives[4],
            },
            slide_background_path(2),
        )
    )

    first_greeting = safe_list_get(s1_focus, 0, {}) or {}
    slides.append(
        make_canva_slide(
            contract,
            "single-focus-card",
            "canva-slide-03-focus",
            "Say hello the Thai way",
            {
                "section_label": "Section 1",
                "title": "Say hello the Thai way",
                "thai": first_greeting.get("thai", "สวัสดี"),
                "translit": first_greeting.get("translit", "sà-wàt-dii"),
                "english": first_greeting.get("english", "hello"),
                "note": "This is your safe, polite greeting for a first meeting.\nTreat politeness as part of the whole line, not as an extra decoration.",
            },
            slide_background_path(3),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "stacked-phrase-cards",
            "canva-slide-04-greeting-stack",
            "The greeting stack",
            {
                "section_label": "Section 1",
                "title": "The greeting stack",
                "subtitle": "Learn the greeting as a full social line, not just the base word.",
                "card_1_thai": (safe_list_get(s1_focus, 0, {}) or {}).get("thai", ""),
                "card_1_translit": (safe_list_get(s1_focus, 0, {}) or {}).get("translit", ""),
                "card_1_english": (safe_list_get(s1_focus, 0, {}) or {}).get("english", ""),
                "card_2_thai": (safe_list_get(s1_focus, 1, {}) or {}).get("thai", ""),
                "card_2_translit": (safe_list_get(s1_focus, 1, {}) or {}).get("translit", ""),
                "card_2_english": (safe_list_get(s1_focus, 1, {}) or {}).get("english", ""),
                "card_3_thai": (safe_list_get(s1_focus, 2, {}) or {}).get("thai", ""),
                "card_3_translit": (safe_list_get(s1_focus, 2, {}) or {}).get("translit", ""),
                "card_3_english": (safe_list_get(s1_focus, 2, {}) or {}).get("english", ""),
            },
            slide_background_path(4),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "dual-card-shared-reply",
            "canva-slide-05-courtesy-pairs",
            "Thank, apologise, and soften",
            {
                "section_label": "Section 2",
                "title": "Thank, apologise, and soften",
                "left_label": "Close warmly",
                "right_label": "Interrupt or repair",
                "left_thai": (safe_list_get(s2_focus, 0, {}) or {}).get("thai", ""),
                "left_translit": (safe_list_get(s2_focus, 0, {}) or {}).get("translit", ""),
                "left_english": (safe_list_get(s2_focus, 0, {}) or {}).get("english", ""),
                "right_thai": (safe_list_get(s2_focus, 1, {}) or {}).get("thai", ""),
                "right_translit": (safe_list_get(s2_focus, 1, {}) or {}).get("translit", ""),
                "right_english": (safe_list_get(s2_focus, 1, {}) or {}).get("english", ""),
                "shared_label": "Natural reply to both",
                "shared_thai": (safe_list_get(s2_focus, 2, {}) or {}).get("thai", ""),
                "shared_translit": (safe_list_get(s2_focus, 2, {}) or {}).get("translit", ""),
                "shared_english": "Smooths the moment after thanks or a small apology.",
            },
            slide_background_path(5),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "social-job-stack",
            "canva-slide-06-social-jobs",
            "How these phrases work together",
            {
                "section_label": "Section 2",
                "title": "How these phrases work together",
                "subtitle": "Each phrase has a social job in the interaction:",
                "row_1_thai": (safe_list_get(s2_focus, 0, {}) or {}).get("thai", ""),
                "row_1_translit": (safe_list_get(s2_focus, 0, {}) or {}).get("translit", ""),
                "row_1_role": "Closes a moment warmly",
                "row_1_usage": "After help, service, or a small kindness",
                "row_2_thai": (safe_list_get(s2_focus, 1, {}) or {}).get("thai", ""),
                "row_2_translit": (safe_list_get(s2_focus, 1, {}) or {}).get("translit", ""),
                "row_2_role": "Opens a repair",
                "row_2_usage": "To interrupt gently or apologise",
                "row_3_thai": (safe_list_get(s2_focus, 2, {}) or {}).get("thai", ""),
                "row_3_translit": (safe_list_get(s2_focus, 2, {}) or {}).get("translit", ""),
                "row_3_role": "Smooths everything out",
                "row_3_usage": "Answers a thank-you or small apology",
            },
            slide_background_path(6),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "polite-ending-stack",
            "canva-slide-07-polite-endings",
            "Pick the right polite ending",
            {
                "section_label": "Section 3",
                "title": "Pick the right polite ending",
                "subtitle": "Thai polite endings do real work - they shape the tone of the whole line.",
                "row_1_thai": (safe_list_get(s3_focus, 0, {}) or {}).get("thai", ""),
                "row_1_support": f"{(safe_list_get(s3_focus, 0, {}) or {}).get('translit', '')}  -  male polite ending\nUse for statements and questions",
                "row_2_thai": (safe_list_get(s3_focus, 1, {}) or {}).get("thai", ""),
                "row_2_support": f"{(safe_list_get(s3_focus, 1, {}) or {}).get('translit', '')}  -  female statement ending\nUse for statements and answers",
                "row_3_thai": (safe_list_get(s3_focus, 2, {}) or {}).get("thai", ""),
                "row_3_support": f"{(safe_list_get(s3_focus, 2, {}) or {}).get('translit', '')}  -  female question ending\nUse when asking questions",
                "rule_note": "Beginner rule: learn ครับ and ค่ะ first for safe polite statements.",
            },
            slide_background_path(7),
        )
    )

    drill_prompts = [
        (safe_list_get(s1_focus, 0, {}) or {}),
        (safe_list_get(s2_focus, 0, {}) or {}),
        (safe_list_get(s2_focus, 1, {}) or {}),
        (safe_list_get(s4_focus, 0, {}) or {}),
        (safe_list_get(s2_focus, 2, {}) or {}),
    ]
    drill_values: dict[str, Any] = {
        "section_label": "Section 3",
        "title": "Substitution drill",
        "subtitle": "Say each line with the correct polite ending:",
    }
    for row_index, prompt in enumerate(drill_prompts, start=1):
        thai = str(prompt.get("thai", "")).strip()
        translit = str(prompt.get("translit", "")).strip()
        drill_values[f"row_{row_index}_prompt_thai"] = f"{row_index}. {thai}___" if thai else ""
        drill_values[f"row_{row_index}_prompt_translit"] = f"{translit} ___" if translit else ""
        drill_values[f"row_{row_index}_instruction"] = "Add ครับ or ค่ะ"

    slides.append(
        make_canva_slide(
            contract,
            "drill-rows",
            "canva-slide-08-drills",
            "Substitution drill",
            drill_values,
            slide_background_path(8),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "response-triad",
            "canva-slide-09-responses",
            "Answer yes, no, and no problem",
            {
                "section_label": "Section 4",
                "title": "Answer yes, no, and no problem",
                "subtitle": "Three response tools - two answer the idea, one softens the moment.",
                "answers_label": "Direct answers",
                "row_1_thai": (safe_list_get(s4_focus, 0, {}) or {}).get("thai", ""),
                "row_1_translit": (safe_list_get(s4_focus, 0, {}) or {}).get("translit", ""),
                "row_1_meaning": "yes / that is right",
                "row_1_usage": "Confirms the idea",
                "row_2_thai": (safe_list_get(s4_focus, 1, {}) or {}).get("thai", ""),
                "row_2_translit": (safe_list_get(s4_focus, 1, {}) or {}).get("translit", ""),
                "row_2_meaning": "no / not right",
                "row_2_usage": "Corrects the idea - not aggressive",
                "softener_label": "Social softener",
                "row_3_thai": (safe_list_get(s4_focus, 2, {}) or {}).get("thai", ""),
                "row_3_translit": (safe_list_get(s4_focus, 2, {}) or {}).get("translit", ""),
                "row_3_meaning": "it is okay / no problem",
                "row_3_usage": "Softens the situation - not a yes/no answer",
            },
            slide_background_path(9),
        )
    )

    roleplay_values: dict[str, Any] = {
        "title": "Roleplay: First check-in",
        "subtitle": "At a reception desk - practise the full exchange",
    }
    for turn_index, line in enumerate(roleplay_lines[:6], start=1):
        roleplay_values[f"turn_{turn_index}_speaker"] = line.get("speaker", f"Turn {turn_index}")
        roleplay_values[f"turn_{turn_index}_thai"] = line.get("thai", "")
        roleplay_values[f"turn_{turn_index}_support"] = line.get("english", "")

    slides.append(
        make_canva_slide(
            contract,
            "roleplay-ladder",
            "canva-slide-10-roleplay",
            "Roleplay: First check-in",
            roleplay_values,
            slide_background_path(10),
        )
    )

    recap_values = {
        "title": "Recap",
        "item_1": safe_list_get(recap_lines, 0, ""),
        "item_2": safe_list_get(recap_lines, 1, ""),
        "item_3": safe_list_get(recap_lines, 2, ""),
        "item_4": safe_list_get(recap_lines, 3, ""),
        "item_5": "Practise the full exchange - greeting to thank-you to no-problem",
    }
    slides.append(
        make_canva_slide(
            contract,
            "recap-checklist",
            "canva-slide-11-recap",
            "Recap",
            recap_values,
            slide_background_path(11),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "closing",
            "canva-slide-12-closing",
            "Closing",
            {
                "lesson_id": lesson_id,
                "title": "You can now greet, thank,\napologise, and respond politely.",
                "line_1": "You have the smallest courtesy toolkit that already works.",
                "line_2": next_lesson_line,
                "brand": "Immersion Thai with Nine",
            },
            slide_background_path(12),
        )
    )

    return slides


def build_l002_reference_canva_slides(
    lesson_id: str,
    script: dict[str, Any],
    row: dict[str, str],
    contract: dict[str, Any],
) -> list[dict[str, Any]]:
    sections = script.get("sections") or []
    s1 = safe_list_get(sections, 0, {}) or {}
    s2 = safe_list_get(sections, 1, {}) or {}
    s3 = safe_list_get(sections, 2, {}) or {}
    s4 = safe_list_get(sections, 3, {}) or {}

    s1_focus = s1.get("languageFocus") or []
    s2_focus = s2.get("languageFocus") or []
    s3_focus = s3.get("languageFocus") or []
    s4_focus = s4.get("languageFocus") or []
    recap_lines = [str(line).strip() for line in script.get("recap", []) if str(line).strip()]
    roleplay_lines = script.get("roleplay", {}).get("lines", []) or []
    lesson_title = row.get("lesson_title") or script.get("title", lesson_id)

    objectives = [
        "1. Ask someone's name with คุณชื่ออะไร",
        "2. Answer with ฉันชื่อ... or ผมชื่อ...",
        "3. Say where you are from with มาจาก...",
        "4. Choose ฉัน and ผม naturally",
        "5. Deliver a short self-introduction",
    ]

    slides: list[dict[str, Any]] = []

    slides.append(
        make_canva_slide(
            contract,
            "opener",
            "canva-slide-01-opener",
            lesson_title,
            {
                "eyebrow": f"{row.get('module_title', 'Module 1')}  ·  {row.get('cefr_band', 'Level A0')}",
                "lesson_id": lesson_id,
                "title": lesson_title,
                "brand": "Immersion Thai with Nine",
            },
            slide_background_path(1),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "objectives-list",
            "canva-slide-02-objectives",
            "What you will learn",
            {
                "title": "What you will learn",
                "objective_1": objectives[0],
                "objective_2": objectives[1],
                "objective_3": objectives[2],
                "objective_4": objectives[3],
                "objective_5": objectives[4],
            },
            slide_background_path(2),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "dual-card-shared-reply",
            "canva-slide-03-name-question",
            "Ask and answer the name question",
            {
                "section_label": "Section 1",
                "title": "Ask and answer the name question",
                "left_label": "Ask",
                "right_label": "Answer",
                "left_thai": (safe_list_get(s1_focus, 1, {}) or {}).get("thai", ""),
                "left_translit": (safe_list_get(s1_focus, 1, {}) or {}).get("translit", ""),
                "left_english": (safe_list_get(s1_focus, 1, {}) or {}).get("english", ""),
                "right_thai": (safe_list_get(s1_focus, 2, {}) or {}).get("thai", ""),
                "right_translit": (safe_list_get(s1_focus, 2, {}) or {}).get("translit", ""),
                "right_english": (safe_list_get(s1_focus, 2, {}) or {}).get("english", ""),
                "shared_label": "Male speaker",
                "shared_thai": (safe_list_get(s1_focus, 3, {}) or {}).get("thai", ""),
                "shared_translit": (safe_list_get(s1_focus, 3, {}) or {}).get("translit", ""),
                "shared_english": "Same answer pattern with a male self-reference.",
            },
            slide_background_path(3),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "response-triad",
            "canva-slide-04-pronoun-choice",
            "Use ฉัน and ผม naturally",
            {
                "section_label": "Section 2",
                "title": "Use ฉัน and ผม naturally",
                "subtitle": "Choose the self-word that matches the speaker, then keep the rest of the sentence stable.",
                "answers_label": "Core self-words",
                "row_1_thai": (safe_list_get(s2_focus, 0, {}) or {}).get("thai", ""),
                "row_1_translit": (safe_list_get(s2_focus, 0, {}) or {}).get("translit", ""),
                "row_1_meaning": (safe_list_get(s2_focus, 0, {}) or {}).get("english", ""),
                "row_1_usage": "Common self-word in this course",
                "row_2_thai": (safe_list_get(s2_focus, 1, {}) or {}).get("thai", ""),
                "row_2_translit": (safe_list_get(s2_focus, 1, {}) or {}).get("translit", ""),
                "row_2_meaning": (safe_list_get(s2_focus, 1, {}) or {}).get("english", ""),
                "row_2_usage": "Use when the speaker is male",
                "softener_label": "Full name line",
                "row_3_thai": "ฉันชื่อ... / ผมชื่อ...",
                "row_3_translit": "chǎn chûue... / phǒm chûue...",
                "row_3_meaning": "my name is ...",
                "row_3_usage": "Swap only the first word, then keep the name frame intact",
            },
            slide_background_path(4),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "stacked-phrase-cards",
            "canva-slide-05-where-from",
            "Say where you are from",
            {
                "section_label": "Section 3",
                "title": "Say where you are from",
                "subtitle": "Keep มาจาก and swap only the country.",
                "card_1_thai": (safe_list_get(s3_focus, 1, {}) or {}).get("thai", ""),
                "card_1_translit": (safe_list_get(s3_focus, 1, {}) or {}).get("translit", ""),
                "card_1_english": (safe_list_get(s3_focus, 1, {}) or {}).get("english", ""),
                "card_2_thai": (safe_list_get(s3_focus, 2, {}) or {}).get("thai", ""),
                "card_2_translit": (safe_list_get(s3_focus, 2, {}) or {}).get("translit", ""),
                "card_2_english": (safe_list_get(s3_focus, 2, {}) or {}).get("english", ""),
                "card_3_thai": (safe_list_get(s3_focus, 3, {}) or {}).get("thai", ""),
                "card_3_translit": (safe_list_get(s3_focus, 3, {}) or {}).get("translit", ""),
                "card_3_english": (safe_list_get(s3_focus, 3, {}) or {}).get("english", ""),
            },
            slide_background_path(5),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "drill-rows",
            "canva-slide-06-intro-sequence",
            "Build a full self-introduction",
            {
                "section_label": "Section 4",
                "title": "Build a full self-introduction",
                "subtitle": "Use the sequence below as one short introduction.",
                "row_1_prompt_thai": "1. สวัสดีครับ / สวัสดีค่ะ",
                "row_1_prompt_translit": "sà-wàt-dii khráp / sà-wàt-dii khâ",
                "row_1_instruction": "Start politely",
                "row_2_prompt_thai": "2. ฉันชื่อ... / ผมชื่อ...",
                "row_2_prompt_translit": "chǎn chûue... / phǒm chûue...",
                "row_2_instruction": "Say your name",
                "row_3_prompt_thai": "3. ฉันมาจาก... / ผมมาจาก...",
                "row_3_prompt_translit": "chǎn maa-jàak... / phǒm maa-jàak...",
                "row_3_instruction": "Add your country",
                "row_4_prompt_thai": (safe_list_get(s4_focus, 1, {}) or {}).get("thai", "ขอบคุณ"),
                "row_4_prompt_translit": (safe_list_get(s4_focus, 1, {}) or {}).get("translit", "khàawp-khun"),
                "row_4_instruction": "Close warmly if needed",
                "row_5_prompt_thai": (safe_list_get(s4_focus, 0, {}) or {}).get("thai", "สวัสดี"),
                "row_5_prompt_translit": (safe_list_get(s4_focus, 0, {}) or {}).get("translit", "sà-wàt-dii"),
                "row_5_instruction": "Keep the whole exchange smooth and polite",
            },
            slide_background_path(6),
        )
    )

    roleplay_values: dict[str, Any] = {
        "title": "Roleplay: First self-introduction",
        "subtitle": "Listen, repeat, then try the exchange yourself.",
    }
    for turn_index, line in enumerate(roleplay_lines[:6], start=1):
        roleplay_values[f"turn_{turn_index}_speaker"] = line.get("speaker", f"Turn {turn_index}")
        roleplay_values[f"turn_{turn_index}_thai"] = line.get("thai", "")
        roleplay_values[f"turn_{turn_index}_support"] = line.get("english", "")

    slides.append(
        make_canva_slide(
            contract,
            "roleplay-ladder",
            "canva-slide-07-roleplay",
            "Roleplay: First self-introduction",
            roleplay_values,
            slide_background_path(7),
        )
    )

    recap_values = {
        "title": "Recap",
        "item_1": safe_list_get(recap_lines, 0, ""),
        "item_2": safe_list_get(recap_lines, 1, ""),
        "item_3": safe_list_get(recap_lines, 2, ""),
        "item_4": safe_list_get(recap_lines, 3, ""),
        "item_5": "Practise one clean greeting, name, and country line until it feels automatic.",
    }
    slides.append(
        make_canva_slide(
            contract,
            "recap-checklist",
            "canva-slide-08-recap",
            "Recap",
            recap_values,
            slide_background_path(8),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "closing",
            "canva-slide-09-closing",
            "Closing",
            {
                "lesson_id": lesson_id,
                "title": "You can now introduce\nyourself in simple Thai.",
                "line_1": "Ask a name, say your name, and say where you are from.",
                "line_2": "Keep it short, polite, and repeatable.",
                "brand": "Immersion Thai with Nine",
            },
            slide_background_path(9),
        )
    )

    return slides


def build_l003_reference_canva_slides(
    lesson_id: str,
    script: dict[str, Any],
    row: dict[str, str],
    contract: dict[str, Any],
) -> list[dict[str, Any]]:
    sections = script.get("sections") or []
    s1 = safe_list_get(sections, 0, {}) or {}
    s2 = safe_list_get(sections, 1, {}) or {}
    s3 = safe_list_get(sections, 2, {}) or {}
    s4 = safe_list_get(sections, 3, {}) or {}

    s1_focus = s1.get("languageFocus") or []
    s2_focus = s2.get("languageFocus") or []
    s3_focus = s3.get("languageFocus") or []
    s4_focus = s4.get("languageFocus") or []
    recap_lines = [str(line).strip() for line in script.get("recap", []) if str(line).strip()]
    roleplay_lines = script.get("roleplay", {}).get("lines", []) or []
    lesson_title = row.get("lesson_title") or script.get("title", lesson_id)

    objectives = [
        "1. Choose ฉัน, ผม, and เขา in the right role",
        "2. Introduce someone with นี่คือ... and เขาคือ...",
        "3. Ask who someone is with เขาคือใคร",
        "4. Confirm or reject an identification politely",
        "5. Build a short identity exchange",
    ]

    slides: list[dict[str, Any]] = []

    slides.append(
        make_canva_slide(
            contract,
            "opener",
            "canva-slide-01-opener",
            lesson_title,
            {
                "eyebrow": f"{row.get('module_title', 'Module 1')}  ·  {row.get('cefr_band', 'Level A0')}",
                "lesson_id": lesson_id,
                "title": lesson_title,
                "brand": "Immersion Thai with Nine",
            },
            slide_background_path(1),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "objectives-list",
            "canva-slide-02-objectives",
            "What you will learn",
            {
                "title": "What you will learn",
                "objective_1": objectives[0],
                "objective_2": objectives[1],
                "objective_3": objectives[2],
                "objective_4": objectives[3],
                "objective_5": objectives[4],
            },
            slide_background_path(2),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "response-triad",
            "canva-slide-03-pronouns",
            "Use pronouns as practical labels",
            {
                "section_label": "Section 1",
                "title": "Use pronouns as practical labels",
                "subtitle": "Pick the word that matches the speaker or person you are pointing to.",
                "answers_label": "Talking about yourself",
                "row_1_thai": (safe_list_get(s1_focus, 0, {}) or {}).get("thai", ""),
                "row_1_translit": (safe_list_get(s1_focus, 0, {}) or {}).get("translit", ""),
                "row_1_meaning": (safe_list_get(s1_focus, 0, {}) or {}).get("english", ""),
                "row_1_usage": "Common self-word in this course",
                "row_2_thai": (safe_list_get(s1_focus, 1, {}) or {}).get("thai", ""),
                "row_2_translit": (safe_list_get(s1_focus, 1, {}) or {}).get("translit", ""),
                "row_2_meaning": (safe_list_get(s1_focus, 1, {}) or {}).get("english", ""),
                "row_2_usage": "Use when the speaker is male",
                "softener_label": "Talking about another person",
                "row_3_thai": (safe_list_get(s1_focus, 2, {}) or {}).get("thai", ""),
                "row_3_translit": (safe_list_get(s1_focus, 2, {}) or {}).get("translit", ""),
                "row_3_meaning": (safe_list_get(s1_focus, 2, {}) or {}).get("english", ""),
                "row_3_usage": "Use for he or she from context",
            },
            slide_background_path(3),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "stacked-phrase-cards",
            "canva-slide-04-introduce-identify",
            "Introduce and identify",
            {
                "section_label": "Section 2",
                "title": "Introduce and identify",
                "subtitle": "Use these frames to point, identify, and check the person.",
                "card_1_thai": (safe_list_get(s2_focus, 0, {}) or {}).get("thai", ""),
                "card_1_translit": (safe_list_get(s2_focus, 0, {}) or {}).get("translit", ""),
                "card_1_english": (safe_list_get(s2_focus, 0, {}) or {}).get("english", ""),
                "card_2_thai": (safe_list_get(s2_focus, 1, {}) or {}).get("thai", ""),
                "card_2_translit": (safe_list_get(s2_focus, 1, {}) or {}).get("translit", ""),
                "card_2_english": (safe_list_get(s2_focus, 1, {}) or {}).get("english", ""),
                "card_3_thai": (safe_list_get(s2_focus, 2, {}) or {}).get("thai", ""),
                "card_3_translit": (safe_list_get(s2_focus, 2, {}) or {}).get("translit", ""),
                "card_3_english": (safe_list_get(s2_focus, 2, {}) or {}).get("english", ""),
            },
            slide_background_path(4),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "dual-card-shared-reply",
            "canva-slide-05-who-questions",
            "Ask who someone is",
            {
                "section_label": "Section 3",
                "title": "Ask who someone is",
                "left_label": "Ask about another person",
                "right_label": "Ask about this person",
                "left_thai": (safe_list_get(s3_focus, 0, {}) or {}).get("thai", ""),
                "left_translit": (safe_list_get(s3_focus, 0, {}) or {}).get("translit", ""),
                "left_english": (safe_list_get(s3_focus, 0, {}) or {}).get("english", ""),
                "right_thai": (safe_list_get(s3_focus, 1, {}) or {}).get("thai", ""),
                "right_translit": (safe_list_get(s3_focus, 1, {}) or {}).get("translit", ""),
                "right_english": (safe_list_get(s3_focus, 1, {}) or {}).get("english", ""),
                "shared_label": "Useful confirmation",
                "shared_thai": (safe_list_get(s3_focus, 3, {}) or {}).get("thai", ""),
                "shared_translit": (safe_list_get(s3_focus, 3, {}) or {}).get("translit", ""),
                "shared_english": "Use after the answer when the identification is correct.",
            },
            slide_background_path(5),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "drill-rows",
            "canva-slide-06-identity-exchange",
            "Build a mini identity exchange",
            {
                "section_label": "Section 4",
                "title": "Build a mini identity exchange",
                "subtitle": "Stack the exchange one line at a time.",
                "row_1_prompt_thai": (safe_list_get(s4_focus, 0, {}) or {}).get("thai", ""),
                "row_1_prompt_translit": (safe_list_get(s4_focus, 0, {}) or {}).get("translit", ""),
                "row_1_instruction": "Open the interaction politely",
                "row_2_prompt_thai": (safe_list_get(s4_focus, 1, {}) or {}).get("thai", ""),
                "row_2_prompt_translit": (safe_list_get(s4_focus, 1, {}) or {}).get("translit", ""),
                "row_2_instruction": "Ask who the person is",
                "row_3_prompt_thai": (safe_list_get(s4_focus, 3, {}) or {}).get("thai", ""),
                "row_3_prompt_translit": (safe_list_get(s4_focus, 3, {}) or {}).get("translit", ""),
                "row_3_instruction": "Add one simple country fact",
                "row_4_prompt_thai": "คุณคือ...ใช่ไหม",
                "row_4_prompt_translit": "khun kheuu... châi mái",
                "row_4_instruction": "Confirm the identity",
                "row_5_prompt_thai": "ไม่ใช่ / ใช่",
                "row_5_prompt_translit": "mâi châi / châi",
                "row_5_instruction": "Accept or repair the guess",
            },
            slide_background_path(6),
        )
    )

    roleplay_values: dict[str, Any] = {
        "title": "Roleplay: Identity exchange",
        "subtitle": "Listen, repeat, then perform the exchange.",
    }
    for turn_index, line in enumerate(roleplay_lines[:6], start=1):
        roleplay_values[f"turn_{turn_index}_speaker"] = line.get("speaker", f"Turn {turn_index}")
        roleplay_values[f"turn_{turn_index}_thai"] = line.get("thai", "")
        roleplay_values[f"turn_{turn_index}_support"] = line.get("english", "")

    slides.append(
        make_canva_slide(
            contract,
            "roleplay-ladder",
            "canva-slide-07-roleplay",
            "Roleplay: Identity exchange",
            roleplay_values,
            slide_background_path(7),
        )
    )

    recap_values = {
        "title": "Recap",
        "item_1": safe_list_get(recap_lines, 0, ""),
        "item_2": safe_list_get(recap_lines, 1, ""),
        "item_3": safe_list_get(recap_lines, 2, ""),
        "item_4": safe_list_get(recap_lines, 3, ""),
        "item_5": "Keep the exchange short: identify, ask who, confirm, and move on.",
    }
    slides.append(
        make_canva_slide(
            contract,
            "recap-checklist",
            "canva-slide-08-recap",
            "Recap",
            recap_values,
            slide_background_path(8),
        )
    )

    slides.append(
        make_canva_slide(
            contract,
            "closing",
            "canva-slide-09-closing",
            "Closing",
            {
                "lesson_id": lesson_id,
                "title": "You can now identify people\nand ask who someone is.",
                "line_1": "Use short pronoun, identity, and who-question lines with confidence.",
                "line_2": "Keep the exchange tidy and easy to repeat.",
                "brand": "Immersion Thai with Nine",
            },
            slide_background_path(9),
        )
    )

    return slides


def pick_generic_family(slide_data: dict[str, Any]) -> str:
    role = slide_data.get("role")
    layout = slide_data.get("layout")
    if role == "opener":
        return "opener"
    if role == "objectives":
        return "objectives-list"
    if role == "roleplay":
        return "roleplay-ladder"
    if role == "recap":
        return "recap-checklist"
    if role == "closing":
        return "closing"
    if layout == "drill-stack":
        return "drill-rows"
    if layout == "contrast-board":
        return "response-triad"
    if layout == "focus-card":
        return "single-focus-card"
    return "stacked-phrase-cards"


def build_generic_canva_slides(
    deck_source: dict[str, Any],
    contract: dict[str, Any],
    lesson_id: str,
    row: dict[str, str],
) -> list[dict[str, Any]]:
    slides: list[dict[str, Any]] = []

    for index, slide_data in enumerate(deck_source.get("slides", []), start=1):
        family = pick_generic_family(slide_data)
        values: dict[str, Any] = {}
        text_lines = first_text_block_lines(slide_data)
        thai_focus = slide_data.get("thaiFocus") or []

        if family == "opener":
            values = {
                "eyebrow": f"{row.get('module_title', '')}  ·  {row.get('cefr_band', '')}".strip(" ·"),
                "lesson_id": lesson_id,
                "title": slide_data.get("title", lesson_id),
                "brand": "Immersion Thai with Nine",
            }
        elif family == "objectives-list":
            values["title"] = slide_data.get("title", "What you will learn")
            for line_index, line in enumerate(text_lines[:5], start=1):
                values[f"objective_{line_index}"] = f"{line_index}. {line}"
        elif family == "single-focus-card":
            first = safe_list_get(thai_focus, 0, {}) or {}
            values = {
                "section_label": f"Section {index - 2}" if index > 2 else "Section",
                "title": slide_data.get("title", "Teaching"),
                "thai": first.get("thai", slide_data.get("title", "")),
                "translit": first.get("translit", ""),
                "english": first.get("english", ""),
                "note": "\n".join(text_lines[:2]) or "\n".join((slide_data.get("speakerNotes") or [])[:2]),
            }
        elif family == "stacked-phrase-cards":
            values = {
                "section_label": f"Section {index - 2}" if index > 2 else "Section",
                "title": slide_data.get("title", "Teaching"),
                "subtitle": slide_data.get("visualStrategy", {}).get("onScreenGoal", ""),
            }
            for card_index, item in enumerate(thai_focus[:3], start=1):
                values[f"card_{card_index}_thai"] = item.get("thai", "")
                values[f"card_{card_index}_translit"] = item.get("translit", "")
                values[f"card_{card_index}_english"] = item.get("english", "")
        elif family == "response-triad":
            values = {
                "section_label": f"Section {index - 2}" if index > 2 else "Section",
                "title": slide_data.get("title", "Responses"),
                "subtitle": slide_data.get("visualStrategy", {}).get("onScreenGoal", ""),
                "answers_label": "Direct answers",
                "softener_label": "Softener",
            }
            for triad_index, item in enumerate(thai_focus[:3], start=1):
                values[f"row_{triad_index}_thai"] = item.get("thai", "")
                values[f"row_{triad_index}_translit"] = item.get("translit", "")
                values[f"row_{triad_index}_meaning"] = item.get("english", "")
                values[f"row_{triad_index}_usage"] = safe_list_get(text_lines, triad_index - 1, "")
        elif family == "drill-rows":
            values = {
                "section_label": f"Section {index - 2}" if index > 2 else "Section",
                "title": slide_data.get("title", "Drill"),
                "subtitle": slide_data.get("visualStrategy", {}).get("onScreenGoal", ""),
            }
            for row_index, item in enumerate(thai_focus[:5], start=1):
                values[f"row_{row_index}_prompt_thai"] = f"{row_index}. {item.get('thai', '')}"
                values[f"row_{row_index}_prompt_translit"] = item.get("translit", "")
                values[f"row_{row_index}_instruction"] = safe_list_get(text_lines, row_index - 1, "Practise aloud")
        elif family == "roleplay-ladder":
            values = {
                "title": slide_data.get("title", "Roleplay"),
                "subtitle": slide_data.get("visualStrategy", {}).get("onScreenGoal", "Listen, repeat, then try the exchange."),
            }
            for turn_index, line in enumerate((slide_data.get("textBlocks") or [{}])[0].get("lines", [])[:6], start=1):
                parts = [part.strip() for part in str(line).split("|", 3)]
                speaker, thai, translit, english = (parts + ["", "", "", ""])[:4]
                values[f"turn_{turn_index}_speaker"] = speaker or f"Turn {turn_index}"
                values[f"turn_{turn_index}_thai"] = thai
                values[f"turn_{turn_index}_support"] = english
        elif family == "recap-checklist":
            values["title"] = slide_data.get("title", "Recap")
            for item_index, line in enumerate(text_lines[:5], start=1):
                values[f"item_{item_index}"] = line
        elif family == "closing":
            values = {
                "lesson_id": lesson_id,
                "title": slide_data.get("title", "Closing"),
                "line_1": safe_list_get(text_lines, 0, safe_list_get(slide_data.get("speakerNotes") or [], 0, "")),
                "line_2": safe_list_get(text_lines, 1, ""),
                "brand": "Immersion Thai with Nine",
            }

        slides.append(
            make_canva_slide(
                contract,
                family,
                f"canva-slide-{index:02d}",
                slide_data.get("title", f"Slide {index}"),
                values,
                slide_background_path(index),
            )
        )

    return slides


def build_canva_content(
    repo_root: Path,
    lesson_id: str,
    script: dict[str, Any],
    row: dict[str, str],
    deck_source: dict[str, Any],
    translit_entries: list[tuple[str, str]],
) -> dict[str, Any]:
    contract = load_canva_contract(repo_root)
    colors = (contract.get("theme") or {}).get("colors") or {}
    if lesson_id == "M01-L001":
        slides = build_l001_reference_canva_slides(lesson_id, script, row, contract)
    elif lesson_id == "M01-L002":
        slides = build_l002_reference_canva_slides(lesson_id, script, row, contract)
    elif lesson_id == "M01-L003":
        slides = build_l003_reference_canva_slides(lesson_id, script, row, contract)
    else:
        slides = build_generic_canva_slides(deck_source, contract, lesson_id, row)

    slides = normalize_canva_slides(slides, translit_entries)

    return {
        "schemaVersion": 1,
        "lessonId": lesson_id,
        "sourceDeck": artifact_name(lesson_id, "deck-source.json"),
        "layoutContract": "course/layout-contracts/canva/l001-reference-v1.json",
        "theme": {
            "thaiFont": contract.get("theme", {}).get("thaiFont", FONT_THAI),
            "latinFont": contract.get("theme", {}).get("latinFont", FONT_LATIN),
            "backgroundColor": colors.get("background", to_hex(BG_IVORY)),
            "rightZoneTint": colors.get("rightZone", to_hex(BG_SAND_LIGHT)),
            "accentColors": [
                colors.get("primary", to_hex(ACCENT_GOLD)),
                colors.get("secondary", to_hex(ACCENT_TEAL)),
                colors.get("tertiary", to_hex(ACCENT_CLAY)),
            ],
        },
        "slides": slides,
    }


def draw_background_element(draw: ImageDraw.ImageDraw, element: dict[str, Any], colors: dict[str, str], total_emu_w: int, total_emu_h: int) -> None:
    fill = colors.get(element.get("fill"), colors.get("background", "#FFFFFF"))
    outline = colors.get(element.get("line"), fill) if element.get("line") else None
    x0 = emu_to_px(element["x"], total_emu_w, CANVA_BG_WIDTH)
    y0 = emu_to_px(element["y"], total_emu_h, CANVA_BG_HEIGHT)
    x1 = emu_to_px(element["x"] + element["w"], total_emu_w, CANVA_BG_WIDTH)
    y1 = emu_to_px(element["y"] + element["h"], total_emu_h, CANVA_BG_HEIGHT)
    shape_type = element.get("type")

    if shape_type == "roundRect":
        radius = max(16, min(x1 - x0, y1 - y0) // 8)
        draw.rounded_rectangle((x0, y0, x1, y1), radius=radius, fill=fill, outline=outline, width=2 if outline else 0)
    elif shape_type == "ellipse":
        draw.ellipse((x0, y0, x1, y1), fill=fill, outline=outline, width=2 if outline else 0)
    else:
        draw.rectangle((x0, y0, x1, y1), fill=fill, outline=outline, width=2 if outline else 0)


def render_canva_backgrounds(repo_root: Path, lesson_root: Path, canva_content: dict[str, Any]) -> None:
    contract = load_canva_contract(repo_root)
    colors = (contract.get("theme") or {}).get("colors") or {}
    family_lookup = contract_family_map(contract)
    total_emu_w = int(contract.get("canvas", {}).get("widthEmu", 12191694))
    total_emu_h = int(contract.get("canvas", {}).get("heightEmu", 6858000))

    for slide in canva_content.get("slides", []):
        family = family_lookup[slide["layoutFamily"]]
        image = Image.new("RGB", (CANVA_BG_WIDTH, CANVA_BG_HEIGHT), hex_to_rgb(colors.get("background", "#FFFFFF")))
        draw = ImageDraw.Draw(image)
        for element in family.get("backgroundElements", []):
            draw_background_element(draw, element, colors, total_emu_w, total_emu_h)

        output_path = lesson_root / slide["backgroundPath"]
        output_path.parent.mkdir(parents=True, exist_ok=True)
        image.save(output_path, format="PNG")


def alignment_for_canva(value: str) -> int:
    if value == "center":
        return PP_ALIGN.CENTER
    if value == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def add_canva_text_element(slide, element: dict[str, Any]) -> None:
    box = slide.shapes.add_textbox(int(element["x"]), int(element["y"]), int(element["w"]), int(element["h"]))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.clear()

    segments = element.get("segments") or []
    if segments and "\n" not in "".join(str(segment.get("text", "")) for segment in segments):
        paragraph = frame.paragraphs[0]
        paragraph.alignment = alignment_for_canva(str(element.get("align", "left")))
        for segment in segments:
            run = paragraph.add_run()
            run.text = str(segment.get("text", ""))
            run.font.name = str(segment.get("fontName", element.get("fontName", FONT_LATIN)))
            run.font.size = Pt(float(segment.get("fontSizePt", element.get("fontSizePt", 14))))
            run.font.color.rgb = color_tuple_to_rgb(hex_to_rgb(str(segment.get("color", element.get("color", "#24333D")))))
            run.font.bold = bool(segment.get("bold", element.get("bold", False)))
            run.font.italic = bool(segment.get("italic", element.get("italic", False)))
            apply_run_font_metadata(run, font_name=run.font.name)
        return

    lines = str(element.get("value", "")).splitlines() or [str(element.get("value", ""))]
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = alignment_for_canva(str(element.get("align", "left")))
        run = paragraph.add_run()
        run.text = line
        run.font.name = str(element.get("fontName", FONT_LATIN))
        run.font.size = Pt(float(element.get("fontSizePt", 14)))
        run.font.color.rgb = color_tuple_to_rgb(hex_to_rgb(str(element.get("color", "#24333D"))))
        run.font.bold = bool(element.get("bold", False))
        run.font.italic = bool(element.get("italic", False))
        apply_run_font_metadata(run, font_name=str(element.get("fontName", FONT_LATIN)))


def render_canva_deck(lesson_root: Path, canva_content: dict[str, Any], output_pptx: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for slide_data in canva_content.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background_path = lesson_root / slide_data["backgroundPath"]
        slide.shapes.add_picture(str(background_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
        for element in slide_data.get("elements", []):
            if element.get("kind") == "image" and element.get("localPath"):
                slide.shapes.add_picture(
                    str(lesson_root / element["localPath"]),
                    int(element["x"]),
                    int(element["y"]),
                    width=int(element["w"]),
                    height=int(element["h"]),
                )
            elif element.get("kind") == "text":
                add_canva_text_element(slide, element)
        # Add PiP placeholder on top of all content
        _add_pip_placeholder(slide)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    patch_theme_fonts(output_pptx)


def write_canva_import_guide(lesson_root: Path, lesson_id: str, canva_content: dict[str, Any], output_path: Path) -> None:
    families = [slide.get("layoutFamily", "") for slide in canva_content.get("slides", [])]
    unique_families = []
    for family in families:
        if family and family not in unique_families:
            unique_families.append(family)

    guide = f"""# Canva Import Guide — {lesson_id}

## What stage 3 generated

- `{artifact_name(lesson_id, "canva-deck.pptx")}`
- `{artifact_name(lesson_id, "canva-content.json")}`
- `{artifact_name(lesson_id, "canva-import-guide.md")}`
- `canva-backgrounds/slide-XX.png`

## One-shot Canva workflow

1. Upload `Sarabun` to Canva Brand Kit.
2. Use `Sarabun` for Thai, transliteration, and English so mixed learner-facing lines stay in one font family.
3. Make sure the custom Thai font is available in Canva before importing the deck.
4. Import `{artifact_name(lesson_id, "canva-deck.pptx")}` into Canva only as a bootstrap copy.
5. Build or update a Canva master template from the same slide families listed below.
6. Treat the PNG backgrounds as locked geometry.
7. Change only text fields and image swaps in Canva.
8. Keep the inline learner-facing format as `Thai (PTM transliteration)` everywhere Thai appears.

## Font reliability rule

- This export writes Thai runs with `th-TH` language metadata and explicit `latin`, `ea`, and `cs` typefaces set to `Sarabun`.
- The deck theme also declares `Sarabun` for both Thai script and the default Latin theme font so mixed lines keep one professional reading style.
- If Canva still substitutes the font on the first import, apply `Sarabun` once with Canva's replace-all behavior inside the master template, then save that template as the stable Canva starting point.

## Slide families used in this lesson

{chr(10).join(f"- `{family}`" for family in unique_families)}

## Roundtrip rule

- Canva is the finishing surface, not the source of truth.
- If you improve spacing or placement in Canva, copy that fix back into the repo layout contract before reusing it.
- Do not rely on “edit in Canva, download, and re-import” as the system.

## Production rule

- Do not use Canva AI slide generation for production layouts.
- It is fine for brainstorming copy, but final geometry should come from the repo export pack and master template.
"""
    output_path.write_text(guide, encoding="utf-8")


def patch_theme_fonts(pptx_path: Path) -> None:
    theme_name = "ppt/theme/theme1.xml"
    with ZipFile(pptx_path, "r") as source_zip:
        if theme_name not in source_zip.namelist():
            return
        theme_xml = source_zip.read(theme_name)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            temp_path = Path(tmp.name)
            with ZipFile(temp_path, "w", compression=ZIP_DEFLATED) as target_zip:
                for member in source_zip.infolist():
                    data = source_zip.read(member.filename)
                    if member.filename == theme_name:
                        root = ET.fromstring(theme_xml)
                        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                        for branch_name in ["majorFont", "minorFont"]:
                            branch = root.find(f".//a:fontScheme/a:{branch_name}", ns)
                            if branch is None:
                                continue

                            for child_name in ["latin", "ea", "cs"]:
                                child = branch.find(f"a:{child_name}", ns)
                                if child is None:
                                    child = ET.SubElement(
                                        branch,
                                        "{http://schemas.openxmlformats.org/drawingml/2006/main}" + child_name,
                                    )
                                child.set("typeface", FONT_LATIN)

                            thai_entry = None
                            for entry in branch.findall("a:font", ns):
                                if entry.get("script") == "Thai":
                                    thai_entry = entry
                                    break
                            if thai_entry is None:
                                thai_entry = ET.SubElement(
                                    branch,
                                    "{http://schemas.openxmlformats.org/drawingml/2006/main}font",
                                )
                                thai_entry.set("script", "Thai")
                            thai_entry.set("typeface", FONT_THAI)

                        data = ET.tostring(root, encoding="utf-8", xml_declaration=True)

                    target_zip.writestr(member, data)

    temp_path.replace(pptx_path)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Render a PPTX lesson deck")
    parser.add_argument("--repo-root", required=True)
    parser.add_argument("--lesson", required=True)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    repo_root = Path(args.repo_root).resolve()
    lesson_id = args.lesson.strip()
    lesson_root = lesson_dir(repo_root, lesson_id)
    script_path = lesson_root / artifact_name(lesson_id, "script-master.json")
    if not script_path.exists():
        print(f"Missing script-master.json for {lesson_id}", file=sys.stderr)
        return 1

    script = read_json(script_path)
    row = read_blueprint_row(repo_root, lesson_id)
    deck_source, provenance = build_deck_source(repo_root, lesson_id, script, row, lesson_root)
    translit_entries = build_translit_entries(script)

    deck_source_path = lesson_root / artifact_name(lesson_id, "deck-source.json")
    provenance_path = lesson_root / artifact_name(lesson_id, "asset-provenance.json")
    output_pptx = lesson_root / artifact_name(lesson_id, "deck.pptx")
    canva_content_path = lesson_root / artifact_name(lesson_id, "canva-content.json")
    canva_deck_path = lesson_root / artifact_name(lesson_id, "canva-deck.pptx")
    canva_import_guide_path = lesson_root / artifact_name(lesson_id, "canva-import-guide.md")

    deck_source_path.write_text(json.dumps(deck_source, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    provenance_path.write_text(json.dumps(provenance, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    # Single deck output — render directly to canva-deck.pptx
    render_deck(deck_source, row, lesson_root, canva_deck_path, translit_entries)
    print(f"Rendered {lesson_id} -> {canva_deck_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
