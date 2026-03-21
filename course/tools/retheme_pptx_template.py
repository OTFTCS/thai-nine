#!/usr/bin/env python3
"""Apply a new color palette and Thai font pass to a reference PPTX deck."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor


FONT_THAI = "Noto Looped Thai UI"

OLD_TO_NEW = {
    "F0EBE3": "DBE7EC",  # right recording zone
    "B8963E": "355C7D",  # primary accent
    "3B7A72": "7A8F54",  # secondary accent
    "4A7C72": "7A8F54",
    "A0522D": "C96F4A",  # tertiary accent
    "FDF5E6": "EEF4F7",  # pale support cards
    "F0E8D8": "D2DEE5",
    "D4BE8A": "C9D7DE",
    "E0D8CE": "D7DEE3",
    "2C2420": "24333D",  # dark text
    "5A4E46": "566873",  # medium text
    "8A7E76": "82909A",  # light text
}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Retheme a PPTX deck")
    parser.add_argument("--input", required=True)
    parser.add_argument("--output", required=True)
    return parser.parse_args()


def looks_thai(text: str) -> bool:
    return any("\u0E00" <= char <= "\u0E7F" for char in text)


def mapped_rgb(rgb: RGBColor | None) -> RGBColor | None:
    if rgb is None:
        return None
    key = str(rgb)
    mapped = OLD_TO_NEW.get(key)
    if not mapped:
        return None
    return RGBColor.from_string(mapped)


def recolor_shape(shape) -> None:
    try:
        if shape.fill.type is not None and shape.fill.fore_color.type is not None:
            rgb = mapped_rgb(shape.fill.fore_color.rgb)
            if rgb is not None:
                shape.fill.fore_color.rgb = rgb
    except Exception:
        pass

    try:
        if shape.line.color.type is not None:
            rgb = mapped_rgb(shape.line.color.rgb)
            if rgb is not None:
                shape.line.color.rgb = rgb
    except Exception:
        pass

    if not hasattr(shape, "text_frame"):
        return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            try:
                if run.font.color.type is not None:
                    rgb = mapped_rgb(run.font.color.rgb)
                    if rgb is not None:
                        run.font.color.rgb = rgb
            except Exception:
                pass

            if looks_thai(run.text):
                run.font.name = FONT_THAI


def main() -> int:
    args = parse_args()
    input_path = Path(args.input).expanduser().resolve()
    output_path = Path(args.output).expanduser().resolve()

    prs = Presentation(input_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            recolor_shape(shape)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Rethemed deck -> {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
