#!/usr/bin/env python3
import csv
import json
import random
import re
import unicodedata
from pathlib import Path
from typing import Dict, Iterable, List, Tuple

from PIL import Image, ImageDraw, ImageFilter, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/Shared/work/thai-nine/Thai images/family-members-template-trial")
MANIFEST_PATH = ROOT / "carousel-manifest.json"
LAYOUT_PATH = ROOT / "template-layout.json"
PNG_DIR = ROOT / "final-png"
PREVIEW_DIR = ROOT / "preview"
PPTX_PATH = ROOT / "family-members-template-trial.pptx"
COPY_MD_PATH = ROOT / "copy-spec.md"
COPY_CSV_PATH = ROOT / "copy-spec.csv"
PPTX_ART_DIR = ROOT / ".pptx-art-cache"
ASSETS_DIR = ROOT / "template-assets"
VALIDATION_REPORT_PATH = ROOT / "transliteration-validation-report.json"

REPO_ROOT = Path("/Users/Shared/work/thai-nine")
STANDARD_PATH = REPO_ROOT / "thai-transliteration-standard.md"
QUESTION_BANK_PATH = REPO_ROOT / "src/lib/quiz/question-banks.ts"
POLICY_PATH = REPO_ROOT / "course/transliteration-policy.md"
RUNTIME_PATH = REPO_ROOT / "src/lib/quiz/transliteration.ts"
VOWELS_PATH = REPO_ROOT / "course/transliteration-ptm-vowels.json"
CONSONANTS_PATH = REPO_ROOT / "course/transliteration-ptm-consonants.json"

FONT_REG = "/System/Library/Fonts/Supplemental/Tahoma.ttf"
FONT_BOLD = "/System/Library/Fonts/Supplemental/Tahoma Bold.ttf"
FONT_FALLBACK = "/System/Library/Fonts/Supplemental/Arial Unicode.ttf"
FONT_THAI = "/System/Library/Fonts/Supplemental/Arial Unicode.ttf"

LATIN_PATTERN = re.compile(r"[A-Za-z]")
FORBIDDEN_PATTERN = re.compile(r"[əɯɤœɨɪʊŋɲɕʔːᴴᴹᴸᴿ]|\^[HMLR]|\b\w+[HMLR]\b")
TOKEN_OVERRIDES = {
    "คือ": "kheuu",
}
TEXTURE_SEEDS = {
    "cover": 11,
    "teaching": 29,
}


def load_json(path: Path):
    return json.loads(path.read_text(encoding="utf-8"))


def strip_diacritics(text: str) -> str:
    replacements = {
        "ʉ": "u",
        "ɔ": "o",
    }
    normalized = text
    for src, dst in replacements.items():
        normalized = normalized.replace(src, dst)
    return "".join(ch for ch in unicodedata.normalize("NFD", normalized) if ch.isascii() and (ch.isalpha() or ch in " -'"))


def load_font(path: str, size: int):
    try:
        return ImageFont.truetype(path, size)
    except OSError:
        return ImageFont.truetype(FONT_FALLBACK, size)


def hex_rgb(value: str) -> Tuple[int, int, int]:
    value = value.lstrip("#")
    return tuple(int(value[i:i + 2], 16) for i in (0, 2, 4))


def rgb_color(value: str) -> RGBColor:
    r, g, b = hex_rgb(value)
    return RGBColor(r, g, b)


def wrap_lines(draw: ImageDraw.ImageDraw, text: str, font, max_width: int) -> List[str]:
    if not text:
        return [""]
    words = text.split()
    lines: List[str] = []
    current = ""
    for word in words:
        candidate = word if not current else f"{current} {word}"
        if draw.textlength(candidate, font=font) <= max_width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def measure_wrapped_text(draw: ImageDraw.ImageDraw, text: str, font, max_width: int, line_gap: int = 6) -> Tuple[List[str], int]:
    lines = text.split("\n") if "\n" in text else wrap_lines(draw, text, font, max_width)
    if not lines:
        lines = [""]
    line_height = font.size + line_gap
    return lines, len(lines) * line_height


def draw_wrapped_text(
    draw: ImageDraw.ImageDraw,
    box: Dict[str, int],
    text: str,
    font,
    fill: str,
    line_gap: int = 6,
    align: str = "left",
    stroke_width: int = 0,
) -> int:
    x = box["x"]
    y = box["y"]
    max_width = box["w"]
    lines, total_height = measure_wrapped_text(draw, text, font, max_width, line_gap)
    line_height = font.size + line_gap
    for idx, line in enumerate(lines):
        line_x = x
        if align == "center":
            line_x = x + max((max_width - draw.textlength(line, font=font)) / 2, 0)
        draw.text((line_x, y + idx * line_height), line, font=font, fill=fill, stroke_width=stroke_width, stroke_fill=fill)
    return total_height


def format_breakdown_lines(entries: Iterable[Dict[str, str]]) -> str:
    return "\n".join(f"{entry['thai']} = {entry['english']} ({entry['translit']})" for entry in entries)


def make_texture(width: int, height: int, colors: Dict[str, str], seed: int) -> Image.Image:
    return Image.new("RGB", (width, height), colors["background"])


def build_template_assets(layout: Dict):
    ASSETS_DIR.mkdir(exist_ok=True)
    cover_bg = ASSETS_DIR / "cover-background.png"
    teaching_bg = ASSETS_DIR / "teaching-background.png"
    make_texture(layout["canvas"]["width"], layout["canvas"]["height"], layout["colors"], TEXTURE_SEEDS["cover"]).save(cover_bg)
    make_texture(layout["canvas"]["width"], layout["canvas"]["height"], layout["colors"], TEXTURE_SEEDS["teaching"]).save(teaching_bg)
    return {
        "cover": cover_bg,
        "teaching": teaching_bg,
    }


def fit_image(image_path: Path, box: Dict[str, int]) -> Image.Image:
    img = ImageOps.exif_transpose(Image.open(image_path).convert("RGB"))
    img = img.crop((0, 0, img.width, int(img.height * 0.84)))
    return ImageOps.fit(img, (box["w"], box["h"]), method=Image.Resampling.LANCZOS, centering=(0.5, 0.0))


def paste_image(base: Image.Image, fitted: Image.Image, box: Dict[str, int]) -> None:
    shadow = Image.new("RGBA", base.size, (0, 0, 0, 0))
    shadow_draw = ImageDraw.Draw(shadow)
    shadow_draw.rounded_rectangle(
        (box["x"] + 8, box["y"] + 12, box["x"] + box["w"] + 8, box["y"] + box["h"] + 12),
        radius=26,
        fill=(70, 52, 36, 28),
    )
    shadow = shadow.filter(ImageFilter.GaussianBlur(18))
    base.alpha_composite(shadow)
    base.alpha_composite(fitted.convert("RGBA"), (box["x"], box["y"]))


def draw_rounded(draw: ImageDraw.ImageDraw, box: Dict[str, int], fill, outline=None, width=1):
    draw.rounded_rectangle(
        (box["x"], box["y"], box["x"] + box["w"], box["y"] + box["h"]),
        radius=box["radius"],
        fill=fill,
        outline=outline,
        width=width,
    )


def draw_centered_segments(draw: ImageDraw.ImageDraw, x: int, y: int, width: int, segments: List[Tuple[str, ImageFont.FreeTypeFont, str, int]]) -> int:
    total = sum(draw.textlength(text, font=font) for text, font, _, _ in segments)
    cursor = x + max((width - total) / 2, 0)
    max_size = 0
    for text, font, color, stroke_width in segments:
        draw.text((cursor, y), text, font=font, fill=color, stroke_width=stroke_width, stroke_fill=color)
        cursor += draw.textlength(text, font=font)
        max_size = max(max_size, font.size)
    return max_size


def draw_handle(draw: ImageDraw.ImageDraw, box: Dict[str, int], text: str, font, color: str):
    draw_wrapped_text(draw, box, text, font, color, 4, "center")


def render_cover(slide: Dict, layout: Dict, colors: Dict[str, str], fonts: Dict[str, Tuple]):
    width = layout["canvas"]["width"]
    height = layout["canvas"]["height"]
    image = Image.open(ASSETS_DIR / "cover-background.png").convert("RGBA").resize((width, height))
    paste_image(image, fit_image(ROOT / slide["image"], layout["cover"]["artFrame"]), layout["cover"]["artFrame"])
    draw = ImageDraw.Draw(image)

    draw_wrapped_text(draw, layout["cover"]["title"], slide["heading"], fonts["coverTitle"], colors["text"], 10, "center")
    draw_wrapped_text(draw, layout["cover"]["thai"], slide["thai"], fonts["coverThai"], colors["coverAccent"], 12, "center", 1)
    draw_wrapped_text(draw, layout["cover"]["translit"], slide["translit"], fonts["coverTranslit"], colors["muted"], 6, "center")
    draw_wrapped_text(draw, layout["cover"]["breakdown"], format_breakdown_lines(slide["breakdown"]), fonts["coverBreakdown"], colors["text"], 8, "center")
    draw_wrapped_text(draw, layout["cover"]["exampleThai"], slide["example"]["thai"], fonts["coverExampleThai"], colors["text"], 10, "center", 1)
    draw_wrapped_text(draw, layout["cover"]["exampleTranslit"], slide["example"]["translit"], fonts["coverExampleBody"], colors["text"], 6, "center")
    draw_wrapped_text(draw, layout["cover"]["exampleEnglish"], slide["example"]["english"], fonts["coverExampleBody"], colors["text"], 6, "center")
    draw_handle(draw, {"x": 390, "y": 1306, "w": 300, "h": 24}, "@thaiwith.nine", fonts["handle"], colors["text"])
    return image.convert("RGB")


def draw_stack_block(
    draw: ImageDraw.ImageDraw,
    x: int,
    y: int,
    width: int,
    text: str,
    font,
    fill: str,
    line_gap: int = 6,
    after_gap: int = 0,
    align: str = "center",
    stroke_width: int = 0,
) -> int:
    height = draw_wrapped_text(draw, {"x": x, "y": y, "w": width, "h": 0}, text, font, fill, line_gap, align, stroke_width)
    return y + height + after_gap


def render_pair(slide: Dict, layout: Dict, colors: Dict[str, str], fonts: Dict[str, Tuple]):
    width = layout["canvas"]["width"]
    height = layout["canvas"]["height"]
    image = Image.open(ASSETS_DIR / "teaching-background.png").convert("RGBA").resize((width, height))
    pair = layout["teachingPair"]
    paste_image(image, fit_image(ROOT / slide["image"], pair["artFrame"]), pair["artFrame"])
    draw = ImageDraw.Draw(image)

    def draw_pair_block(item: Dict, block: Dict[str, int], accent: str):
        x = block["x"]
        y = block["y"]
        w = block["w"]
        y = draw_stack_block(draw, x, y, w, item["english"], fonts["english"], colors["text"], 8, 8, "center")
        draw_centered_segments(
            draw,
            x,
            y,
            w,
            [
                (item["thai"], fonts["thai"], accent, 1),
                (f"({item['translit']})", fonts["translitInline"], accent, 0),
            ],
        )
        y += max(fonts["thai"].size, fonts["translitInline"].size) + 18
        if len(item["breakdown"]) > 1:
            y = draw_stack_block(draw, x, y, w, format_breakdown_lines(item["breakdown"]), fonts["breakdown"], colors["text"], 8, 12, "center")
        y = draw_stack_block(draw, x, y, w, item["example"]["thai"], fonts["exampleThai"], colors["text"], 10, 8, "center", 1)
        y = draw_stack_block(draw, x, y, w, item["example"]["translit"], fonts["exampleBody"], colors["text"], 6, 2, "center")
        draw_stack_block(draw, x, y, w, item["example"]["english"], fonts["exampleBody"], colors["text"], 6, 0, "center")

    draw_pair_block(slide["items"][0], pair["leftBlock"], colors["leftAccent"])
    draw_pair_block(slide["items"][1], pair["rightBlock"], colors["rightAccent"])
    draw_handle(draw, pair["handle"], "@thaiwith.nine", fonts["handle"], colors["text"])
    return image.convert("RGB")


def render_single(slide: Dict, layout: Dict, colors: Dict[str, str], fonts: Dict[str, Tuple]):
    width = layout["canvas"]["width"]
    height = layout["canvas"]["height"]
    image = Image.open(ASSETS_DIR / "teaching-background.png").convert("RGBA").resize((width, height))
    single = layout["teachingSingle"]
    paste_image(image, fit_image(ROOT / slide["image"], single["artFrame"]), single["artFrame"])
    draw = ImageDraw.Draw(image)
    item = slide["item"]
    block = single["contentBlock"]
    x = block["x"]
    y = block["y"]
    w = block["w"]
    y = draw_stack_block(draw, x, y, w, item["english"], fonts["english"], colors["text"], 8, 8, "center")
    draw_centered_segments(
        draw,
        x,
        y,
        w,
        [
            (item["thai"], fonts["thai"], colors["rightAccent"], 1),
            (f"({item['translit']})", fonts["translitInline"], colors["muted"], 0),
        ],
    )
    y += max(fonts["thai"].size, fonts["translitInline"].size) + 18
    y = draw_stack_block(draw, x, y, w, format_breakdown_lines(item["breakdown"]), fonts["breakdown"], colors["text"], 8, 12, "center")
    y = draw_stack_block(draw, x, y, w, item["example"]["thai"], fonts["exampleThai"], colors["text"], 10, 8, "center", 1)
    y = draw_stack_block(draw, x, y, w, item["example"]["translit"], fonts["exampleBody"], colors["text"], 6, 2, "center")
    draw_stack_block(draw, x, y, w, item["example"]["english"], fonts["exampleBody"], colors["text"], 6, 0, "center")
    draw_handle(draw, single["handle"], "@thaiwith.nine", fonts["handle"], colors["text"])
    return image.convert("RGB")


def make_fonts(layout: Dict):
    spec = layout["fonts"]
    fonts = {}
    for key, value in spec.items():
        default_path = FONT_BOLD if value["bold"] else FONT_REG
        if key in {"coverThai", "thai", "exampleThai"}:
            default_path = FONT_THAI
        if key in {"coverBreakdown", "breakdown"}:
            default_path = FONT_FALLBACK
        fonts[key] = load_font(default_path, value["size"])
    return fonts


def repo_translit_overrides() -> Dict[str, str]:
    _ = STANDARD_PATH.read_text(encoding="utf-8")
    source = QUESTION_BANK_PATH.read_text(encoding="utf-8")
    overrides: Dict[str, str] = {}
    for thai, translit in re.findall(r'thai:\s*"([^"]+)",\s*translit:\s*"([^"]+)"', source):
        overrides.setdefault(thai, translit)
    for thai, translit in re.findall(r'"([^"]+)":\s*\{\s*translit:\s*"([^"]+)"', source):
        overrides.setdefault(thai, translit)

    policy = POLICY_PATH.read_text(encoding="utf-8")
    for thai, translit in re.findall(r"-\s+(.+?)\s+→\s+`([^`]+)`", policy):
        overrides.setdefault(thai.strip(), translit.strip())
    return overrides


def allowed_character_inventory() -> Tuple[set[str], set[str]]:
    runtime_source = RUNTIME_PATH.read_text(encoding="utf-8")
    match = re.search(r"DIACRITIC_PATTERN = /(\[[^/]+?\])/u;", runtime_source)
    runtime_chars = set()
    if match:
        runtime_chars.update(ch for ch in match.group(1) if not ch.isspace() and ch not in "[]\\p{}M")

    vowels = load_json(VOWELS_PATH)
    consonants = load_json(CONSONANTS_PATH)
    ascii_letters = set("abcdefghijklmnopqrstuvwxyz")
    for section in ("short", "long"):
        for entry in vowels[section]:
            ascii_letters.update(ch for ch in entry["translit"] if ch.isalpha())
    for token in consonants["initialTokens"] + consonants["finalTokens"]:
        ascii_letters.update(ch for ch in token if ch.isalpha())

    return ascii_letters, runtime_chars


def validate_translit(
    thai: str,
    translit: str,
    context: str,
    overrides: Dict[str, str],
    allowed_ascii: set[str],
    allowed_runtime_chars: set[str],
):
    normalized = translit.strip()
    if not normalized:
        raise ValueError(f"Missing transliteration for {context}")
    if not LATIN_PATTERN.search(normalized):
        raise ValueError(f"Transliteration lacks Latin characters for {context}: {translit}")
    if FORBIDDEN_PATTERN.search(normalized):
        raise ValueError(f"Forbidden transliteration pattern for {context}: {translit}")
    approved = overrides.get(thai)
    if approved and approved != normalized:
        raise ValueError(f'Expected approved transliteration "{approved}" for {context}, got "{translit}"')
    for thai_token, translit_token in TOKEN_OVERRIDES.items():
        if thai_token in thai and translit_token not in normalized:
            raise ValueError(f'Expected "{translit_token}" in transliteration for {context}, got "{translit}"')

    stripped = strip_diacritics(normalized)
    for ch in stripped:
        if ch in " -'":
            continue
        if ch.lower() not in allowed_ascii:
            raise ValueError(f'Unexpected transliteration character "{ch}" for {context}: {translit}')

    for ch in normalized:
        if ord(ch) < 128 or ch.isspace() or ch in "-'":
            continue
        if ch in allowed_runtime_chars:
            continue
        if "\u0300" <= ch <= "\u036f":
            continue
        raise ValueError(f'Unexpected non-ASCII transliteration character "{ch}" for {context}: {translit}')


def validate_manifest(manifest: Dict):
    overrides = repo_translit_overrides()
    allowed_ascii, allowed_runtime_chars = allowed_character_inventory()
    report = {
        "sources": [
            str(STANDARD_PATH),
            str(POLICY_PATH),
            str(RUNTIME_PATH),
            str(VOWELS_PATH),
            str(CONSONANTS_PATH),
        ],
        "notes": [
            "Validation prefers exact repo-approved transliterations when overrides exist.",
            "Runtime-allowed characters from src/lib/quiz/transliteration.ts take precedence over stale prose examples in policy docs.",
            'Hard token override preserves "kheuu" for คือ.',
        ],
        "slidesValidated": [],
    }
    for slide in manifest["slides"]:
        if slide["kind"] == "cover":
            validate_translit(slide["thai"], slide["translit"], f"slide-{slide['index']}/cover", overrides, allowed_ascii, allowed_runtime_chars)
            for idx, entry in enumerate(slide["breakdown"]):
                validate_translit(entry["thai"], entry["translit"], f"slide-{slide['index']}/breakdown-{idx}", overrides, allowed_ascii, allowed_runtime_chars)
            validate_translit(slide["example"]["thai"], slide["example"]["translit"], f"slide-{slide['index']}/example", overrides, allowed_ascii, allowed_runtime_chars)
        elif slide["kind"] == "teaching-pair":
            for pos, item in enumerate(slide["items"]):
                validate_translit(item["thai"], item["translit"], f"slide-{slide['index']}/item-{pos}", overrides, allowed_ascii, allowed_runtime_chars)
                for idx, entry in enumerate(item["breakdown"]):
                    validate_translit(entry["thai"], entry["translit"], f"slide-{slide['index']}/item-{pos}/breakdown-{idx}", overrides, allowed_ascii, allowed_runtime_chars)
                validate_translit(item["example"]["thai"], item["example"]["translit"], f"slide-{slide['index']}/item-{pos}/example", overrides, allowed_ascii, allowed_runtime_chars)
        else:
            item = slide["item"]
            validate_translit(item["thai"], item["translit"], f"slide-{slide['index']}/item", overrides, allowed_ascii, allowed_runtime_chars)
            for idx, entry in enumerate(item["breakdown"]):
                validate_translit(entry["thai"], entry["translit"], f"slide-{slide['index']}/breakdown-{idx}", overrides, allowed_ascii, allowed_runtime_chars)
            validate_translit(item["example"]["thai"], item["example"]["translit"], f"slide-{slide['index']}/example", overrides, allowed_ascii, allowed_runtime_chars)
        report["slidesValidated"].append({"index": slide["index"], "kind": slide["kind"]})
    VALIDATION_REPORT_PATH.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")


def export_copy_specs(manifest: Dict):
    md_lines = ["# Family Members Template Trial Copy", ""]
    csv_rows = [[
        "slide", "kind", "left_or_single_english", "left_or_single_thai", "left_or_single_translit",
        "left_or_single_breakdown", "left_or_single_example_thai", "left_or_single_example_translit",
        "left_or_single_example_english", "right_english", "right_thai", "right_translit", "right_breakdown",
        "right_example_thai", "right_example_translit", "right_example_english"
    ]]

    for slide in manifest["slides"]:
        md_lines.append(f"## Slide {slide['index']}")
        md_lines.append("")
        md_lines.append(f"- Image: `{slide['image']}`")
        md_lines.append(f"- Kind: `{slide['kind']}`")
        md_lines.append(f"- Prompt: {slide['image_prompt']}")
        if slide["kind"] == "cover":
            md_lines.append(f"- Thai: {slide['thai']} | {slide['translit']}")
            md_lines.append(f"- Breakdown: {format_breakdown_lines(slide['breakdown']).replace(chr(10), '; ')}")
            md_lines.append(f"- Example: {slide['example']['thai']} | {slide['example']['translit']} | {slide['example']['english']}")
            csv_rows.append([
                slide["index"], slide["kind"], slide["heading"], slide["thai"], slide["translit"],
                format_breakdown_lines(slide["breakdown"]).replace("\n", " | "), slide["example"]["thai"], slide["example"]["translit"], slide["example"]["english"],
                "", "", "", "", "", "", ""
            ])
        elif slide["kind"] == "teaching-pair":
            left, right = slide["items"]
            md_lines.append(f"- Left: {left['english']} | {left['thai']} | {left['translit']}")
            md_lines.append(f"- Left breakdown: {format_breakdown_lines(left['breakdown']).replace(chr(10), '; ')}")
            md_lines.append(f"- Left example: {left['example']['thai']} | {left['example']['translit']} | {left['example']['english']}")
            md_lines.append(f"- Right: {right['english']} | {right['thai']} | {right['translit']}")
            md_lines.append(f"- Right breakdown: {format_breakdown_lines(right['breakdown']).replace(chr(10), '; ')}")
            md_lines.append(f"- Right example: {right['example']['thai']} | {right['example']['translit']} | {right['example']['english']}")
            csv_rows.append([
                slide["index"], slide["kind"], left["english"], left["thai"], left["translit"],
                format_breakdown_lines(left["breakdown"]).replace("\n", " | "), left["example"]["thai"], left["example"]["translit"], left["example"]["english"],
                right["english"], right["thai"], right["translit"], format_breakdown_lines(right["breakdown"]).replace("\n", " | "),
                right["example"]["thai"], right["example"]["translit"], right["example"]["english"]
            ])
        else:
            item = slide["item"]
            md_lines.append(f"- Item: {item['english']} | {item['thai']} | {item['translit']}")
            md_lines.append(f"- Breakdown: {format_breakdown_lines(item['breakdown']).replace(chr(10), '; ')}")
            md_lines.append(f"- Example: {item['example']['thai']} | {item['example']['translit']} | {item['example']['english']}")
            csv_rows.append([
                slide["index"], slide["kind"], item["english"], item["thai"], item["translit"],
                format_breakdown_lines(item["breakdown"]).replace("\n", " | "), item["example"]["thai"], item["example"]["translit"], item["example"]["english"],
                "", "", "", "", "", "", ""
            ])
        md_lines.append("")

    COPY_MD_PATH.write_text("\n".join(md_lines), encoding="utf-8")
    with COPY_CSV_PATH.open("w", encoding="utf-8", newline="") as file:
        writer = csv.writer(file)
        writer.writerows(csv_rows)


def create_preview(out_files: List[Path], out_path: Path):
    thumb_w, thumb_h = 216, 270
    cols = 3
    rows = (len(out_files) + cols - 1) // cols
    margin = 20
    label_h = 28
    sheet = Image.new("RGB", (margin + cols * (thumb_w + margin), margin + rows * (thumb_h + label_h + margin)), "white")
    font = load_font(FONT_REG, 16)
    for idx, path in enumerate(out_files):
        img = Image.open(path).convert("RGB")
        img.thumbnail((thumb_w, thumb_h))
        cell = Image.new("RGB", (thumb_w, thumb_h), "#f5f3ef")
        x = (thumb_w - img.width) // 2
        y = (thumb_h - img.height) // 2
        cell.paste(img, (x, y))
        row, col = divmod(idx, cols)
        sx = margin + col * (thumb_w + margin)
        sy = margin + row * (thumb_h + label_h + margin)
        sheet.paste(cell, (sx, sy))
        ImageDraw.Draw(sheet).text((sx, sy + thumb_h + 6), path.stem, font=font, fill="#333333")
    sheet.save(out_path, quality=92)


def build_pngs(manifest: Dict, layout: Dict):
    colors = layout["colors"]
    fonts = make_fonts(layout)
    PNG_DIR.mkdir(exist_ok=True)
    PREVIEW_DIR.mkdir(exist_ok=True)

    out_files = []
    for slide in manifest["slides"]:
        if slide["kind"] == "cover":
            canvas = render_cover(slide, layout, colors, fonts)
        elif slide["kind"] == "teaching-pair":
            canvas = render_pair(slide, layout, colors, fonts)
        else:
            canvas = render_single(slide, layout, colors, fonts)
        out_path = PNG_DIR / f"{slide['index']:02d}.png"
        canvas.save(out_path, quality=95)
        out_files.append(out_path)

    create_preview(out_files, PREVIEW_DIR / "contact-sheet.jpg")


def px_to_in(px: float, slide_width_px: int) -> float:
    return 8.0 * (px / slide_width_px)


def add_textbox(
    slide,
    left_px: int,
    top_px: int,
    width_px: int,
    height_px: int,
    text: str,
    size: int,
    bold: bool,
    color: str,
    name: str,
    canvas_width: int,
    align: str = "center",
    font_name: str = "Tahoma",
):
    shape = slide.shapes.add_textbox(Inches(px_to_in(left_px, canvas_width)), Inches(px_to_in(top_px, canvas_width)),
                                     Inches(px_to_in(width_px, canvas_width)), Inches(px_to_in(height_px, canvas_width)))
    shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
        run = paragraph.runs[0]
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb_color(color)


def add_shape_rect(slide, box: Dict[str, int], fill: str, line: str, transparency: float, name: str, canvas_width: int):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(px_to_in(box["x"], canvas_width)),
        Inches(px_to_in(box["y"], canvas_width)),
        Inches(px_to_in(box["w"], canvas_width)),
        Inches(px_to_in(box["h"], canvas_width)),
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb_color(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb_color(line)
    shape.line.width = Pt(1.5)
    shape.adjustments[0] = 0.08
    return shape


def add_picture(slide, image_path: Path, box: Dict[str, int], canvas_width: int):
    slide.shapes.add_picture(
        str(image_path),
        Inches(px_to_in(box["x"], canvas_width)),
        Inches(px_to_in(box["y"], canvas_width)),
        width=Inches(px_to_in(box["w"], canvas_width)),
        height=Inches(px_to_in(box["h"], canvas_width)),
    )


def box_args(spec: Dict[str, int]) -> Tuple[int, int, int, int]:
    return spec["x"], spec["y"], spec["w"], spec["h"]


def build_pptx(manifest: Dict, layout: Dict):
    prs = Presentation()
    prs.slide_width = Inches(8)
    prs.slide_height = Inches(10)
    blank = prs.slide_layouts[6]
    for slide_data in manifest["slides"]:
        slide = prs.slides.add_slide(blank)
        png_path = PNG_DIR / f"{slide_data['index']:02d}.png"
        slide.shapes.add_picture(str(png_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(str(PPTX_PATH))


def main():
    manifest = load_json(MANIFEST_PATH)
    layout = load_json(LAYOUT_PATH)
    build_template_assets(layout)
    validate_manifest(manifest)
    export_copy_specs(manifest)
    build_pngs(manifest, layout)
    build_pptx(manifest, layout)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
