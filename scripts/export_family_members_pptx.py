#!/usr/bin/env python3
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/Shared/work/thai-nine/Thai images/family-members-carousel-v1")
DATA_PATH = ROOT / "carousel-data.json"
OUT_PATH = ROOT / "family-members-carousel-v1.pptx"

SLIDE_W = Inches(7.5)
SLIDE_H = Inches(13.333333)
Y_SHIFT = 150

COLORS = {
    "bg": RGBColor(243, 237, 228),
    "panel": RGBColor(252, 248, 242),
    "line": RGBColor(224, 208, 189),
    "card": RGBColor(255, 250, 244),
    "text": RGBColor(32, 25, 19),
    "muted": RGBColor(95, 82, 74),
    "accent": RGBColor(169, 93, 27),
    "chip": RGBColor(244, 183, 110),
    "chip_text": RGBColor(75, 46, 31),
}

PX_TO_IN = 7.5 / 1080.0


def px(value: float) -> int:
    return Inches(value * PX_TO_IN)


def add_textbox(slide, left, top, width, height, text, *, font_size=20, bold=False, color=None, name=None):
    shape = slide.shapes.add_textbox(left, top, width, height)
    if name:
      shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraphs = text.split("\n")
    for idx, line in enumerate(paragraphs):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Tahoma"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color or COLORS["text"]
    return shape


def add_chip(slide, text):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, px(56), px(54), px(240), px(52))
    chip.fill.solid()
    chip.fill.fore_color.rgb = COLORS["chip"]
    chip.line.fill.background()
    chip.adjustments[0] = 0.9
    add_textbox(slide, px(74), px(64), px(220), px(30), text, font_size=14, bold=True, color=COLORS["chip_text"], name="chip-text")


def add_panel(slide):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, px(28), px(820 + Y_SHIFT), px(1024), px(1072 - Y_SHIFT))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["panel"]
    panel.fill.transparency = 0.70
    panel.line.color.rgb = COLORS["line"]
    panel.line.width = Pt(1.6)
    panel.adjustments[0] = 0.08
    return panel


def add_pair_cards(slide):
    boxes = [
        ("left-card", px(52), px(930 + Y_SHIFT), px(468), px(934 - Y_SHIFT)),
        ("right-card", px(560), px(930 + Y_SHIFT), px(468), px(934 - Y_SHIFT)),
    ]
    for name, left, top, width, height in boxes:
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.name = name
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["card"]
        card.fill.transparency = 0.70
        card.line.color.rgb = RGBColor(230, 215, 199)
        card.line.width = Pt(1.25)
        card.adjustments[0] = 0.07


def add_single_card(slide):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, px(74), px(946 + Y_SHIFT), px(932), px(910 - Y_SHIFT))
    card.name = "single-card"
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["card"]
    card.fill.transparency = 0.70
    card.line.color.rgb = RGBColor(230, 215, 199)
    card.line.width = Pt(1.25)
    card.adjustments[0] = 0.07


def add_image(slide, image_path: Path):
    # Fit to slide width and let lower area sit behind the text panel.
    slide.shapes.add_picture(str(image_path), 0, 0, width=SLIDE_W)


def render_cover(prs, layout, slide_data):
    slide = prs.slides.add_slide(layout)
    add_image(slide, ROOT / slide_data["image"])
    add_chip(slide, "THAI WITH NINE")
    add_panel(slide)
    add_textbox(slide, px(84), px(930 + Y_SHIFT), px(900), px(104), slide_data["heading"], font_size=32, bold=True, name="title")
    add_textbox(slide, px(84), px(1068 + Y_SHIFT), px(900), px(72), slide_data["thai"], font_size=24, bold=True, color=COLORS["accent"], name="thai")
    add_textbox(slide, px(84), px(1148 + Y_SHIFT), px(860), px(42), slide_data["translit"], font_size=17, bold=True, color=COLORS["muted"], name="translit")
    add_textbox(slide, px(84), px(1212 + Y_SHIFT), px(860), px(136), "\n".join(slide_data["breakdown"]), font_size=15, bold=True, color=COLORS["muted"], name="breakdown")
    add_textbox(slide, px(84), px(1384 + Y_SHIFT), px(860), px(42), slide_data["example_thai"], font_size=17, bold=True, name="example-thai")
    add_textbox(slide, px(84), px(1440 + Y_SHIFT), px(860), px(38), slide_data["example_translit"], font_size=14, bold=True, color=COLORS["muted"], name="example-translit")
    add_textbox(slide, px(84), px(1486 + Y_SHIFT), px(860), px(38), slide_data["example_english"], font_size=14, bold=True, color=COLORS["muted"], name="example-english")


def render_pair(prs, layout, slide_data):
    slide = prs.slides.add_slide(layout)
    add_image(slide, ROOT / slide_data["image"])
    add_chip(slide, "FAMILY MEMBERS")
    add_panel(slide)
    add_pair_cards(slide)
    add_textbox(slide, px(58), px(866 + Y_SHIFT), px(520), px(40), slide_data["heading"], font_size=20, bold=True, name="heading")

    def column(prefix, item, base_x):
        add_textbox(slide, px(base_x), px(972 + Y_SHIFT), px(300), px(36), item["english"], font_size=17, bold=True, color=COLORS["accent"], name=f"{prefix}-english")
        add_textbox(slide, px(base_x), px(1028 + Y_SHIFT), px(280), px(48), item["thai"], font_size=24, bold=True, name=f"{prefix}-thai")
        add_textbox(slide, px(base_x), px(1088 + Y_SHIFT), px(280), px(36), item["translit"], font_size=14, bold=True, color=COLORS["muted"], name=f"{prefix}-translit")
        add_textbox(slide, px(base_x), px(1134 + Y_SHIFT), px(350), px(102), "\n".join(item["breakdown"]), font_size=12.5, bold=True, color=COLORS["muted"], name=f"{prefix}-breakdown")
        add_textbox(slide, px(base_x), px(1242 + Y_SHIFT), px(350), px(52), item["example_thai"], font_size=14, bold=True, name=f"{prefix}-example-thai")
        add_textbox(slide, px(base_x), px(1294 + Y_SHIFT), px(356), px(56), item["example_translit"], font_size=11.8, bold=True, color=COLORS["muted"], name=f"{prefix}-example-translit")
        add_textbox(slide, px(base_x), px(1352 + Y_SHIFT), px(340), px(56), item["example_english"], font_size=11.8, bold=True, color=COLORS["muted"], name=f"{prefix}-example-english")

    column("left", slide_data["left"], 80)
    column("right", slide_data["right"], 588)


def render_single(prs, layout, slide_data):
    slide = prs.slides.add_slide(layout)
    add_image(slide, ROOT / slide_data["image"])
    add_chip(slide, "FAMILY MEMBERS")
    add_panel(slide)
    add_single_card(slide)
    add_textbox(slide, px(58), px(866 + Y_SHIFT), px(520), px(40), slide_data["heading"], font_size=20, bold=True, name="heading")
    add_textbox(slide, px(110), px(996 + Y_SHIFT), px(300), px(40), slide_data["english"], font_size=18, bold=True, color=COLORS["accent"], name="english")
    add_textbox(slide, px(110), px(1058 + Y_SHIFT), px(500), px(54), slide_data["thai"], font_size=28, bold=True, name="thai")
    add_textbox(slide, px(110), px(1128 + Y_SHIFT), px(480), px(40), slide_data["translit"], font_size=15, bold=True, color=COLORS["muted"], name="translit")
    add_textbox(slide, px(110), px(1188 + Y_SHIFT), px(470), px(184), "\n".join(slide_data["breakdown"]), font_size=13, bold=True, color=COLORS["muted"], name="breakdown")
    add_textbox(slide, px(110), px(1400 + Y_SHIFT), px(540), px(52), slide_data["example_thai"], font_size=15, bold=True, name="example-thai")
    add_textbox(slide, px(110), px(1458 + Y_SHIFT), px(560), px(56), slide_data["example_translit"], font_size=12.2, bold=True, color=COLORS["muted"], name="example-translit")
    add_textbox(slide, px(110), px(1508 + Y_SHIFT), px(360), px(42), slide_data["example_english"], font_size=12.2, bold=True, color=COLORS["muted"], name="example-english")


def main():
    data = json.loads(DATA_PATH.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    layout = prs.slide_layouts[6]

    for slide_data in data["slides"]:
        if slide_data["kind"] == "cover":
            render_cover(prs, layout, slide_data)
        elif slide_data["kind"] == "pair":
            render_pair(prs, layout, slide_data)
        else:
            render_single(prs, layout, slide_data)

    prs.save(str(OUT_PATH))
    print(OUT_PATH)


if __name__ == "__main__":
    main()
